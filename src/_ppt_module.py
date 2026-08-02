# ==============================================================================
# MODULE: Presentation Generator
# ==============================================================================

def render_presentation_generator() -> None:
    """🎞️ AI Presentation Generator — Create beautiful, downloadable PowerPoint slides from any prompt."""
    st.markdown("### 🎞️ AI Presentation Generator")
    st.markdown(
        '<div style="color:var(--sub); font-size:0.88rem; margin-bottom:16px;">'
        "Type a topic or detailed description — EduSphere AI will build a stunning, fully-coloured "
        "PowerPoint presentation with high-contrast text containers that solve readability issues on premium wave backgrounds."
        "</div>",
        unsafe_allow_html=True,
    )

    col_left, col_right = st.columns([1.6, 1], gap="medium")

    with col_right:
        _card_open()
        st.markdown("#### ⚙️ Slide Options")
        slide_count = st.slider("Number of Content Slides", min_value=4, max_value=14, value=7)
        include_chart = st.checkbox("📊 Include Chart Slide", value=True)
        include_table = st.checkbox("📋 Include Data Table Slide", value=True)
        theme_choice = st.selectbox(
            "🎨 Presentation Design & Theme",
            [
                "🌊 Rhythm Blue Wave (Premium)",
                "🌌 Deep Space (Dark)",
                "🌅 Sunrise (Orange)",
                "🌿 Nature (Green)",
                "💜 Neon Violet",
                "🔷 Corporate Blue",
            ],
        )
        st.markdown("---")
        st.markdown(
            '<div style="font-size:0.76rem; color:var(--sub);">'
            "✨ Dynamic high-contrast glassmorphism stencils<br>"
            "💧 All slides include <b>EduSphere AI</b> watermark<br>"
            "📥 Download as <b>.pptx</b> — opens in MS PowerPoint, Google Slides, LibreOffice"
            "</div>",
            unsafe_allow_html=True,
        )
        _card_close()

    with col_left:
        _card_open()
        st.markdown("#### 📝 Describe Your Presentation")
        user_topic = st.text_area(
            "Prompt",
            placeholder=(
                "Examples:\n"
                "• Make a presentation about climate change and its global impact\n"
                "• Slides on Machine Learning: types, applications, and future\n"
                "• Business pitch for a food delivery startup in Nepal\n"
                "• History of Nepal — geography, culture, achievements"
            ),
            height=160,
            label_visibility="collapsed",
        )
        gen_btn = st.button("✨ Generate Presentation", use_container_width=True, type="primary")
        _card_close()

    if gen_btn and user_topic.strip():
        with logo_spinner("🤖 EduSphere AI is designing your presentation…"):
            _generate_and_show_ppt(
                user_topic.strip(), slide_count, include_chart, include_table, theme_choice
            )
    elif gen_btn:
        st.warning("⚠️ Please enter a topic or description first.")


# ─── PPT helpers ──────────────────────────────────────────────────────────────

def _ppt_theme_colors(theme_choice: str) -> dict:
    themes = {
        "🌊 Rhythm Blue Wave (Premium)": dict(
            bg_title=(8, 10, 36), bg_content=(12, 14, 45), bg_alt=(16, 18, 55),
            accent=(0, 240, 255), accent2=(236, 72, 153),
            text=(248, 250, 252), sub_text=(148, 163, 184),
            chart_colors=["#00f0ff", "#ec4899", "#a855f7", "#10b981", "#fb923c", "#38bdf8"],
        ),
        "🌌 Deep Space (Dark)": dict(
            bg_title=(5, 5, 20), bg_content=(10, 12, 35), bg_alt=(15, 10, 40),
            accent=(0, 180, 255), accent2=(139, 92, 246),
            text=(240, 248, 255), sub_text=(180, 200, 230),
            chart_colors=["#00b4ff", "#8b5cf6", "#ec4899", "#10b981", "#f97316", "#f59e0b"],
        ),
        "🌅 Sunrise (Orange)": dict(
            bg_title=(20, 10, 5), bg_content=(35, 18, 8), bg_alt=(45, 20, 5),
            accent=(249, 115, 22), accent2=(251, 191, 36),
            text=(255, 250, 235), sub_text=(220, 180, 130),
            chart_colors=["#f97316", "#fbbf24", "#ef4444", "#ec4899", "#8b5cf6", "#06b6d4"],
        ),
        "🌿 Nature (Green)": dict(
            bg_title=(5, 18, 10), bg_content=(8, 28, 15), bg_alt=(10, 35, 20),
            accent=(52, 211, 153), accent2=(16, 185, 129),
            text=(236, 255, 244), sub_text=(160, 220, 180),
            chart_colors=["#34d399", "#10b981", "#06b6d4", "#84cc16", "#fbbf24", "#f97316"],
        ),
        "💜 Neon Violet": dict(
            bg_title=(8, 5, 20), bg_content=(12, 8, 30), bg_alt=(18, 10, 40),
            accent=(168, 85, 247), accent2=(236, 72, 153),
            text=(250, 245, 255), sub_text=(200, 170, 240),
            chart_colors=["#a855f7", "#ec4899", "#8b5cf6", "#f43f5e", "#06b6d4", "#10b981"],
        ),
        "🔷 Corporate Blue": dict(
            bg_title=(5, 15, 40), bg_content=(8, 22, 60), bg_alt=(12, 30, 75),
            accent=(59, 130, 246), accent2=(14, 165, 233),
            text=(240, 246, 255), sub_text=(170, 200, 240),
            chart_colors=["#3b82f6", "#0ea5e9", "#6366f1", "#8b5cf6", "#10b981", "#f97316"],
        ),
    }
    return themes.get(theme_choice, themes["🌌 Deep Space (Dark)"])


def _wm(slide, _tc):
    """Add EduSphere AI watermark to bottom-right of every slide."""
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    tb = slide.shapes.add_textbox(Inches(9.8), Inches(6.9), Inches(3.0), Inches(0.35))
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = "edusphere ai"
    r.font.name = "Montserrat"
    r.font.size = Pt(8.5)
    r.font.color.rgb = RGBColor(120, 140, 180)
    r.font.bold = True
    r.font.italic = True


def _set_bg(slide, rgb):
    from pptx.dml.color import RGBColor
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*rgb)


def _rect(slide, l, t, w, h, rgb):
    from pptx.util import Inches
    from pptx.dml.color import RGBColor
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = RGBColor(*rgb)
    s.line.fill.background()
    return s


def _tb(slide, text, l, t, w, h, size, bold=False, color=(255, 255, 255), align=None, italic=False):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    if align is None:
        align = PP_ALIGN.LEFT
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Space Grotesk"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = RGBColor(*color)
    return box


def _apply_animation_tag(slide):
    try:
        slide.slide_show_transition.type = 3  # MSO_ANIMATION.WHEEL / morph simulation
        slide.slide_show_transition.speed = 1 # fast
    except Exception:
        pass


def _generate_and_show_ppt(
    topic: str, slide_count: int, include_chart: bool, include_table: bool, theme_choice: str
) -> None:
    import io
    import json
    import re
    import os
    import datetime as dtm

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
    except ImportError as exc:
        st.error(f"❌ Missing library: {exc}. Run `pip install python-pptx matplotlib`")
        return

    tc = _ppt_theme_colors(theme_choice)
    A = tc["accent"]
    A2 = tc["accent2"]
    TX = tc["text"]
    SB = tc["sub_text"]
    BG0 = tc["bg_title"]
    BG1 = tc["bg_content"]
    BG2 = tc["bg_alt"]
    CC = tc["chart_colors"]

    CENTER = PP_ALIGN.CENTER
    LEFT = PP_ALIGN.LEFT

    # ── AI: slide outline with extra bullet point details ────────────────────────
    ai_prompt = (
        f'Create a high-quality PowerPoint presentation outline for: "{topic}"\n'
        f"Generate exactly {slide_count} slides.\n"
        "Each content slide MUST have: \n"
        "1. A short, punchy title\n"
        "2. Exactly 3-4 structured, descriptive bullet points. Every bullet point MUST contain a bold title followed by a detail explanation, like: '**Topic Header:** Detailed explanation text.'\n"
        "Return ONLY a valid JSON array format, no other text:\n"
        '[{"slide_num":1,"title":"Title","points":["**Header 1:** Detailed explanation","**Header 2:** Detailed explanation"]}]'
    )
    try:
        raw = groq_chat(ai_prompt, system="Return ONLY valid JSON arrays. No markdown formatting, no commentary.")
        raw = re.sub(r"```(?:json)?", "", raw).strip().rstrip("`").strip()
        slides_data = json.loads(raw)
        if not isinstance(slides_data, list):
            raise ValueError("Not a list")
    except Exception:
        slides_data = [
            {
                "slide_num": i + 1,
                "title": f"Core Pillar {i + 1}",
                "points": [
                    f"**Strategic Objective {i + 1}:** Detailed explanation of the primary goal and milestone targets.",
                    "**Data Grounding:** Evidence-based statistics supporting this core pillar and framework.",
                    "**Market Alignment:** Exploring consumer fitment, demographics, and execution strategies.",
                ],
            }
            for i in range(slide_count)
        ]

    # ── AI: chart data ─────────────────────────────────────────────────────────
    chart_data = None
    if include_chart:
        try:
            cr = groq_chat(
                f'For "{topic}", create 5 data points for a bar chart. '
                'Return ONLY JSON: {"title":"Chart Title","labels":["A","B","C","D","E"],"values":[45,78,62,91,55]}',
                system="Return ONLY valid JSON. No markdown.",
            )
            cr = re.sub(r"```(?:json)?", "", cr).strip().rstrip("`").strip()
            chart_data = json.loads(cr)
        except Exception:
            chart_data = {
                "title": f"Key Metrics: {topic[:35]}",
                "labels": ["Pillar 1", "Pillar 2", "Pillar 3", "Pillar 4", "Pillar 5"],
                "values": [65, 82, 55, 90, 73],
            }

    # ── AI: table data ─────────────────────────────────────────────────────────
    table_data = None
    if include_table:
        try:
            tr = groq_chat(
                f'For "{topic}", create a comparison table. '
                'Return ONLY JSON: {"title":"Table Title","headers":["Col1","Col2","Col3"],'
                '"rows":[["v1","v2","v3"],["v4","v5","v6"],["v7","v8","v9"],["v10","v11","v12"]]}',
                system="Return ONLY valid JSON. No markdown.",
            )
            tr = re.sub(r"```(?:json)?", "", tr).strip().rstrip("`").strip()
            table_data = json.loads(tr)
        except Exception:
            table_data = {
                "title": f"Overview: {topic[:35]}",
                "headers": ["Aspect", "Details", "Impact Assessment"],
                "rows": [
                    ["Pillar 1", "Implementation Stage A", "High Return"],
                    ["Pillar 2", "Implementation Stage B", "Medium Stability"],
                    ["Pillar 3", "Implementation Stage C", "High Return"],
                    ["Pillar 4", "Implementation Stage D", "Low Risk"],
                ],
            }

    # ── Build PowerPoint ───────────────────────────────────────────────────────
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    bg_img_path = "assets/ppt_wave_bg.png"
    has_custom_bg = os.path.exists(bg_img_path)

    # Helper function to place the premium blue wave background image
    def apply_slide_bg(slide, bg_color):
        if theme_choice == "🌊 Rhythm Blue Wave (Premium)" and has_custom_bg:
            slide.shapes.add_picture(bg_img_path, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
        else:
            _set_bg(slide, bg_color)

    # === Slide 1: Title ===
    sl = prs.slides.add_slide(blank)
    apply_slide_bg(sl, BG0)
    _apply_animation_tag(sl)

    # Semi-transparent glassmorphic title card to separate text from background wave flares
    title_glass = sl.shapes.add_shape(1, Inches(0.4), Inches(1.8), Inches(12.53), Inches(3.2))
    title_glass.fill.solid()
    title_glass.fill.fore_color.rgb = RGBColor(15, 23, 42) # dark slate base
    title_glass.line.color.rgb = RGBColor(*A)
    title_glass.line.width = Pt(1.5)

    _rect(sl, 0, 0, 13.33, 0.09, A)           # top bar
    _rect(sl, 0, 0.09, 0.09, 7.41, A2)         # left strip
    _tb(sl, topic.upper(), 0.5, 2.3, 12.33, 1.6, 38, bold=True, color=TX, align=CENTER)
    _rect(sl, 4.5, 3.95, 4.3, 0.05, A)
    _tb(sl, "Powered by EduSphere AI  ·  Premium Presentation Studio", 0.5, 4.15, 12.33, 0.55,
        15, italic=True, color=SB, align=CENTER)
    _tb(sl, str(dtm.datetime.now().year), 0.5, 6.6, 12.33, 0.4, 11, color=SB, align=CENTER, italic=True)
    _wm(sl, tc)

    # === Content Slides ===
    for idx, sinfo in enumerate(slides_data):
        sl = prs.slides.add_slide(blank)
        _apply_animation_tag(sl)
        bg_c = BG2 if idx % 2 == 0 else BG1
        apply_slide_bg(sl, bg_c)
        _rect(sl, 0, 0, 13.33, 0.07, A)

        # Slide number badge
        badge = sl.shapes.add_shape(1, Inches(11.9), Inches(0.12), Inches(1.2), Inches(0.35))
        badge.fill.solid()
        badge.fill.fore_color.rgb = RGBColor(*A2)
        badge.line.fill.background()
        _tb(sl, f"{idx+1}/{len(slides_data)}", 11.95, 0.14, 1.1, 0.3, 8,
            bold=True, color=BG0, align=CENTER)

        slide_title = sinfo.get("title", f"Slide {idx+1}")
        
        # High contrast title box backing
        title_backing = sl.shapes.add_shape(1, Inches(0.28), Inches(0.15), Inches(11.4), Inches(0.85))
        title_backing.fill.solid()
        title_backing.fill.fore_color.rgb = RGBColor(12, 15, 35)
        title_backing.line.fill.background()
        
        _tb(sl, slide_title, 0.38, 0.22, 11.2, 0.85, 27, bold=True, color=A)
        _rect(sl, 0.38, 0.95, 5.0, 0.04, A2)

        # Content Rendering (Left bar visual + bullet details parser)
        _rect(sl, 0.35, 1.5, 0.05, 5.0, A2)
        points = sinfo.get("points", [])
        
        for b_idx, pt in enumerate(points[:5]):
            top_pos = 1.48 + b_idx * 1.35
            
            # Semi-transparent backing panel for each bullet point to ensure max readability on light waves/stars
            panel = sl.shapes.add_shape(1, Inches(0.45), Inches(top_pos - 0.08), Inches(12.38), Inches(1.15))
            panel.fill.solid()
            panel.fill.fore_color.rgb = RGBColor(16, 20, 48) # high contrast dark navy card
            panel.line.color.rgb = RGBColor(40, 50, 90)
            panel.line.width = Pt(1.0)
            
            # Sub-elements for bullet points (Bold Title + Details layout)
            if "**" in pt:
                parts = pt.split("**", 2)
                bold_header = parts[1].replace(":", "").strip()
                detail_text = parts[2].strip()
            else:
                bold_header = "Objective"
                detail_text = pt

            # Bullet Indicator icon
            _tb(sl, "⚡", 0.58, top_pos, 0.4, 0.4, 13, bold=True, color=A)
            
            # Content textbox with header in accent, body in white
            tb_box = sl.shapes.add_textbox(Inches(0.95), Inches(top_pos - 0.05), Inches(11.8), Inches(1.1))
            tf = tb_box.text_frame
            tf.word_wrap = True
            
            # Paragraph 1: Bold subheader
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = bold_header
            r.font.name = "Montserrat"
            r.font.size = Pt(15.5)
            r.font.bold = True
            r.font.color.rgb = RGBColor(*A)

            # Paragraph 2: Detailed text body
            p2 = tf.add_paragraph()
            r2 = p2.add_run()
            r2.text = detail_text
            r2.font.name = "Space Grotesk"
            r2.font.size = Pt(13)
            r2.font.color.rgb = RGBColor(*TX)

        _wm(sl, tc)

    # === Chart Slide ===
    if include_chart and chart_data:
        sl = prs.slides.add_slide(blank)
        _apply_animation_tag(sl)
        apply_slide_bg(sl, BG0)
        _rect(sl, 0, 0, 13.33, 0.07, A)
        
        # High contrast title box backing
        title_backing = sl.shapes.add_shape(1, Inches(0.28), Inches(0.1), Inches(12.77), Inches(0.85))
        title_backing.fill.solid()
        title_backing.fill.fore_color.rgb = RGBColor(12, 15, 35)
        title_backing.line.fill.background()

        _tb(sl, "📊 " + chart_data.get("title", "Data Visualization"),
            0.35, 0.14, 12.8, 0.78, 24, bold=True, color=A)
        _rect(sl, 0.35, 0.85, 4.5, 0.04, A2)

        fig, ax = plt.subplots(figsize=(11.5, 5.5))
        fig.patch.set_facecolor(f"#{BG0[0]:02x}{BG0[1]:02x}{BG0[2]:02x}")
        ax.set_facecolor(f"#{BG1[0]:02x}{BG1[1]:02x}{BG1[2]:02x}")
        labels = chart_data.get("labels", [])
        values = chart_data.get("values", [])
        bar_colors_list = [CC[i % len(CC)] for i in range(len(labels))]
        bars = ax.bar(labels, values, color=bar_colors_list, edgecolor="none", width=0.55)
        for br, val in zip(bars, values):
            ax.text(
                br.get_x() + br.get_width() / 2, br.get_height() + 1,
                str(val), ha="center", va="bottom",
                color="white", fontsize=11, fontweight="bold",
            )
        ax.set_xticks(range(len(labels)))
        ax.set_xticklabels(labels, color=f"#{TX[0]:02x}{TX[1]:02x}{TX[2]:02x}", fontsize=10)
        ax.tick_params(axis="y", colors=f"#{SB[0]:02x}{SB[1]:02x}{SB[2]:02x}")
        for sp in ["top", "right"]:
            ax.spines[sp].set_visible(False)
        ax.spines["left"].set_color(f"#{SB[0]:02x}{SB[1]:02x}{SB[2]:02x}")
        ax.spines["bottom"].set_color(f"#{SB[0]:02x}{SB[1]:02x}{SB[2]:02x}")
        ax.yaxis.grid(True, alpha=0.2, color="white")
        ax.set_axisbelow(True)
        plt.tight_layout(pad=0.5)
        chart_buf = io.BytesIO()
        plt.savefig(chart_buf, format="png", dpi=150, bbox_inches="tight",
                    facecolor=fig.get_facecolor())
        plt.close(fig)
        chart_buf.seek(0)
        
        # High contrast picture framing card
        pic_frame = sl.shapes.add_shape(1, Inches(0.35), Inches(0.95), Inches(12.63), Inches(6.25))
        pic_frame.fill.solid()
        pic_frame.fill.fore_color.rgb = RGBColor(10, 12, 32)
        pic_frame.line.color.rgb = RGBColor(40, 50, 90)
        pic_frame.line.width = Pt(1.5)

        sl.shapes.add_picture(chart_buf, Inches(0.4), Inches(1.0), Inches(12.5), Inches(6.1))
        _wm(sl, tc)

    # === Table Slide ===
    if include_table and table_data:
        from pptx.enum.text import PP_ALIGN as _PA
        sl = prs.slides.add_slide(blank)
        _apply_animation_tag(sl)
        apply_slide_bg(sl, BG1)
        _rect(sl, 0, 0, 13.33, 0.07, A)
        
        # High contrast title box backing
        title_backing = sl.shapes.add_shape(1, Inches(0.28), Inches(0.1), Inches(12.77), Inches(0.85))
        title_backing.fill.solid()
        title_backing.fill.fore_color.rgb = RGBColor(12, 15, 35)
        title_backing.line.fill.background()

        _tb(sl, "📋 " + table_data.get("title", "Data Overview"),
            0.35, 0.14, 12.8, 0.78, 24, bold=True, color=A)
        _rect(sl, 0.35, 0.85, 4.5, 0.04, A2)
        headers = table_data.get("headers", [])
        rows = table_data.get("rows", [])
        if headers and rows:
            nc = len(headers)
            nr = len(rows) + 1
            
            # High contrast table framing panel
            tbl_frame = sl.shapes.add_shape(1, Inches(0.35), Inches(0.95), Inches(12.6), Inches(min(6.2, nr * 0.94 + 0.2)))
            tbl_frame.fill.solid()
            tbl_frame.fill.fore_color.rgb = RGBColor(10, 12, 32)
            tbl_frame.line.color.rgb = RGBColor(40, 50, 90)
            tbl_frame.line.width = Pt(1.5)

            tbl = sl.shapes.add_table(
                nr, nc, Inches(0.45), Inches(1.05), Inches(12.4), Inches(min(5.8, nr * 0.9))
            ).table
            # Header row
            for ci, hdr in enumerate(headers):
                cell = tbl.cell(0, ci)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(*A)
                p = cell.text_frame.paragraphs[0]
                p.alignment = _PA.CENTER
                r = p.add_run()
                r.text = str(hdr)
                r.font.name = "Montserrat"
                r.font.size = Pt(13)
                r.font.bold = True
                r.font.color.rgb = RGBColor(*BG0)
            # Data rows
            for ri, row_vals in enumerate(rows):
                row_bg = (
                    (min(A[0] // 8, 255), min(A[1] // 8, 255), min(A[2] // 8, 255))
                    if ri % 2 == 0
                    else (min(BG2[0] + 15, 255), min(BG2[1] + 15, 255), min(BG2[2] + 15, 255))
                )
                for ci in range(min(nc, len(row_vals))):
                    cell = tbl.cell(ri + 1, ci)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(*row_bg)
                    p = cell.text_frame.paragraphs[0]
                    p.alignment = _PA.CENTER
                    r = p.add_run()
                    r.text = str(row_vals[ci])
                    r.font.name = "Space Grotesk"
                    r.font.size = Pt(12)
                    r.font.color.rgb = RGBColor(*TX)
        _wm(sl, tc)

    # === Thank You Slide ===
    sl = prs.slides.add_slide(blank)
    _apply_animation_tag(sl)
    apply_slide_bg(sl, BG0)
    _rect(sl, 0, 0, 13.33, 0.08, A)
    _rect(sl, 0, 0.08, 0.08, 7.42, A2)

    # High contrast card backing for closing slide text
    close_card = sl.shapes.add_shape(1, Inches(2.5), Inches(1.8), Inches(8.33), Inches(3.8))
    close_card.fill.solid()
    close_card.fill.fore_color.rgb = RGBColor(12, 15, 35)
    close_card.line.color.rgb = RGBColor(*A2)
    close_card.line.width = Pt(1.5)

    _tb(sl, "🙏 Thank You", 0.5, 2.3, 12.3, 1.0, 42, bold=True, color=A, align=CENTER)
    _rect(sl, 4.5, 3.55, 4.3, 0.05, A2)
    _tb(sl, topic, 0.5, 3.65, 12.3, 0.6, 18, italic=True, color=SB, align=CENTER)
    _tb(sl, "Powered by EduSphere AI", 0.5, 4.5, 12.3, 0.45, 13, italic=True, color=SB, align=CENTER)
    _wm(sl, tc)

    # ── Save & download ────────────────────────────────────────────────────────
    ppt_buf = io.BytesIO()
    prs.save(ppt_buf)
    ppt_buf.seek(0)

    total = (
        1
        + len(slides_data)
        + (1 if include_chart and chart_data else 0)
        + (1 if include_table and table_data else 0)
        + 1
    )
    st.success(f"✅ Presentation generated! **{total} slides** ready.")
    st.markdown(
        f"""
        <div style="background:var(--card); border:1px solid rgba(255,255,255,0.12); border-radius:14px;
                    padding:22px; margin:14px 0; text-align:center;">
            <div style="font-size:2.8rem; margin-bottom:8px;">🎉</div>
            <div style="font-size:1.1rem; font-weight:700; color:var(--text); margin-bottom:6px;">
                Your Presentation is Ready!
            </div>
            <div style="font-size:0.85rem; color:var(--sub);">
                {total} slides &nbsp;·&nbsp; {theme_choice} theme &nbsp;·&nbsp; EduSphere AI watermark on every slide
            </div>
        </div>
        """,
        unsafe_allow_html=True,
    )
    safe = "".join(c if c.isalnum() or c in " _-" else "_" for c in topic[:40]).strip().replace(" ", "_")
    st.download_button(
        label="📥 Download Presentation (.pptx)",
        data=ppt_buf.getvalue(),
        file_name=f"EduSphere_{safe}.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        use_container_width=True,
        type="primary",
    )

    # Slide structure preview
    st.markdown("#### 📋 Slide Structure Preview")
    all_titles = ["🎯 Title Slide"] + [s.get("title", f"Slide {i+1}") for i, s in enumerate(slides_data)]
    if include_chart and chart_data:
        all_titles.append("📊 Chart Slide")
    if include_table and table_data:
        all_titles.append("📋 Table Slide")
    all_titles.append("🙏 Thank You")

    pcols = st.columns(min(4, len(all_titles)))
    for i, ttl in enumerate(all_titles[:16]):
        with pcols[i % len(pcols)]:
            st.markdown(
                f'<div style="background:var(--card); border:1px solid rgba(255,255,255,0.08);'
                f'border-radius:8px; padding:8px 10px; margin-bottom:8px; font-size:0.74rem; color:var(--sub);">'
                f'<span style="color:var(--accent); font-weight:700;">#{i+1}</span> {ttl}'
                f"</div>",
                unsafe_allow_html=True,
            )
