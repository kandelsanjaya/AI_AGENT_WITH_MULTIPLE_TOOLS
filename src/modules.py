"""
modules.py
==========
All 12 feature-module UI implementations for EduSphere AI.

Each module is a standalone function. All modules use the refactored
utility helpers (guardrails, GROQ client, RAG pipeline, etc.).

Module 12 (Resume Builder) is new — generates professional resumes
via LLM and exports them as downloadable PDFs.
"""

from __future__ import annotations

import datetime
import io
import logging
from typing import Optional

import streamlit as st
from PIL import Image

from .config import CRISIS_RESOURCES
from .rag import VectorIndex, VectorStoreManager, extract_pdf_chunks, load_embedding_model
from .utils import (
    evaluate_guardrails,
    fetch_website_content,
    get_download_link,
    get_pdf_download_link,
    get_audio_download_link,
    groq_chat,
    groq_chat_with_history,
    stream_to_ui,
    duckduckgo_search,
    detect_greeting,
    get_greeting_response,
    search_lict_knowledge,
    log_interaction_to_json,
    logo_spinner,
)

log = logging.getLogger("edusphere.modules")

# Shared RAG manager (model is already cached by Streamlit)
_rag_manager = VectorStoreManager(load_embedding_model())
# ==============================================================================
# Shared helpers
# ==============================================================================

def _card_open(extra_style: str = "") -> None:
    st.markdown(f'<div class="g-card" style="{extra_style}">', unsafe_allow_html=True)


def _card_close() -> None:
    st.markdown("</div>", unsafe_allow_html=True)


def _render_chat_history() -> None:
    """Render the conversation history as a messenger-style chat using st.chat_message."""
    history = st.session_state.get("chat_history", [])
    if not history:
        st.markdown(
            '<div style="text-align:center;color:#475569;margin-top:80px;font-size:0.95rem;">'
            '🎓 Ask me anything — type below or use the mic</div>',
            unsafe_allow_html=True,
        )
        return
    for msg in history:
        role = "user" if msg["role"] == "user" else "assistant"
        avatar = "👤" if role == "user" else "🎓"
        with st.chat_message(role, avatar=avatar):
            st.markdown(msg["msg"])


def _render_export_buttons(content: str, base_name: str = "output") -> None:
    """Render a file type dropdown and download button for exporting content."""
    _card_open("margin-top: 10px;")
    st.markdown("**💾 Export Output**")
    col_sel, col_btn = st.columns([2, 1], vertical_alignment="bottom")
    with col_sel:
        import hashlib
        key_hash = hashlib.md5(f"{base_name}_{content[:50]}".encode("utf-8")).hexdigest()
        format_choice = st.selectbox(
            "Select Format",
            ["Plain Text (.txt)", "PDF Document (.pdf)", "Audio Podcast (.mp3)"],
            key=f"fmt_choice_{key_hash}"
        )
    with col_btn:
        if format_choice == "Plain Text (.txt)":
            st.download_button(
                label="⬇️ Download",
                data=content,
                file_name=f"{base_name}.txt",
                mime="text/plain",
                key=f"dl_txt_{key_hash}",
                use_container_width=True
            )
        elif format_choice == "PDF Document (.pdf)":
            from .utils import generate_pdf_bytes
            pdf_bytes = generate_pdf_bytes(content)
            st.download_button(
                label="⬇️ Download",
                data=pdf_bytes,
                file_name=f"{base_name}.pdf",
                mime="application/pdf",
                key=f"dl_pdf_{key_hash}",
                use_container_width=True
            )
        elif format_choice == "Audio Podcast (.mp3)":
            from .utils import generate_audio_bytes
            audio_bytes = generate_audio_bytes(content[:2000])
            if audio_bytes:
                st.download_button(
                    label="⬇️ Download",
                    data=audio_bytes,
                    file_name=f"{base_name}.mp3",
                    mime="audio/mpeg",
                    key=f"dl_mp3_{key_hash}",
                    use_container_width=True
                )
            else:
                st.info("🔇 MP3 requires gTTS", icon="🔇")
    _card_close()




# ==============================================================================
# MODULE 0 — Dashboard (Lucy-AI style workspace with EduSphere cards)
# ==============================================================================

def render_dashboard() -> None:
    """🏠 EduSphere AI Dashboard — Premium UI Workspace."""
    user = st.session_state.user_info or {}
    name = user.get("name", "User")
    role = user.get("role", "Student")
    now = datetime.datetime.now()
    hour = now.hour
    greeting = "Good morning" if hour < 12 else ("Good afternoon" if hour < 17 else "Good evening")

    total_q = st.session_state.get("total_queries", 0)
    chat_len = len(st.session_state.get("chat_history", []))
    session_mins = int((now - st.session_state.get("session_start", now)).total_seconds() / 60)

    # ── Premium Welcome Banner ──
    st.markdown(
        f"""
        <div style="background: linear-gradient(135deg, rgba(59, 130, 246, 0.15) 0%, rgba(139, 92, 246, 0.15) 100%);
                    border: 1px solid rgba(255,255,255,0.1); border-radius: 16px; padding: 32px 40px;
                    margin-bottom: 24px; position: relative; overflow: hidden; backdrop-filter: blur(10px);">
            <div style="position: absolute; top: -50%; left: -10%; width: 50%; height: 200%; 
                        background: radial-gradient(circle, rgba(139, 92, 246, 0.2) 0%, transparent 70%); 
                        transform: rotate(-45deg); pointer-events: none;"></div>
            <div style="position: relative; z-index: 1;">
                <h1 style="margin:0; font-size: 2.4rem; font-family: 'Space Grotesk', sans-serif; font-weight: 800; letter-spacing: -0.5px;">
                    {greeting}, <span style="background: linear-gradient(90deg, #60a5fa, #a78bfa); -webkit-background-clip: text; -webkit-text-fill-color: transparent;">{name}</span> 👋
                </h1>
                <p style="margin: 8px 0 0 0; font-size: 1.05rem; color: var(--sub); max-width: 600px;">
                    Welcome to your EduSphere AI workspace. All systems are online and ready to assist you.
                </p>
            </div>
        </div>
        """,
        unsafe_allow_html=True
    )

    # ── Quick Ask Integrated Bar ──
    quick_q = st.chat_input("Ask EduSphere AI anything right here...")
    if quick_q:
        with logo_spinner("Thinking..."):
            answer = groq_chat(
                quick_q.strip(),
                system="You are EduSphere AI, a helpful educational assistant. Give a concise but thorough answer."
            )
            st.session_state.total_queries = total_q + 1
        st.markdown(
            f"""
            <div style="background: rgba(255,255,255,0.05); border-left: 4px solid var(--accent); padding: 16px; border-radius: 8px; margin-bottom: 24px;">
                <strong style="color: var(--accent);">EduSphere AI:</strong><br><br>{answer}
            </div>
            """, 
            unsafe_allow_html=True
        )


    # ── Main area + right panel ──
    left_col, right_col = st.columns([2.4, 1], gap="large")

    with left_col:
        # ── Stat Row (Glassmorphic) ──
        st.markdown("#### 📊 Real-time Metrics")
        st.markdown(
            f"""
            <div style="display:grid; grid-template-columns:repeat(4,1fr); gap:16px; margin-bottom:28px;">
                <div style="background: rgba(255,255,255,0.03); border: 1px solid rgba(255,255,255,0.08);
                            border-radius: 16px; padding: 20px; text-align: center; box-shadow: 0 4px 30px rgba(0, 0, 0, 0.1);
                            backdrop-filter: blur(5px); transition: transform 0.2s ease;">
                    <div style="font-size:1.8rem; margin-bottom: 4px;">⚡</div>
                    <div style="font-size:2.2rem; font-weight:800; font-family:'Orbitron',monospace; 
                                background: linear-gradient(180deg, #fff, #94a3b8); -webkit-background-clip: text; -webkit-text-fill-color: transparent;">
                        {total_q}
                    </div>
                    <div style="font-size:0.75rem; color:var(--sub); font-weight:600; text-transform:uppercase; letter-spacing:1px; margin-top:4px;">Queries</div>
                </div>
                
                <div style="background: rgba(255,255,255,0.03); border: 1px solid rgba(255,255,255,0.08);
                            border-radius: 16px; padding: 20px; text-align: center; box-shadow: 0 4px 30px rgba(0, 0, 0, 0.1);
                            backdrop-filter: blur(5px); transition: transform 0.2s ease;">
                    <div style="font-size:1.8rem; margin-bottom: 4px;">💬</div>
                    <div style="font-size:2.2rem; font-weight:800; font-family:'Orbitron',monospace; 
                                background: linear-gradient(180deg, #fff, #94a3b8); -webkit-background-clip: text; -webkit-text-fill-color: transparent;">
                        {chat_len}
                    </div>
                    <div style="font-size:0.75rem; color:var(--sub); font-weight:600; text-transform:uppercase; letter-spacing:1px; margin-top:4px;">Messages</div>
                </div>

                <div style="background: rgba(255,255,255,0.03); border: 1px solid rgba(255,255,255,0.08);
                            border-radius: 16px; padding: 20px; text-align: center; box-shadow: 0 4px 30px rgba(0, 0, 0, 0.1);
                            backdrop-filter: blur(5px); transition: transform 0.2s ease;">
                    <div style="font-size:1.8rem; margin-bottom: 4px;">⏱️</div>
                    <div style="font-size:2.2rem; font-weight:800; font-family:'Orbitron',monospace; 
                                background: linear-gradient(180deg, #fff, #94a3b8); -webkit-background-clip: text; -webkit-text-fill-color: transparent;">
                        {session_mins}m
                    </div>
                    <div style="font-size:0.75rem; color:var(--sub); font-weight:600; text-transform:uppercase; letter-spacing:1px; margin-top:4px;">Session</div>
                </div>

                <div style="background: rgba(255,255,255,0.03); border: 1px solid rgba(255,255,255,0.08);
                            border-radius: 16px; padding: 20px; text-align: center; box-shadow: 0 4px 30px rgba(0, 0, 0, 0.1);
                            backdrop-filter: blur(5px); transition: transform 0.2s ease;">
                    <div style="font-size:1.8rem; margin-bottom: 4px;">👤</div>
                    <div style="font-size:1.6rem; font-weight:800; font-family:'Orbitron',monospace; 
                                background: linear-gradient(180deg, #fff, #94a3b8); -webkit-background-clip: text; -webkit-text-fill-color: transparent; margin: 6px 0;">
                        {role[:5].upper()}
                    </div>
                    <div style="font-size:0.75rem; color:var(--sub); font-weight:600; text-transform:uppercase; letter-spacing:1px; margin-top:4px;">Access</div>
                </div>
            </div>
            """,
            unsafe_allow_html=True
        )

        # ── Quick Module Access (Premium Grid) ──
        st.markdown("#### 🚀 Quick Module Access")
        modules_grid = [
            ("📝", "Summariser", "📝 Executive Summariser", "linear-gradient(135deg, #0ea5e9, #2563eb)"),
            ("🎨", "Image Gen", "🎨 AI Image Generator", "linear-gradient(135deg, #8b5cf6, #d946ef)"),
            ("💻", "Code Lab", "💻 Code Lab & Explainer", "linear-gradient(135deg, #10b981, #059669)"),
            ("🌍", "Translator", "🌍 Academic Translator", "linear-gradient(135deg, #f59e0b, #ea580c)"),
            ("📊", "Data Analyst", "📝 Executive Summariser", "linear-gradient(135deg, #ec4899, #e11d48)"),
            ("📋", "Resume", "📋 Resume Builder", "linear-gradient(135deg, #64748b, #334155)"),
        ]
        
        mcols = st.columns(3)
        for idx, (icon, label, nav_key, bg_grad) in enumerate(modules_grid):
            with mcols[idx % 3]:
                # Render beautiful custom cards with a real Streamlit overlay button
                with st.container():
                    st.markdown(
                        f"""
                        <div style="background: rgba(255,255,255,0.03); border: 1px solid rgba(255,255,255,0.06);
                                    border-radius: 16px; padding: 20px; display: flex; align-items: center; gap: 16px;
                                    transition: transform 0.2s, background 0.2s; cursor: pointer; min-height: 90px;
                                    box-shadow: 0 4px 6px -1px rgba(0, 0, 0, 0.1);">
                            <div style="background: {bg_grad}; width: 48px; height: 48px; border-radius: 12px;
                                        display: flex; align-items: center; justify-content: center; font-size: 1.5rem;
                                        box-shadow: 0 4px 12px rgba(0,0,0,0.3); flex-shrink: 0;">
                                {icon}
                            </div>
                            <div>
                                <div style="font-size: 1rem; font-weight: 700; color: #f1f5f9;">{label}</div>
                                <div style="font-size: 0.75rem; color: #94a3b8; margin-top: 2px;">Open Tool ↗</div>
                            </div>
                        </div>
                        """,
                        unsafe_allow_html=True
                    )
                    # Overlay invisible button to make the entire visual card act as a button
                    if st.button(" ", key=f"quick_btn_{idx}", help=f"Open {label}", use_container_width=True):
                        st.session_state.selected_menu = nav_key
                        st.rerun()

    with right_col:
        # ── Memory / Vector Index Status (Animated) ──
        _card_open("padding: 24px;")
        st.markdown("#### 🗃️ Memory Status")
        vi = st.session_state.get("vector_index")
        doc_count = len(vi.texts) if vi else 0
        mem_pct = min(int((doc_count / 200) * 100), 100)
        status_color = "#10b981" if mem_pct < 60 else ("#f59e0b" if mem_pct < 85 else "#ef4444")
        st.markdown(
            f"""
            <div style="margin-top: 8px;">
                <div style="display:flex; justify-content:space-between; font-size:0.85rem; font-weight: 600;
                            color:var(--text); margin-bottom:8px;">
                    <span>Usage</span>
                    <span style="color:{status_color};">{mem_pct}%</span>
                </div>
                <div style="background:rgba(255,255,255,0.05); border-radius:99px; height:8px; overflow: hidden; position: relative;">
                    <div style="background: {status_color}; width: {max(mem_pct, 2)}%; height: 100%; 
                                border-radius: 99px; transition: width 1s ease-in-out;
                                box-shadow: 0 0 10px {status_color};"></div>
                </div>
                <div style="font-size:0.75rem; color:var(--sub); margin-top:10px; display: flex; justify-content: space-between;">
                    <span>{doc_count} Documents</span>
                    <span>200 Max</span>
                </div>
            </div>
            """,
            unsafe_allow_html=True
        )
        _card_close()

        # ── Recent chat activity ──
        st.markdown("<div style='margin-top:16px;'></div>", unsafe_allow_html=True)
        _card_open("padding: 24px;")
        st.markdown("#### 💬 Recent Activity")
        history = st.session_state.get("chat_history", [])
        if history:
            recents = [m for m in history if m["role"] == "user"][-4:][::-1]
            st.markdown("<div style='display: flex; flex-direction: column; gap: 12px; margin-top: 12px;'>", unsafe_allow_html=True)
            for i, msg in enumerate(recents):
                preview = msg["msg"][:35] + "..." if len(msg["msg"]) > 35 else msg["msg"]
                opacity = 1.0 - (i * 0.2)
                st.markdown(
                    f"""
                    <div style="background: rgba(255,255,255,0.03); border: 1px solid rgba(255,255,255,0.05);
                                padding: 12px; border-radius: 12px; display: flex; align-items: center; gap: 12px;
                                opacity: {opacity}; transition: opacity 0.2s;">
                        <div style="background: rgba(59, 130, 246, 0.2); color: #60a5fa; width: 32px; height: 32px; 
                                    border-radius: 50%; display: flex; align-items: center; justify-content: center; font-size: 0.9rem;">
                            👤
                        </div>
                        <div style="flex: 1; overflow: hidden;">
                            <div style="font-size:0.85rem; color:var(--text); white-space: nowrap; text-overflow: ellipsis; overflow: hidden;">
                                {preview}
                            </div>
                            <div style="font-size:0.7rem; color:var(--sub); margin-top:2px;">Just now</div>
                        </div>
                    </div>
                    """,
                    unsafe_allow_html=True
                )
            st.markdown("</div>", unsafe_allow_html=True)
        else:
            st.info("No activity yet. Use the Quick Ask or EduChat to get started!")
        _card_close()

        # ── Model selector (display only) ──
        st.markdown("<div style='margin-top:16px;'></div>", unsafe_allow_html=True)
        _card_open("padding: 24px;")
        st.markdown("#### 🤖 Engine Status")
        from .config import AVAILABLE_MODELS
        model_options = list(AVAILABLE_MODELS.keys()) if isinstance(AVAILABLE_MODELS, dict) else AVAILABLE_MODELS
        if "dashboard_model" not in st.session_state:
            st.session_state.dashboard_model = model_options[0] if model_options else "llama-3.1-8b-instant"
        
        st.markdown(
            f"""
            <div style="background: rgba(16, 185, 129, 0.1); border: 1px solid rgba(16, 185, 129, 0.2); 
                        padding: 12px; border-radius: 10px; display: flex; align-items: center; gap: 12px; margin-top: 12px;">
                <div style="width: 10px; height: 10px; background: #10b981; border-radius: 50%; box-shadow: 0 0 8px #10b981;"></div>
                <div>
                    <div style="font-size: 0.85rem; font-weight: 700; color: #10b981;">ONLINE & READY</div>
                    <div style="font-size: 0.75rem; color: var(--sub); margin-top: 2px;">Accelerated Inference</div>
                </div>
            </div>
            """, unsafe_allow_html=True
        )
        
        st.markdown("<div style='margin-top:12px; margin-bottom:4px; font-size:0.8rem; color:var(--sub);'>Active Model</div>", unsafe_allow_html=True)
        chosen = st.selectbox(
            "Model",
            model_options,
            index=model_options.index(st.session_state.dashboard_model) if st.session_state.dashboard_model in model_options else 0,
            key="dashboard_model_sel",
            label_visibility="collapsed"
        )
        st.session_state.dashboard_model = chosen
        _card_close()

        # ── EduSphere info card ──
        st.markdown("<div style='margin-top:16px;'></div>", unsafe_allow_html=True)
        _card_open("padding: 24px;")
        st.markdown("#### 🎓 Platform Overview")
        st.markdown(
            """
            <div style="font-size:0.82rem; color:var(--sub); line-height:1.8; margin-top: 12px;">
                <div style="display: flex; align-items: center; gap: 8px; margin-bottom: 4px;">
                    <span style="color: #60a5fa;">✓</span> 14 AI-powered modules
                </div>
                <div style="display: flex; align-items: center; gap: 8px; margin-bottom: 4px;">
                    <span style="color: #60a5fa;">✓</span> FAISS RAG document search
                </div>
                <div style="display: flex; align-items: center; gap: 8px; margin-bottom: 4px;">
                    <span style="color: #60a5fa;">✓</span> Groq LLaMA 3.1 inference
                </div>
                <div style="display: flex; align-items: center; gap: 8px;">
                    <span style="color: #60a5fa;">✓</span> Enterprise-grade security
                </div>
            </div>
            """,
            unsafe_allow_html=True
        )
        _card_close()



# ==============================================================================
# MODULE 1 — RAG Chatbot (with LICT Knowledge Integration)
# ==============================================================================

def render_educhat() -> None:
    """🧠 EduChat & RAG Studio — document-aware conversational AI with LICT campus knowledge."""
    # Inject JavaScript listener in the parent window to handle voice prompts from the iframe
    st.markdown(
        """
        <script>
            if (!window.voicePromptListenerAdded) {
                window.addEventListener('message', function(event) {
                    if (event.data && event.data.type === 'voice_prompt') {
                        const text = event.data.text;
                        const chatInput = document.querySelector('textarea[data-testid="stChatInputTextArea"]') || document.querySelector('.stChatInput textarea');
                        if (chatInput) {
                            const nativeSetter = Object.getOwnPropertyDescriptor(HTMLTextAreaElement.prototype, 'value').set;
                            nativeSetter.call(chatInput, text);
                            chatInput.dispatchEvent(new Event('input', { bubbles: true }));
                            setTimeout(() => {
                                const sendBtn = chatInput.closest('[data-testid="stChatInput"]').querySelector('button');
                                if (sendBtn) {
                                    sendBtn.click();
                                }
                            }, 100);
                        }
                    }
                });
                window.voicePromptListenerAdded = true;
            }
        </script>
        """,
        unsafe_allow_html=True
    )

    st.markdown("### 🧠 EduChat & RAG Search Engine")

    col_chat, col_panel = st.columns([2, 1])

    # ── Right panel: document management ──────────────────────────────────────
    with col_panel:
        _card_open()
        st.markdown("**📄 Document Vectorisation Hub**")

        uploaded_pdf = st.file_uploader(
            "Upload Academic PDF",
            type=["pdf"],
            key="rag_pdf_uploader",
            help="Upload a PDF to enable RAG-powered answers.",
        )

        if uploaded_pdf:
            if st.button("⚡ Index Document", key="btn_index_doc"):
                with logo_spinner("Extracting & Embedding PDF chunks…"):
                    try:
                        chunks = extract_pdf_chunks(uploaded_pdf)
                        vi = _rag_manager.build_index(chunks)
                        st.session_state.vector_index = vi
                        st.success(f"✅ Indexed **{vi.size}** chunks into FAISS!")
                        log.info("Document indexed: %d chunks.", vi.size)
                    except Exception as exc:
                        st.error(f"❌ Indexing failed: {exc}")
                        log.error("Indexing error: %s", exc)

        vi: Optional[VectorIndex] = st.session_state.get("vector_index")
        if vi and vi.size > 0:
            st.info(f"📚 FAISS Active — **{vi.size}** chunks loaded.")

        if st.button("🗑️ Clear Chat History", key="btn_clear_history"):
            st.session_state.chat_history = []
            st.rerun()

        _card_close()

    # ── Left panel: chat interface ─────────────────────────────────────────────
    with col_chat:
        # Unified chat box container (outer card style)
        with st.container(border=True):
            # Scrollable chat history area (no border)
            chat_container = st.container(height=380, border=False)
            with chat_container:
                _render_chat_history()
    
            st.markdown("<hr style='margin: 8px 0; border-top: 1px solid rgba(255, 255, 255, 0.08); border-bottom: none;'>", unsafe_allow_html=True)
    
            # Also check for voice query param
            voice_param = st.query_params.get("voice_prompt", "")
            
            # Display transcribed audio notification if received
            if voice_param:
                st.info(f"🎙️ **Transcribed Audio:** {voice_param}")
                
            prompt = st.chat_input("💬 Ask EduSphere anything…", key="edu_chat_input")
            
            st.components.v1.html(
                """
                <style>
                    html, body {
                        margin: 0;
                        padding: 0;
                        background: transparent;
                        overflow: hidden;
                    }
                    .voice-wave {
                        display: inline-flex;
                        align-items: center;
                        gap: 3px;
                        height: 12px;
                        margin-left: 6px;
                    }
                    .voice-bar {
                        width: 2px;
                        height: 100%;
                        background: #ef4444;
                        border-radius: 1px;
                        animation: bounce 0.8s ease-in-out infinite alternate;
                    }
                    .voice-bar:nth-child(2) { animation-delay: 0.15s; }
                    .voice-bar:nth-child(3) { animation-delay: 0.3s; }
                    .voice-bar:nth-child(4) { animation-delay: 0.45s; }
                    .voice-bar:nth-child(5) { animation-delay: 0.6s; }

                    @keyframes bounce {
                        0% { transform: scaleY(0.3); }
                        100% { transform: scaleY(1.1); }
                    }
                </style>
                <div style="font-family:'Inter',sans-serif; display:flex; flex-direction:column;
                            background-color: rgb(38, 39, 48); border: 1px solid transparent; position: relative;
                            padding: 0.4rem 0.6rem; border-radius: 0.5rem; box-sizing: border-box; width: 100%; margin: 0 auto;">
                    <div style="display:flex; flex-direction:row; align-items:center; gap:0.5rem;">
                        <button type="button" id="mic-btn"
                            style="background:transparent;border:none;font-size:1.15rem;color:#ececf1;
                                   cursor:pointer;outline:none;padding:0;transition:all .2s;display:flex;align-items:center;">🎙️</button>
                        <select id="lang-select"
                            style="background:transparent;color:#ececf1;border:none;
                                   font-size:0.8rem;outline:none;cursor:pointer;padding:1px 2px; width:45px;">
                            <option value="en-US" style="background:#262730;color:#ececf1;">EN</option>
                            <option value="ne-NP" style="background:#262730;color:#ececf1;">NE</option>
                            <option value="hi-IN" style="background:#262730;color:#ececf1;">HI</option>
                            <option value="es-ES" style="background:#262730;color:#ececf1;">ES</option>
                            <option value="fr-FR" style="background:#262730;color:#ececf1;">FR</option>
                        </select>
                        <span id="listening-indicator" style="color: #ef4444; font-size: 0.75rem; font-style: italic; display: none; align-items:center; gap:4px;">
                            🔴 Listening
                            <span class="voice-wave">
                                <span class="voice-bar"></span>
                                <span class="voice-bar"></span>
                                <span class="voice-bar"></span>
                                <span class="voice-bar"></span>
                                <span class="voice-bar"></span>
                            </span>
                        </span>
                    </div>
                    
                    <textarea id="live-transcription-box" placeholder="Transcribed text will appear here..." style="width: 100%; height: 32px; background: rgba(0,0,0,0.25); color: #ececf1; border: 1px solid rgba(255,255,255,0.08); border-radius: 6px; padding: 4px 6px; font-size: 0.85rem; outline: none; resize: none; margin-top: 4px; box-sizing: border-box; font-family: inherit;"></textarea>
                    
                    <div style="display: flex; gap: 8px; margin-top: 4px; justify-content: flex-end;">
                        <button type="button" id="copy-btn" style="background: rgba(255,255,255,0.06); color: #ececf1; border: 1px solid rgba(255,255,255,0.1); border-radius: 4px; padding: 2px 6px; font-size: 0.72rem; cursor: pointer; outline: none;">📋 Copy</button>
                        <button type="button" id="ask-btn" style="background: #00f0ff; color: #0f172a; border: none; border-radius: 4px; padding: 2px 8px; font-size: 0.72rem; font-weight: bold; cursor: pointer; outline: none;">💬 Ask Bot</button>
                    </div>
                </div>
                <script>
                    const micBtn = document.getElementById('mic-btn');
                    const langSel = document.getElementById('lang-select');
                    const indicator = document.getElementById('listening-indicator');
                    const textBox = document.getElementById('live-transcription-box');
                    const copyBtn = document.getElementById('copy-btn');
                    const askBtn = document.getElementById('ask-btn');
            
                    const SR = window.SpeechRecognition || window.webkitSpeechRecognition;
                    if (!SR) {
                        micBtn.style.opacity = '0.3';
                        micBtn.title = 'Speech recognition not supported';
                    } else {
                        const rec = new SR();
                        rec.continuous = false;
                        rec.interimResults = true;
            
                        micBtn.addEventListener('click', () => {
                            rec.lang = langSel.value;
                            rec.start();
                            micBtn.style.color = '#ef4444';
                            micBtn.style.transform = 'scale(1.2)';
                            indicator.style.display = 'inline-flex';
                            textBox.value = '';
                        });
            
                        rec.onresult = (e) => {
                            let text = '';
                            for (let i = e.resultIndex; i < e.results.length; ++i) {
                                text += e.results[i][0].transcript;
                            }
                            textBox.value = text;
                        };
            
                        rec.onspeechend = () => rec.stop();
                        rec.onerror = (e) => {
                            micBtn.style.color = '#ececf1';
                            micBtn.style.transform = 'none';
                            indicator.style.display = 'none';
                        };
                        rec.onend = () => {
                            micBtn.style.color = '#ececf1';
                            micBtn.style.transform = 'none';
                            indicator.style.display = 'none';
                        };
                    }

                    copyBtn.addEventListener('click', () => {
                        textBox.select();
                        navigator.clipboard.writeText(textBox.value).then(() => {
                            copyBtn.textContent = '✅ Copied!';
                            setTimeout(() => { copyBtn.textContent = '📋 Copy'; }, 2000);
                        }).catch(() => {
                            // Fallback
                            document.execCommand('copy');
                            copyBtn.textContent = '✅ Copied!';
                            setTimeout(() => { copyBtn.textContent = '📋 Copy'; }, 2000);
                        });
                    });

                    askBtn.addEventListener('click', () => {
                        const text = textBox.value.trim();
                        if (text) {
                            try {
                                const parentDoc = window.parent.document;
                                const chatInputContainer = parentDoc.querySelector('div[data-testid="stChatInput"]') || parentDoc.querySelector('.stChatInput');
                                if (chatInputContainer) {
                                    const chatInput = chatInputContainer.querySelector('textarea, input');
                                    if (chatInput) {
                                        const proto = chatInput.tagName === 'TEXTAREA' ? window.parent.HTMLTextAreaElement.prototype : window.parent.HTMLInputElement.prototype;
                                        const nativeSetter = Object.getOwnPropertyDescriptor(proto, 'value').set;
                                        nativeSetter.call(chatInput, text);
                                        chatInput.dispatchEvent(new Event('input', { bubbles: true }));
                                        setTimeout(() => {
                                            const sendBtn = chatInputContainer.querySelector('button');
                                            if (sendBtn) {
                                                sendBtn.click();
                                            }
                                        }, 100);
                                        textBox.value = '';
                                    }
                                } else {
                                    fallbackRedirect(text);
                                }
                            } catch (e) {
                                fallbackRedirect(text);
                            }
                        }
                    });

                    function fallbackRedirect(text) {
                        let parentUrl = document.referrer;
                        if (!parentUrl) {
                            try {
                                parentUrl = window.parent.location.href;
                            } catch (e) {
                                parentUrl = window.location.origin;
                            }
                        }
                        const url = new URL(parentUrl);
                        url.searchParams.set('voice_prompt', text);
                        window.parent.location = url.toString();
                        textBox.value = '';
                    }
                </script>
                """,
                height=100,
            )

        # Trigger chat processing if voice prompt was received
        if voice_param and not prompt:
            prompt = voice_param

        if not prompt:
            return

        st.session_state.total_queries = st.session_state.get("total_queries", 0) + 1
        now_time = datetime.datetime.now().strftime("%H:%M")

        # ── Smart greeting detection (LICT Campus greetings) ─────────────────
        greeting_key = detect_greeting(prompt)
        if greeting_key:
            greeting_resp = get_greeting_response(greeting_key)
            st.session_state.chat_history.append({"role": "user", "msg": prompt, "time": now_time})
            st.session_state.chat_history.append({"role": "assistant", "msg": greeting_resp, "time": now_time})
            log_interaction_to_json(
                st.session_state.get("user_info", {}).get("name", "user"),
                "user", prompt,
            )
            log_interaction_to_json(
                st.session_state.get("user_info", {}).get("name", "user"),
                "assistant", greeting_resp,
            )
            if "voice_prompt" in st.query_params:
                st.query_params.clear()
            st.rerun()

        # Guardrails
        guards = evaluate_guardrails(prompt)
        if guards.crisis:
            st.session_state.blocked_count = st.session_state.get("blocked_count", 0) + 1
            st.markdown(
                f'<div class="g-card" style="border-color:#ef4444;">{CRISIS_RESOURCES}</div>',
                unsafe_allow_html=True,
            )
            if "voice_prompt" in st.query_params:
                st.query_params.clear()
            st.rerun()
        if guards.harmful or guards.private:
            st.session_state.blocked_count = st.session_state.get("blocked_count", 0) + 1
            st.error("🚫 Prompt blocked: Violates System Security or Data Privacy Policy.")
            if "voice_prompt" in st.query_params:
                st.query_params.clear()
            st.rerun()

        # Record user message
        st.session_state.chat_history.append(
            {"role": "user", "msg": prompt, "time": now_time}
        )

        # ── LICT Knowledge Base context injection ────────────────────────────
        context = ""
        used_web_search = False
        lict_context = search_lict_knowledge(prompt)

        # RAG context retrieval (document-based)
        if vi and vi.size > 0:
            try:
                context = _rag_manager.similarity_search(prompt, vi)
            except Exception as exc:
                st.warning(f"⚠️ RAG retrieval failed, answering without context: {exc}")
                log.warning("RAG search error: %s", exc)

        # Prepend College knowledge if relevant
        if lict_context:
            context = f"\n\n🏫 **College Knowledge Base:**\n{lict_context}\n\n{context}"

        is_personal_query = any(phrase in prompt.lower() for phrase in ["my name", "who am i", "do you know me", "who i am"])

        if not context and not is_personal_query:
            # Fallback to web search
            with logo_spinner("🔍 Context not found in local documents. Searching the web..."):
                web_results = duckduckgo_search(prompt)
                if web_results:
                    used_web_search = True
                    web_context_pieces = []
                    for idx, res in enumerate(web_results, 1):
                        web_context_pieces.append(
                            f"Source [{idx}]: {res['title']}\n"
                            f"URL: {res['link']}\n"
                            f"Snippet: {res['snippet']}"
                        )
                    context = "\n\n📄 **Live Web Search Results:**\n" + "\n\n".join(web_context_pieces) + "\n"

        full_prompt = f"{context}\n\nUser Question: {prompt}" if context else prompt

        # Get user profile information for personalization (excluding sensitive details like email or ID)
        user_info = st.session_state.get("user_info", {})
        user_name = user_info.get("name", "User")
        user_role = user_info.get("role", "Student")

        custom_ai_name = st.session_state.get("custom_ai_name", "EduSphere AI")

        identity_context = (
            f"\n\nYou are talking to a user named '{user_name}' who holds the role '{user_role}'. "
            f"Your name is currently '{custom_ai_name}'. "
            "If they ask who you are or what your name is, state that your name is 'EduSphere AI' (or your custom name if they previously assigned one), and tell them 'you can call me anything you'd like'. "
            "If the user gives you a new name (e.g. saying 'I will call you Jarvis' or 'your name is now Jarvis'), you must happily accept it, and you MUST include the tag `[SET_AI_NAME: <NewName>]` (for example: `[SET_AI_NAME: Jarvis]`) anywhere in your reply so the system remembers it. "
            "If they ask who they are, what their name is, or if you know them, answer them naturally like a human "
            "using this information. Do NOT under any circumstances disclose any other personal details like their email, "
            "password hash, or user ID."
        )

        # ── Detect intent for specialized handling ─────────────────────────────
        prompt_lower = prompt.lower()
        is_poetry_request = any(w in prompt_lower for w in [
            "poem", "poetry", "rhyme", "rhyming", "verse", "sonnet", "haiku",
            "kavita", "कविता", "شعر", "poème", "gedicht", "poema", "poesia"
        ])
        is_song_request = any(w in prompt_lower for w in [
            "song", "lyrics", "chorus", "verse", "bridge", "anthem",
            "गीत", "chanson", "canção", "lied", "cancion"
        ])
        is_literature_request = any(w in prompt_lower for w in [
            "story", "novel", "essay", "narrative", "prose", "fiction",
            "short story", "write a", "कथा", "निबंध", "histoire", "cuento"
        ])
        is_medical_request = any(w in prompt_lower for w in [
            "medicine", "drug", "medication", "symptom", "disease", "treatment",
            "dose", "tablet", "pill", "prescription", "diagnos", "health",
            "illness", "fever", "pain", "headache", "infection", "antibiotic",
            "doctor", "hospital", "patient", "cure", "remedy", "दवाई", "औषधि"
        ])
        is_code_request = any(w in prompt_lower for w in [
            "code", "program", "function", "algorithm", "script", "debug",
            "python", "java", "javascript", "c++", "c#", "html", "css",
            "sql", "rust", "go", "kotlin", "swift", "php", "ruby"
        ])

        # ── Build specialized system role ───────────────────────────────────────
        universal_base = (
            "You are EduSphere AI — a supremely knowledgeable, world-class AI assistant. "
            "You can expertly answer questions from ANY field: science, math, history, geography, "
            "law, economics, philosophy, psychology, art, music, sports, technology, and more. "
            "Always give accurate, thorough, well-structured responses. "
        )

        if is_poetry_request:
            creative_instruction = (
                "POETRY MODE ACTIVE: You are a master poet fluent in all languages and literary traditions. "
                "Write beautiful, emotionally resonant poetry with STRICT rhyming schemes (AABB, ABAB, or ABCB as appropriate). "
                "Use vivid imagery, metaphors, alliteration, and assonance. "
                "For multi-stanza poems, ensure each stanza has consistent meter and rhyme. "
                "Match the language and cultural style of the user's request. "
                "If writing in non-English, use authentic cultural idioms and poetic traditions. "
            )
        elif is_song_request:
            creative_instruction = (
                "SONG WRITING MODE ACTIVE: You are a world-class songwriter and lyricist. "
                "Structure the song properly with: [Verse 1], [Pre-Chorus] (if applicable), [Chorus], [Verse 2], [Bridge], [Final Chorus]. "
                "Ensure the CHORUS has strong rhyming (AA BB or ABAB) and is emotionally catchy and memorable. "
                "Verses should tell a story with natural flow. The Bridge should contrast emotionally. "
                "Match the genre requested (pop, rock, folk, classical, etc.) or infer from context. "
                "Write in the language of the request with culturally authentic expressions. "
            )
        elif is_literature_request:
            creative_instruction = (
                "LITERATURE MODE ACTIVE: You are a celebrated literary author. "
                "Write with compelling narrative structure, vivid scene-setting, rich character development, "
                "and beautiful prose. Use literary devices: foreshadowing, symbolism, dialogue, sensory details. "
                "Match the tone and style requested (dark, humorous, romantic, philosophical, etc.). "
                "Write in the language of the user's request. "
            )
        elif is_medical_request:
            creative_instruction = (
                "MEDICAL ADVISORY MODE: You are a knowledgeable medical information assistant. "
                "Provide general information about symptoms, conditions, common medications, dosages, and treatments "
                "based on established medical knowledge. "
                "ALWAYS include this disclaimer at the end: '⚠️ MEDICAL DISCLAIMER: This information is for educational purposes only. "
                "Always consult a qualified doctor or healthcare professional before taking any medication or making health decisions.' "
                "Do NOT diagnose definitively — say 'This may indicate...' or 'Common causes include...'. "
                "Mention both generic and brand names of medicines where applicable. "
            )
        elif is_code_request:
            creative_instruction = (
                "CODE EXPERT MODE: You are a senior software engineer with expertise in all programming languages. "
                "Write clean, efficient, well-commented code. Always:"
                "1. Add clear inline comments explaining logic "
                "2. Include example usage or test cases "
                "3. Handle edge cases and errors "
                "4. Follow best practices for the language "
                "5. Explain the code after writing it "
                "Format code in proper markdown code blocks with language identifier. "
            )
        else:
            creative_instruction = (
                "UNIVERSAL EXPERT MODE: Answer questions from any field with depth and precision. "
                "Structure your response clearly with headings, bullet points, or numbered lists as appropriate. "
                "Provide examples, analogies, and context to make complex topics accessible. "
            )

        if used_web_search:
            system_role = (
                universal_base
                + creative_instruction
                + "Answer based on the Live Web Search Results context provided. "
                "Cite your web sources (e.g. Source [1], Source [2], etc.) and links in the response. "
                "CRITICAL: You MUST write your ENTIRE response in the same language as the user's question. "
                "If the user writes in Nepali, respond in Nepali. If French, respond in French, etc."
                f"{identity_context}"
            )
        else:
            system_role = (
                universal_base
                + creative_instruction
                + "Answer based on context provided if available. "
                "If college knowledge context is provided, ground your answer in it. "
                "If not, answer from your vast training knowledge. "
                "CRITICAL: You MUST write your ENTIRE response in the same language as the user's question. "
                "If the user writes in Nepali, respond in Nepali. If French, respond in French. Never switch languages."
                f"{identity_context}"
            )

        # ── Build conversation history for memory-aware response ─────────────
        # Include up to last 20 turns (40 messages) so the AI remembers context
        MAX_HISTORY_TURNS = 20
        prior_history = st.session_state.get("chat_history", [])[:-1]  # exclude the just-added user msg
        recent_history = prior_history[-(MAX_HISTORY_TURNS * 2):]  # last N turn pairs

        messages_for_llm = []
        for h in recent_history:
            role = "user" if h["role"] == "user" else "assistant"
            messages_for_llm.append({"role": role, "content": h["msg"]})

        # Append the current user message (with context injected)
        messages_for_llm.append({"role": "user", "content": full_prompt})

        # Collect response using memory-aware history call
        with logo_spinner("🤖 EduSphere is thinking…"):
            response_text = groq_chat_with_history(messages_for_llm, system=system_role)

        # Parse custom name request from response if present
        import re
        name_match = re.search(r'\[SET_AI_NAME:\s*([^\]]+)\]', response_text)
        if name_match:
            st.session_state.custom_ai_name = name_match.group(1).strip()
            response_text = re.sub(r'\[SET_AI_NAME:\s*([^\]]+)\]', '', response_text).strip()

        st.session_state.chat_history.append(
            {"role": "assistant", "msg": response_text, "time": now_time}
        )

        # Log interaction
        log_interaction_to_json(
            st.session_state.get("user_info", {}).get("name", "user"),
            "user", prompt,
        )
        log_interaction_to_json(
            st.session_state.get("user_info", {}).get("name", "user"),
            "assistant", response_text,
        )

        if "voice_prompt" in st.query_params:
            st.query_params.clear()
        st.rerun()


# ==============================================================================
# MODULE 2 — Study Planner
# ==============================================================================

def render_study_planner() -> None:
    """📚 Automated Study Plan Generator."""
    st.markdown("### 📚 Automated Study Plan Generator")
    col1, col2 = st.columns([1, 1])

    with col1:
        _card_open()
        topic = st.text_input("Target Subject / Course", "Quantum Mechanics & Computing")
        days = st.slider("Preparation Duration (Days)", 3, 30, 7)
        hours_per_day = st.number_input("Study Hours / Day", min_value=1, max_value=12, value=4)
        level = st.selectbox("Current Knowledge Level", ["Beginner", "Intermediate", "Advanced"])
        _card_close()

    with col2:
        if st.button("🚀 Generate Study Roadmap", key="btn_study_plan"):
            if not topic.strip():
                st.warning("Please enter a topic.")
                return
            prompt = (
                f"Generate a structured {days}-day study plan for '{topic}'. "
                f"Daily budget: {hours_per_day} hours. Target level: {level}. "
                "Break into daily topics, key resources, and revision checkpoints. "
                "Use clear headings and bullet points."
            )
            with logo_spinner("Synthesising optimal study plan…"):
                result = groq_chat(prompt, system="You are an expert curriculum designer.")

            _card_open()
            st.markdown(result)
            _card_close()
            _render_export_buttons(result, "study_plan")


# ==============================================================================
# MODULE 3 — Socratic Clarifier
# ==============================================================================

def render_socratic_clarifier() -> None:
    """🔬 Socratic Concept Deconstructor."""
    st.markdown("### 🔬 Socratic Concept Deconstructor")
    st.write("Deconstruct complex theories through guided Socratic questioning.")

    concept = st.text_input("Enter Concept to Master", "Fourier Transforms")

    if st.button("Deconstruct Concept", key="btn_socratic"):
        if not concept.strip():
            st.warning("Please enter a concept.")
            return
        prompt = (
            f"Explain '{concept}' using the Socratic method:\n"
            "1. **Core Intuition** – provide a real-world analogy.\n"
            "2. **Mathematical / Logical Breakdown** – step-by-step derivation.\n"
            "3. **Three Socratic Probing Questions** – to test deep comprehension.\n"
            "4. **Common Misconceptions** – pitfalls learners should avoid."
        )
        with logo_spinner("Deconstructing concept…"):
            result = groq_chat(prompt, system="You are a Socratic Academic Mentor.")
        _card_open()
        st.markdown(result)
        _card_close()
        _render_export_buttons(result, "socratic_analysis")


# ==============================================================================
# MODULE 4 — Quiz Generator
# ==============================================================================

def render_quiz_generator() -> None:
    """🧪 Automated Quiz & University Exam Paper Generator."""
    st.markdown("### 🧪 Automated Quiz & University Exam Generator")
    st.markdown(
        '<div style="color:var(--sub); font-size:0.88rem; margin-bottom:16px;">'
        "Generate formal University Examination Question Papers, Class Practice Papers, or Quick Quizzes. "
        "Customize university branding, structural guidelines, and download print-ready PDFs."
        "</div>",
        unsafe_allow_html=True,
    )

    col1, col2 = st.columns([1.2, 1], gap="medium")

    with col1:
        _card_open()
        st.markdown("#### 🎓 Exam Configuration")
        
        exam_mode = st.selectbox(
            "📝 Assessment Mode",
            [
                "🏫 Formal University Exam Paper",
                "📖 Classroom Practice Paper",
                "⚡ Quick Self-Assessment Quiz"
            ]
        )
        
        quiz_topic = st.text_input("Course Subject / Specific Topics", "Quantum Computing & Cryptography")
        
        difficulty = st.selectbox(
            "🎓 Academic Standard / Target Grade",
            options=[
                "Primary School (Grade 1-5)",
                "Middle School (Grade 6-8)",
                "High School (Grade 9-12)",
                "Undergraduate (BSc/BE/BA)",
                "Postgraduate (MSc/PhD)",
                "Job Interview Preparation"
            ],
            index=3,
        )

        nepal_faculty = st.selectbox(
            "🏛️ Nepal Faculty / Board Affiliation",
            options=[
                "None / General Assessment",
                "Tribhuvan University (TU - IOE/IOM/FOHSS/FMS)",
                "Kathmandu University (KU - SOE/SOM/SOS)",
                "Pokhara University (PoU)",
                "Purbanchal University (PU)",
                "National Examinations Board (NEB - Grade 11/12)",
                "CTEVT (Diploma/Technical)"
            ],
            index=1
        )
        _card_close()

    with col2:
        _card_open()
        st.markdown("#### ⚙️ Structural Details")
        
        col_sub1, col_sub2 = st.columns(2)
        with col_sub1:
            num_q = st.slider("Total Questions", 3, 20, 8)
            time_limit = st.selectbox("Exam Time Limit", ["No Limit", "45 Minutes", "90 Minutes", "3 Hours"])
        with col_sub2:
            q_type = st.selectbox(
                "Question Style", 
                [
                    "Mixed Standards (Long & Short & MCQs)", 
                    "Long Questions (10 Marks each)", 
                    "Short Questions (5 Marks each)",
                    "Long Subjective / Analytical Only", 
                    "Multiple Choice (A–D) Only", 
                    "Short Answers Only"
                ]
            )
            hardness = st.select_slider(
                "🔥 Difficulty Rating",
                options=["Easy", "Medium", "Hard", "Expert / Analytical"],
                value="Hard"
            )
            include_keys = st.checkbox("🔑 Include Answer Key & Details", value=True)

        univ_name = st.text_input("University / Institutional Title", "EDUSPHERE ACADEMIC BOARD")
        _card_close()

    st.markdown("<div style='margin-top:14px;'></div>", unsafe_allow_html=True)

    if st.button("🎯 Generate University Exam Paper", key="btn_quiz", use_container_width=True, type="primary"):
        if not quiz_topic.strip():
            st.warning("⚠️ Please specify a course subject or syllabus topic first.")
            return

        with logo_spinner("✍️ Formulating exam paper structure and answers..."):
            # Structure prompt specifically for high-standard university tests
            prompt = f"""
Create a formal academic {exam_mode} on the subject/topic: '{quiz_topic}'.
Target Grade/Level: {difficulty}.
Nepal Board/Faculty Affiliation: {nepal_faculty}.
Difficulty/Hardness Rating: {hardness}.
Total Questions: {num_q}.
Time Limit: {time_limit}.
Question Style: {q_type}.
Institution Name: {univ_name.upper()}.

Format Requirements:
1. Start with a clean Header including Time Allowed, Board Affiliation ({nepal_faculty}), Difficulty level, Target Standard, and general Instructions.
2. Structure the questions clearly under sections matching the selected level (e.g., SECTION A: Objective, SECTION B: Subjective).
3. Adjust question complexity to match both the level ({difficulty}), board guidelines ({nepal_faculty}), and difficulty rating ({hardness}).
4. {"IMPORTANT: After the question paper, provide a clean page break marker '---' followed by a comprehensive 'ANSWER KEY & EXPLANATIONS' section." if include_keys else "Do not include the answers."}

Ensure appropriate complexity, clarity, and precise scientific terminology matching Nepali educational board standards where applicable. Return Markdown.
"""
            result = groq_chat(prompt, system="You are an Academic Registrar and University Course Evaluator.")

            # Set up session state variables for rendered outputs
            st.session_state.current_quiz_text = result
            st.session_state.current_quiz_title = f"{univ_name.replace(' ', '_')}_{quiz_topic.replace(' ', '_')}"

    # Show generated quiz if it exists in session state
    if "current_quiz_text" in st.session_state:
        st.markdown("### 📝 Generated Question Paper Preview")
        _card_open()
        st.markdown(st.session_state.current_quiz_text)
        _card_close()

        # Render custom export selector for University PDFs
        _card_open("margin-top: 14px;")
        st.markdown("#### 💾 Download Exam Paper")
        
        col_sel, col_btn = st.columns([2, 1], vertical_alignment="bottom")
        with col_sel:
            export_format = st.selectbox(
                "Export Format Presets",
                ["Formal Exam PDF (.pdf)", "Plain Text Document (.txt)"],
                key="quiz_export_selector"
            )
        with col_btn:
            if export_format == "Plain Text Document (.txt)":
                st.download_button(
                    label="⬇️ Download Document",
                    data=st.session_state.current_quiz_text,
                    file_name=f"{st.session_state.current_quiz_title}.txt",
                    mime="text/plain",
                    use_container_width=True,
                    type="primary"
                )
            else:
                # Custom high-quality PDF generator with Exam Header overlays
                from .utils import generate_pdf_bytes
                
                # Prepend formal header before generating PDF bytes
                import datetime as dtm
                header_prefix = (
                    f"=========================================================\n"
                    f"               {univ_name.upper()}\n"
                    f"               OFFICIAL UNIVERSITY EVALUATION\n"
                    f"---------------------------------------------------------\n"
                    f"Course: {quiz_topic}                  Standard: {difficulty}\n"
                    f"Allowed Time: {time_limit}            Date: {dtm.datetime.now().strftime('%d %b %Y')}\n"
                    f"=========================================================\n\n"
                )
                
                # If we have an answer key, make it separate
                full_content_for_pdf = header_prefix + st.session_state.current_quiz_text
                pdf_bytes = generate_pdf_bytes(full_content_for_pdf)
                
                st.download_button(
                    label="⬇️ Download Print-Ready PDF",
                    data=pdf_bytes,
                    file_name=f"{st.session_state.current_quiz_title}.pdf",
                    mime="application/pdf",
                    use_container_width=True,
                    type="primary"
                )
        _card_close()

    # ── NEW: .txt to PDF Document Converter Workspace ───────────────────────
    st.markdown("<hr style='margin: 20px 0; border: none; border-top: 1px solid rgba(255,255,255,0.08);'>", unsafe_allow_html=True)
    st.markdown("### 📄 Document Converter (.txt ➔ .pdf)")
    st.markdown(
        '<div style="color:var(--sub); font-size:0.85rem; margin-bottom:12px;">'
        "Upload any plain text file or syllabus paper to compile and transform it into a formatted, Latin1-sanitized PDF document."
        "</div>",
        unsafe_allow_html=True,
    )
    
    conv_col1, conv_col2 = st.columns([1, 1], gap="medium")
    with conv_col1:
        _card_open()
        st.markdown("##### 📥 Upload Text File")
        uploaded_txt = st.file_uploader("Choose a plain text file (.txt)", type=["txt"], key="txt_to_pdf_uploader")
        _card_close()
        
    with conv_col2:
        if uploaded_txt is not None:
            _card_open()
            st.markdown("##### ⚙️ Compilation Options")
            # Read uploaded file content
            txt_string = uploaded_txt.read().decode("utf-8", errors="replace")
            
            # Show premium styled scrollable document preview
            st.markdown("##### 📝 Document Content Preview")
            st.markdown(
                f'<div style="background: rgba(255, 255, 255, 0.02); border: 1px solid rgba(255, 255, 255, 0.08); border-radius: 8px; padding: 12px; height: 180px; overflow-y: auto; font-family: \'Space Grotesk\', sans-serif; font-size: 0.85rem; color: var(--text-color); margin-bottom: 12px; white-space: pre-wrap; line-height: 1.5;">'
                f'{txt_string}'
                f'</div>',
                unsafe_allow_html=True
            )
            
            # Download PDF Conversion Trigger
            from .utils import generate_pdf_bytes
            converted_pdf = generate_pdf_bytes(txt_string)
            orig_name = uploaded_txt.name.rsplit(".", 1)[0]
            
            st.markdown("<div style='margin-top:10px;'></div>", unsafe_allow_html=True)
            st.download_button(
                label="📥 Convert & Download PDF",
                data=converted_pdf,
                file_name=f"{orig_name}_converted.pdf",
                mime="application/pdf",
                use_container_width=True,
                type="primary"
            )
            _card_close()
        else:
            _card_open("height: 100%; display: flex; align-items: center; justify-content: center; text-align: center;")
            st.markdown(
                '<div style="color:var(--sub); font-size:0.86rem; padding: 24px 0;">'
                "Waiting for file upload...<br>Converted PDF will be instantly available for download here."
                "</div>",
                unsafe_allow_html=True
            )
            _card_close()

# ==============================================================================
# MODULE 5 — Code Lab
# ==============================================================================

def render_code_lab() -> None:
    """💻 Intelligent Code Debugger & Explainer."""
    st.markdown("### 💻 Intelligent Code Debugger & Explainer")

    code_snippet = st.text_area(
        "Paste Code Snippet",
        value=(
            "def binary_search(arr, target):\n"
            "    # Bug: high should be len(arr) - 1\n"
            "    low, high = 0, len(arr)\n"
            "    while low < high:\n"
            "        mid = (low + high) // 2\n"
            "        if arr[mid] == target: return mid\n"
            "        elif arr[mid] < target: low = mid   # Bug: should be mid + 1\n"
            "        else: high = mid\n"
            "    return -1"
        ),
        height=180,
    )

    lang = st.selectbox("Language", ["Python", "JavaScript", "Java", "C++", "Go", "Rust"])
    task = st.selectbox(
        "Analysis Task",
        [
            "Explain Line-by-Line",
            "Find & Fix All Bugs",
            "Optimise Time & Space Complexity",
            f"Translate to C++",
            "Write Unit Tests",
            "Add Type Hints / Annotations",
        ],
    )

    if st.button("⚡ Execute Analysis", key="btn_code_lab"):
        if not code_snippet.strip():
            st.warning("Please paste a code snippet.")
            return
        prompt = (
            f"Task: {task}\n\n"
            f"Language: {lang}\n\n"
            f"Code:\n```{lang.lower()}\n{code_snippet}\n```\n\n"
            "Provide a thorough, structured analysis with corrected code where applicable."
        )
        with logo_spinner("Analysing code structure…"):
            result = groq_chat(prompt, system="You are an expert Software Engineer & CS Professor.")
        _card_open()
        st.markdown(result)
        _card_close()
        _render_export_buttons(result, "code_analysis")


# ==============================================================================
# MODULE 6 — Academic Translator
# ==============================================================================

def render_translator() -> None:
    """🌍 Multi-Lingual Academic Translator."""
    st.markdown("### 🌍 Multi-Lingual Academic Translator")

    # ── Voice Input handling for Translator ──
    translate_voice = st.query_params.get("translate_voice_prompt", "")
    default_text = "Neural networks utilise backpropagation to update weights based on gradient loss."
    should_auto_translate = False
    
    if "translate_input" not in st.session_state:
        st.session_state.translate_input = default_text
        
    if translate_voice:
        st.session_state.translate_input = translate_voice
        should_auto_translate = True
        st.query_params.clear()

    # ── Translator Voice Input component ──
    st.components.v1.html(
        """
        <style>
            html, body {
                margin: 0;
                padding: 0;
                background: transparent;
                overflow: hidden;
            }
            .voice-wave {
                display: inline-flex;
                align-items: center;
                gap: 3px;
                height: 12px;
                margin-left: 6px;
            }
            .voice-bar {
                width: 2px;
                height: 100%;
                background: #ef4444;
                border-radius: 1px;
                animation: bounce 0.8s ease-in-out infinite alternate;
            }
            .voice-bar:nth-child(2) { animation-delay: 0.15s; }
            .voice-bar:nth-child(3) { animation-delay: 0.3s; }
            .voice-bar:nth-child(4) { animation-delay: 0.45s; }
            .voice-bar:nth-child(5) { animation-delay: 0.6s; }

            @keyframes bounce {
                0% { transform: scaleY(0.3); }
                100% { transform: scaleY(1.1); }
            }
        </style>
        <div style="font-family:'Inter',sans-serif; display:flex; flex-direction:column;
                    background-color: rgb(38, 39, 48); border: 1px solid transparent; position: relative;
                    padding: 0.4rem 0.6rem; border-radius: 0.5rem; box-sizing: border-box; width: 100%; margin-bottom: 6px;">
            <div style="display:flex; flex-direction:row; align-items:center; gap:0.5rem;">
                <button type="button" id="trans-mic-btn"
                    style="background:transparent;border:none;font-size:1.15rem;color:#ececf1;
                           cursor:pointer;outline:none;padding:0;transition:all .2s;display:flex;align-items:center;gap:6px;">🎙️ Speak to Translate</button>
                <select id="trans-lang-select"
                    style="background:transparent;color:#ececf1;border:none;
                           font-size:0.8rem;outline:none;cursor:pointer;padding:1px 2px; width:45px;">
                    <option value="en-US" style="background:#262730;color:#ececf1;">EN</option>
                    <option value="ne-NP" style="background:#262730;color:#ececf1;">NE</option>
                    <option value="hi-IN" style="background:#262730;color:#ececf1;">HI</option>
                    <option value="es-ES" style="background:#262730;color:#ececf1;">ES</option>
                    <option value="fr-FR" style="background:#262730;color:#ececf1;">FR</option>
                </select>
                <span id="trans-indicator" style="color: #ef4444; font-size: 0.75rem; font-style: italic; display: none; align-items:center; gap:4px;">
                    🔴 Listening
                    <span class="voice-wave">
                        <span class="voice-bar"></span>
                        <span class="voice-bar"></span>
                        <span class="voice-bar"></span>
                        <span class="voice-bar"></span>
                        <span class="voice-bar"></span>
                    </span>
                </span>
            </div>
            
            <textarea id="trans-transcription-box" placeholder="Transcribed text will appear here..." style="width: 100%; height: 32px; background: rgba(0,0,0,0.25); color: #ececf1; border: 1px solid rgba(255,255,255,0.08); border-radius: 6px; padding: 4px 6px; font-size: 0.85rem; outline: none; resize: none; margin-top: 4px; box-sizing: border-box; font-family: inherit;"></textarea>
            
            <div style="display: flex; gap: 8px; margin-top: 4px; justify-content: flex-end;">
                <button type="button" id="trans-copy-btn" style="background: rgba(255,255,255,0.06); color: #ececf1; border: 1px solid rgba(255,255,255,0.1); border-radius: 4px; padding: 2px 6px; font-size: 0.72rem; cursor: pointer; outline: none;">📋 Copy</button>
                <button type="button" id="trans-submit-btn" style="background: #00f0ff; color: #0f172a; border: none; border-radius: 4px; padding: 2px 8px; font-size: 0.72rem; font-weight: bold; cursor: pointer; outline: none;">🌐 Translate</button>
            </div>
        </div>
        <script>
            const micBtn = document.getElementById('trans-mic-btn');
            const langSel = document.getElementById('trans-lang-select');
            const indicator = document.getElementById('trans-indicator');
            const textBox = document.getElementById('trans-transcription-box');
            const copyBtn = document.getElementById('trans-copy-btn');
            const submitBtn = document.getElementById('trans-submit-btn');
    
            const SR = window.SpeechRecognition || window.webkitSpeechRecognition;
            if (!SR) {
                micBtn.style.opacity = '0.3';
                micBtn.title = 'Speech recognition not supported';
            } else {
                const rec = new SR();
                rec.continuous = false;
                rec.interimResults = true;
    
                micBtn.addEventListener('click', () => {
                    rec.lang = langSel.value;
                    rec.start();
                    micBtn.style.color = '#ef4444';
                    indicator.style.display = 'inline-flex';
                    textBox.value = '';
                });
    
                rec.onresult = (e) => {
                    let text = '';
                    for (let i = e.resultIndex; i < e.results.length; ++i) {
                        text += e.results[i][0].transcript;
                    }
                    textBox.value = text;
                };
    
                rec.onspeechend = () => rec.stop();
                rec.onerror = (e) => {
                    micBtn.style.color = '#ececf1';
                    indicator.style.display = 'none';
                };
                rec.onend = () => {
                    micBtn.style.color = '#ececf1';
                    indicator.style.display = 'none';
                };
            }

            copyBtn.addEventListener('click', () => {
                textBox.select();
                navigator.clipboard.writeText(textBox.value).then(() => {
                    copyBtn.textContent = '✅ Copied!';
                    setTimeout(() => { copyBtn.textContent = '📋 Copy'; }, 2000);
                }).catch(() => {
                    document.execCommand('copy');
                    copyBtn.textContent = '✅ Copied!';
                    setTimeout(() => { copyBtn.textContent = '📋 Copy'; }, 2000);
                });
            });

            submitBtn.addEventListener('click', () => {
                const text = textBox.value.trim();
                if (text) {
                    try {
                        const parentDoc = window.parent.document;
                        const transInput = parentDoc.querySelector('textarea[aria-label="Source Text"]') || parentDoc.querySelector('.stTextArea textarea');
                        if (transInput) {
                            const nativeSetter = Object.getOwnPropertyDescriptor(window.parent.HTMLTextAreaElement.prototype, 'value').set;
                            nativeSetter.call(transInput, text);
                            transInput.dispatchEvent(new Event('input', { bubbles: true }));
                            setTimeout(() => {
                                const buttons = Array.from(parentDoc.querySelectorAll('button'));
                                const translateBtn = buttons.find(btn => btn.textContent.includes('Translate Text'));
                                if (translateBtn) {
                                    translateBtn.click();
                                }
                            }, 100);
                            textBox.value = '';
                        } else {
                            fallbackRedirect(text);
                        }
                    } catch (e) {
                        fallbackRedirect(text);
                    }
                }
            });

            function fallbackRedirect(text) {
                let parentUrl = document.referrer;
                if (!parentUrl) {
                    try {
                        parentUrl = window.parent.location.href;
                    } catch (e) {
                        parentUrl = window.location.origin;
                    }
                }
                const url = new URL(parentUrl);
                url.searchParams.set('translate_voice_prompt', text);
                window.parent.location = url.toString();
                textBox.value = '';
            }
        </script>
        """,
        height=100,
    )

    input_text = st.text_area(
        "Source Text",
        value=st.session_state.translate_input,
        height=120,
        key="translator_source_text"
    )
    
    # Sync typed text back to state
    st.session_state.translate_input = input_text

    col1, col2 = st.columns([1, 1])
    with col1:
        target_lang = st.selectbox(
            "Target Language",
            ["Nepali", "Spanish", "French", "German", "Chinese (Mandarin)", "Japanese", "Hindi", "Arabic"],
        )
    with col2:
        formality = st.selectbox("Formality Level", ["Academic / Formal", "Conversational", "Technical"])

    if st.button("🌐 Translate Text", key="btn_translate") or should_auto_translate:
        if not input_text.strip():
            st.warning("Please provide source text.")
            return
        prompt = (
            f"Translate the following academic text into {target_lang} "
            f"using a {formality} tone. Preserve specialised terminology. "
            "If any terms have no direct equivalent, provide a footnote explanation.\n\n"
            f"Text:\n{input_text}"
        )
        with logo_spinner("Translating…"):
            result = groq_chat(prompt)
        _card_open()
        st.subheader(f"Translation → {target_lang}")
        st.write(result)
        _card_close()
        _render_export_buttons(result, "translation")


# ==============================================================================
# MODULE 7 — Executive Summariser
# ==============================================================================

def render_summariser() -> None:
    """📝 Academic Text & Paper Summariser."""
    st.markdown("### 📝 Executive Summariser & Data Analyst")

    tab1, tab2 = st.tabs(["📝 Text Summariser", "📊 Data Analyst"])

    with tab1:
        long_text = st.text_area("Input Academic Passage", height=220)
        col1, col2 = st.columns([1, 1])
        with col1:
            format_type = st.selectbox(
                "Summary Format",
                ["Bullet Core Takeaways", "Executive One-Pager", "Abstract Style", "ELI5 (Simple Explanation)"],
            )
        with col2:
            word_limit = st.slider("Target Word Count", 50, 500, 150)

        if st.button("📑 Summarise", key="btn_summarise"):
            if not long_text.strip():
                st.warning("Please provide input text.")
                return
            prompt = (
                f"Summarise the following text using the '{format_type}' format. "
                f"Target approximately {word_limit} words. "
                "Be precise and preserve key academic insights.\n\n"
                f"Text:\n{long_text}"
            )
            with logo_spinner("Summarising…"):
                result = groq_chat(prompt)
            _card_open()
            st.markdown(result)
            _card_close()
            _render_export_buttons(result, "summary")

    with tab2:
        st.markdown("Upload a dataset for automated comprehensive AI analysis and visualization.")
        uploaded_file = st.file_uploader("Upload Data (CSV, JSON, XLSX, TXT, MD)", type=["csv", "json", "xlsx", "xls", "txt", "md"])
        
        if uploaded_file is not None:
            try:
                import pandas as pd
                import matplotlib.pyplot as plt
                import io
                
                # Load the dataset
                file_extension = uploaded_file.name.split(".")[-1].lower()
                
                with logo_spinner("Parsing dataset..."):
                    if file_extension == "csv":
                        df = pd.read_csv(uploaded_file)
                    elif file_extension in ["xls", "xlsx"]:
                        df = pd.read_excel(uploaded_file)
                    elif file_extension == "json":
                        df = pd.read_json(uploaded_file)
                    elif file_extension in ["txt", "md"]:
                        df = pd.read_csv(uploaded_file, sep="\t") # Assume tab-separated for text files as fallback
                    else:
                        st.error("Unsupported file format.")
                        return

                st.success("✅ Dataset loaded successfully!")
                
                with st.expander("🔍 Dataset Preview"):
                    st.dataframe(df.head(10), use_container_width=True)
                    st.caption(f"Rows: {df.shape[0]} | Columns: {df.shape[1]}")
                
                if st.button("🚀 Generate Comprehensive Analysis", key="btn_data_analyst"):
                    # 1. Prepare statistical summary for the LLM (Security: don't send raw data)
                    desc = df.describe(include='all').to_string()
                    missing_info = df.isnull().sum().to_string()
                    types_info = df.dtypes.to_string()
                    
                    data_context = (
                        f"Dataset Shape: {df.shape[0]} rows, {df.shape[1]} columns.\n\n"
                        f"Data Types:\n{types_info}\n\n"
                        f"Missing Values:\n{missing_info}\n\n"
                        f"Statistical Summary:\n{desc}"
                    )
                    
                    # 2. Get AI Analysis
                    prompt = (
                        "You are an expert Data Analyst AI. I am providing you with the statistical summary and schema of a dataset (not the raw data). "
                        "Please provide a comprehensive, structured data analysis report. "
                        "Include: \n"
                        "- High-level overview of the data\n"
                        "- Key statistical insights and anomalies\n"
                        "- Potential implications or trends based on the statistics\n"
                        "- Recommendations for further investigation\n\n"
                        f"Data Summary Context:\n{data_context}"
                    )
                    
                    with logo_spinner("AI is analyzing the data..."):
                        analysis_result = groq_chat(prompt)
                    
                    # 3. Generate Charts (Matplotlib)
                    with logo_spinner("Generating charts..."):
                        image_bytes_list = []
                        num_cols = df.select_dtypes(include=['number']).columns.tolist()
                        cat_cols = df.select_dtypes(include=['object', 'category']).columns.tolist()
                        
                        _card_open()
                        st.markdown(analysis_result)
                        _card_close()
                        
                        st.markdown("### 📈 Visualizations")
                        col1, col2 = st.columns(2)
                        
                        plots_generated = 0
                        # Try to set a modern style if available
                        try:
                            plt.style.use("dark_background")
                        except:
                            pass
                        
                        # Plot 1: Correlation Matrix (if > 1 numeric col)
                        if len(num_cols) > 1:
                            fig1, ax1 = plt.subplots(figsize=(6, 4))
                            corr = df[num_cols].corr()
                            cax = ax1.matshow(corr, cmap='coolwarm')
                            fig1.colorbar(cax)
                            ax1.set_xticks(range(len(num_cols)))
                            ax1.set_yticks(range(len(num_cols)))
                            ax1.set_xticklabels(num_cols, rotation=45, ha='left')
                            ax1.set_yticklabels(num_cols)
                            ax1.set_title("Correlation Heatmap", pad=20)
                            fig1.tight_layout()
                            
                            buf1 = io.BytesIO()
                            fig1.savefig(buf1, format="png", dpi=150)
                            buf1.seek(0)
                            image_bytes_list.append(buf1.read())
                            col1.image(buf1, caption="Correlation Heatmap")
                            plt.close(fig1)
                            plots_generated += 1
                            
                        # Plot 2: Distribution of first numeric col
                        if len(num_cols) > 0:
                            fig2, ax2 = plt.subplots(figsize=(6, 4))
                            col_name = num_cols[0]
                            ax2.hist(df[col_name].dropna(), bins=20, color='skyblue', edgecolor='black')
                            ax2.set_title(f"Distribution of {col_name}")
                            fig2.tight_layout()
                            
                            buf2 = io.BytesIO()
                            fig2.savefig(buf2, format="png", dpi=150)
                            buf2.seek(0)
                            image_bytes_list.append(buf2.read())
                            
                            target_col = col1 if plots_generated % 2 == 0 else col2
                            target_col.image(buf2, caption=f"Distribution of {col_name}")
                            plt.close(fig2)
                            plots_generated += 1
                            
                        # Plot 3: Distribution of second numeric col OR first categorical
                        if len(num_cols) > 1:
                            fig3, ax3 = plt.subplots(figsize=(6, 4))
                            col_name = num_cols[1]
                            ax3.boxplot(df[col_name].dropna())
                            ax3.set_title(f"Boxplot of {col_name}")
                            fig3.tight_layout()
                            
                            buf3 = io.BytesIO()
                            fig3.savefig(buf3, format="png", dpi=150)
                            buf3.seek(0)
                            image_bytes_list.append(buf3.read())
                            
                            target_col = col1 if plots_generated % 2 == 0 else col2
                            target_col.image(buf3, caption=f"Boxplot of {col_name}")
                            plt.close(fig3)
                            plots_generated += 1
                        elif len(cat_cols) > 0:
                            fig3, ax3 = plt.subplots(figsize=(6, 4))
                            col_name = cat_cols[0]
                            counts = df[col_name].value_counts().head(10)
                            counts.plot(kind='bar', ax=ax3, color='lightgreen')
                            ax3.set_title(f"Top 10 categories in {col_name}")
                            fig3.tight_layout()
                            
                            buf3 = io.BytesIO()
                            fig3.savefig(buf3, format="png", dpi=150)
                            buf3.seek(0)
                            image_bytes_list.append(buf3.read())
                            
                            target_col = col1 if plots_generated % 2 == 0 else col2
                            target_col.image(buf3, caption=f"Bar chart of {col_name}")
                            plt.close(fig3)
                            plots_generated += 1
                        
                        if plots_generated == 0:
                            st.info("Dataset doesn't have sufficient numeric/categorical variation for automatic plotting.")
                        
                        # 4. Provide PDF Export
                        st.markdown("---")
                        from .utils import get_analytical_pdf_download_link
                        
                        download_html = get_analytical_pdf_download_link(
                            content=analysis_result,
                            image_bytes_list=image_bytes_list,
                            filename="data_analysis_report.pdf",
                            label="Download Data Analyst PDF Report"
                        )
                        st.markdown(download_html, unsafe_allow_html=True)

            except Exception as e:
                log.error("Data Analyst error: %s", e)
                st.error(f"❌ Error processing dataset: {e}")


# ==============================================================================
# MODULE 8 — URL & Visual Intelligence
# ==============================================================================

def render_url_intelligence() -> None:
    """🖼️ Web & Image Scraping Engine."""
    st.markdown("### 🖼️ Web Content Intelligence Engine")

    target_url = st.text_input(
        "Web URL to Analyse",
        "https://en.wikipedia.org/wiki/Artificial_intelligence",
    )

    if st.button("🔍 Analyse URL", key="btn_analyse_url"):
        if not target_url.strip():
            st.warning("Please enter a URL.")
            return
        with logo_spinner("Scraping webpage structure…"):
            data = fetch_website_content(target_url)

        if "error" in data:
            st.error(f"❌ Failed to extract content: {data['error']}")
            return

        _card_open()
        st.markdown(f"**🌐 Title:** {data['title']}")
        st.markdown(f"**🔗 URL:** {data['url']}")

        if data["headings"]:
            st.markdown("**📌 Headings Discovered:**")
            for h in data["headings"]:
                st.markdown(f"- {h}")

        if data["paragraphs"]:
            st.markdown("**📄 Key Paragraphs:**")
            for p in data["paragraphs"]:
                st.caption(f"• {p}")
        _card_close()

        combined_text = "\n".join(data["headings"] + data["paragraphs"])
        _render_export_buttons(combined_text, "url_analysis")

        if st.button("🤖 AI-Summarise This Page", key="btn_ai_summarise_url"):
            combined = "\n".join(data["paragraphs"])
            prompt = f"Summarise the key academic insights from this web content:\n\n{combined}"
            with logo_spinner("Generating AI summary…"):
                summary = groq_chat(prompt)
            _card_open()
            st.markdown("**🤖 AI Summary:**")
            st.markdown(summary)
            _card_close()
            _render_export_buttons(summary, "url_summary")


# ==============================================================================
# MODULE 9 — Background Remover
# ==============================================================================

def render_bg_remover() -> None:
    """🧹 Image Background Eraser."""
    st.markdown("### 🧹 AI-Powered Background Eraser")
    st.caption("Upload a JPG or PNG image to remove its background instantly.")

    img_file = st.file_uploader("Upload Image", type=["jpg", "png", "jpeg"], key="bg_uploader")

    if not img_file:
        return

    try:
        img = Image.open(img_file).convert("RGBA")
    except Exception as exc:
        st.error(f"❌ Could not open image: {exc}")
        return

    col1, col2 = st.columns(2)
    with col1:
        st.image(img, caption="Original Image", use_container_width=True)

    if st.button("✂️ Remove Background", key="btn_remove_bg"):
        with logo_spinner("Processing image tensors…"):
            try:
                from rembg import remove
                output_img = remove(img)
            except Exception as exc:
                st.error(f"❌ Background removal failed: {exc}")
                log.error("rembg error: %s", exc)
                return

        with col2:
            st.image(output_img, caption="Background Removed", use_container_width=True)

        # Download as PNG
        buf_png = io.BytesIO()
        output_img.save(buf_png, format="PNG")
        st.download_button(
            "⬇️ Download PNG",
            data=buf_png.getvalue(),
            file_name="bg_removed.png",
            mime="image/png",
        )

        # Download as JPEG
        jpg_img = output_img.convert("RGB")
        buf_jpg = io.BytesIO()
        jpg_img.save(buf_jpg, format="JPEG", quality=95)
        st.download_button(
            "⬇️ Download JPEG",
            data=buf_jpg.getvalue(),
            file_name="bg_removed.jpg",
            mime="image/jpeg",
        )

        st.toast("✅ Background removed successfully!", icon="✂️")


# ==============================================================================
# MODULE 10 — Analytics
# ==============================================================================

def render_analytics() -> None:
    """📊 Platform Usage & Telemetry."""
    st.markdown("### 📊 Platform Usage & Telemetry Dashboard")

    vi = st.session_state.get("vector_index")
    faiss_size = vi.size if vi else 0

    col_a, col_b, col_c = st.columns(3)
    col_a.metric("🔢 Total Queries", st.session_state.get("total_queries", 0))
    col_b.metric("🚫 Security Blocks", st.session_state.get("blocked_count", 0))
    col_c.metric("📚 FAISS Chunk Count", faiss_size)

    _card_open()
    st.markdown("**Session Metadata:**")
    session_start = st.session_state.get("session_start", datetime.datetime.now())
    elapsed = datetime.datetime.now() - session_start
    hours, rem = divmod(int(elapsed.total_seconds()), 3600)
    minutes, seconds = divmod(rem, 60)

    user_info = st.session_state.get("user_info", {})
    st.json(
        {
            "session_start": str(session_start),
            "session_duration": f"{hours:02d}:{minutes:02d}:{seconds:02d}",
            "user_name": user_info.get("name", "—"),
            "user_role": user_info.get("role", "—"),
            "theme": st.session_state.get("theme", "—"),
            "rag_active": faiss_size > 0,
            "chat_messages": len(st.session_state.get("chat_history", [])),
        }
    )
    _card_close()


# ==============================================================================
# MODULE 11 — Architecture Blueprint
# ==============================================================================

def render_architecture() -> None:
    """🏛️ Executive Vision & System Architecture Documentation."""
    st.markdown("## 🏛️ Executive Vision & System Architecture")

    st.markdown("""
---
### Executive Vision & Core Goals

**EduSphere AI** is an enterprise-grade, privacy-first educational ecosystem capable of:
- **Adaptive Learning Paths** — personalised study plans and Socratic dialogue.
- **Real-Time RAG** — FAISS-backed context extraction from uploaded documents.
- **Multi-Modal Problem Solving** — code analysis, translation, summarisation, and image processing.
- **College AI Assistant Integration** — built-in knowledge base for college queries.
- **Professional Resume Builder** — AI-powered resume generation with PDF export.

---
### System Architecture

```
┌─────────────────────────────────────────────────────┐
│                    USER INTERFACE                   │
│            (Streamlit — Multi-Module SPA)           │
└─────────────────────────┬───────────────────────────┘
                          │
                          ▼
┌─────────────────────────────────────────────────────┐
│         INPUT EVALUATION & GUARDRAIL ENGINE         │
│   (Regex Safety · Privacy · Crisis Classifier)      │
└──────────┬──────────────────────────────┬───────────┘
           │ BLOCKED                      │ SAFE
           ▼                              ▼
┌──────────────────┐           ┌──────────────────────┐
│  BLOCK & REPORT  │           │   ROUTING ENGINE     │
│ Crisis Resources │           │  (12 Feature Modules)│
└──────────────────┘           └────────┬─────────────┘
                                        │
                          ┌─────────────┴──────────────┐
                          │                            │
                    DIRECT LLM                    RAG MODE
                          │                            │
                          │               ┌────────────▼────────────┐
                          │               │  FAISS VECTOR INDEX     │
                          │               │  (SentenceTransformers) │
                          │               └────────────┬────────────┘
                          │                            │ Top-K Chunks
                          │               ┌────────────▼────────────┐
                          │               │  COLLEGE KNOWLEDGE BASE │
                          │               │  + CONTEXT AUGMENTER    │
                          │               └────────────┬────────────┘
                          │                            │
                          └──────────────┬─────────────┘
                                         │
                          ┌──────────────▼──────────────┐
                          │  GROQ LLaMA-3.1 INFERENCE   │
                          │   (8B / 70B Instant API)    │
                          └──────────────┬──────────────┘
                                         │
                          ┌──────────────▼──────────────┐
                          │   STREAMING RESPONSE UI     │
                          │  + MULTI-FORMAT EXPORT       │
                          │  (TXT · PDF · MP3 · PNG)    │
                          └─────────────────────────────┘
```

---
### Security Architecture

| Layer | Mechanism |
|---|---|
| **Authentication** | bcrypt password hashing (cost=12) + SQLite persistence |
| **Input Safety** | Multi-pattern regex guardrails |
| **Privacy** | PII / credential pattern detection |
| **API Keys** | `.env` file, never logged or exposed |
| **Logging** | Structured — no sensitive data written |
| **Chat History** | SQLite with per-user session isolation |

---
### Database Schema

```sql
CREATE TABLE users (
    id            INTEGER PRIMARY KEY AUTOINCREMENT,
    username      TEXT UNIQUE NOT NULL,
    email         TEXT UNIQUE NOT NULL,
    full_name     TEXT,
    salt          TEXT NOT NULL,
    password_hash TEXT NOT NULL,
    theme         TEXT DEFAULT 'dark',
    created_at    TEXT NOT NULL
);

CREATE TABLE sessions (
    session_id    TEXT PRIMARY KEY,
    user_id       INTEGER NOT NULL,
    title         TEXT,
    created_at    TEXT NOT NULL,
    updated_at    TEXT NOT NULL
);

CREATE TABLE chats (
    id            INTEGER PRIMARY KEY AUTOINCREMENT,
    user_id       INTEGER NOT NULL,
    session_id    TEXT NOT NULL,
    role          TEXT NOT NULL,
    content       TEXT NOT NULL,
    created_at    TEXT NOT NULL,
    FOREIGN KEY (user_id) REFERENCES users (id)
);

CREATE TABLE document_embeddings (
    chunk_id       UUID PRIMARY KEY,
    session_id     UUID REFERENCES sessions(session_id),
    document_name  VARCHAR(255) NOT NULL,
    chunk_index    INT NOT NULL,
    chunk_content  TEXT NOT NULL,
    embedding_vec  VECTOR(384)  -- all-MiniLM-L6-v2 dimension
);
```

---
### Project Structure

```text
edusphere-ai/
├── .env                  # API keys (never commit)
├── DasaAI.py             # Streamlit entry point
├── data.json             # College knowledge base + greetings + interaction log
├── users.db              # SQLite user/session/chat store (auto-created)
├── src/
│   ├── __init__.py
│   ├── auth.py           # bcrypt authentication
│   ├── config.py         # constants & env loading (6 themes)
│   ├── database.py       # SQLite CRUD layer
│   ├── exceptions.py     # custom exception hierarchy
│   ├── main.py           # Streamlit entrypoint
│   ├── modules.py        # all 12 feature modules
│   ├── rag.py            # FAISS RAG pipeline
│   └── utils.py          # GROQ client, guardrails, exports & helpers
├── assets/
│   ├── styles.css        # premium CSS stylesheet
│   └── nepal_map.png     # login background
├── tests/
│   ├── test_auth.py
│   ├── test_rag.py
│   └── test_utils.py
├── logs/
│   └── app.log
└── requirements.txt
```
    """)


# ==============================================================================
# MODULE 12 — Resume Builder + Portfolio Website Builder
# ==============================================================================

# ──────────────────────────────────────────────────────────────────────────────
# PORTFOLIO HTML TEMPLATE ENGINE
# ──────────────────────────────────────────────────────────────────────────────

def _generate_portfolio_html(data: dict, theme: str) -> str:
    """Generate a complete single-file animated HTML portfolio website."""

    themes = {
        "🌌 Nebula Dark": {
            "bg": "#050a14",
            "card": "rgba(15,25,50,0.7)",
            "accent": "#00f0ff",
            "accent2": "#7c3aed",
            "accent3": "#ff007f",
            "text": "#e2e8f0",
            "muted": "#64748b",
            "border": "rgba(0,240,255,0.2)",
            "glow": "rgba(0,240,255,0.3)",
            "particle_color": "0,240,255",
            "loader_color": "#00f0ff",
            "gradient": "linear-gradient(135deg,#050a14 0%,#0d1b3e 50%,#050a14 100%)",
            "aurora1": "rgba(0,240,255,0.08)",
            "aurora2": "rgba(124,58,237,0.08)",
            "aurora3": "rgba(255,0,127,0.06)",
        },
        "🔥 Crimson Forge": {
            "bg": "#0d0500",
            "card": "rgba(30,10,5,0.75)",
            "accent": "#ff4500",
            "accent2": "#ff8c00",
            "accent3": "#ffd700",
            "text": "#fde8d8",
            "muted": "#8b6355",
            "border": "rgba(255,69,0,0.25)",
            "glow": "rgba(255,69,0,0.35)",
            "particle_color": "255,80,0",
            "loader_color": "#ff4500",
            "gradient": "linear-gradient(135deg,#0d0500 0%,#200800 50%,#0d0500 100%)",
            "aurora1": "rgba(255,69,0,0.08)",
            "aurora2": "rgba(255,140,0,0.07)",
            "aurora3": "rgba(255,215,0,0.05)",
        },
        "🌊 Ocean Pulse": {
            "bg": "#020d1a",
            "card": "rgba(5,25,50,0.72)",
            "accent": "#00d4ff",
            "accent2": "#0077ff",
            "accent3": "#00ffb3",
            "text": "#cde8f5",
            "muted": "#4a7a96",
            "border": "rgba(0,212,255,0.2)",
            "glow": "rgba(0,212,255,0.3)",
            "particle_color": "0,212,255",
            "loader_color": "#00d4ff",
            "gradient": "linear-gradient(135deg,#020d1a 0%,#041e38 50%,#020d1a 100%)",
            "aurora1": "rgba(0,212,255,0.07)",
            "aurora2": "rgba(0,119,255,0.07)",
            "aurora3": "rgba(0,255,179,0.05)",
        },
        "🌿 Emerald Matrix": {
            "bg": "#010d05",
            "card": "rgba(5,25,12,0.72)",
            "accent": "#00ff41",
            "accent2": "#39d353",
            "accent3": "#00e5ff",
            "text": "#ccffd8",
            "muted": "#3a7a50",
            "border": "rgba(0,255,65,0.2)",
            "glow": "rgba(0,255,65,0.3)",
            "particle_color": "0,255,65",
            "loader_color": "#00ff41",
            "gradient": "linear-gradient(135deg,#010d05 0%,#041a0c 50%,#010d05 100%)",
            "aurora1": "rgba(0,255,65,0.07)",
            "aurora2": "rgba(57,211,83,0.06)",
            "aurora3": "rgba(0,229,255,0.05)",
        },
    }

    t = themes.get(theme, themes["🌌 Nebula Dark"])

    name       = data.get("name", "Your Name")
    title      = data.get("title", "Full-Stack Developer & AI Enthusiast")
    email      = data.get("email", "")
    phone      = data.get("phone", "")
    location   = data.get("location", "")
    github     = data.get("github", "")
    linkedin   = data.get("linkedin", "")
    about      = data.get("about", "I build thoughtful digital products and experiences.")
    skills_raw = data.get("skills", "Python, JavaScript, React, Machine Learning")
    projects   = data.get("projects", "")
    education  = data.get("education", "")
    certs      = data.get("certs", "")

    # Parse skills into list
    skills_list = [s.strip() for s in skills_raw.replace("\n", ",").split(",") if s.strip()]

    # Build skills pills HTML
    skills_pills = "".join(
        f'<span class="skill-pill" style="animation-delay:{i*0.07:.2f}s">{s}</span>'
        for i, s in enumerate(skills_list)
    )

    # Build skill bars (first 6 skills with fake percentages for visual effect)
    skill_bars_html = ""
    bar_vals = [92, 85, 78, 88, 72, 80, 76, 90, 82, 70]
    for i, s in enumerate(skills_list[:8]):
        pct = bar_vals[i % len(bar_vals)]
        skill_bars_html += f"""
        <div class="skill-bar-item" style="animation-delay:{i*0.1:.1f}s">
          <div class="skill-bar-label"><span>{s}</span><span>{pct}%</span></div>
          <div class="skill-bar-track"><div class="skill-bar-fill" style="width:{pct}%;animation-delay:{i*0.1+0.3:.1f}s"></div></div>
        </div>"""

    # Build project cards HTML
    project_cards_html = ""
    if projects.strip():
        for i, line in enumerate(projects.strip().split("\n")):
            if line.strip() and not line.startswith("•"):
                parts = line.split("—") if "—" in line else line.split("-")
                proj_name = parts[0].strip()
                proj_desc = parts[1].strip() if len(parts) > 1 else "A remarkable project."
                project_cards_html += f"""
                <div class="proj-card" style="animation-delay:{i*0.1:.1f}s">
                  <div class="proj-glow"></div>
                  <h3>{proj_name}</h3>
                  <p>{proj_desc}</p>
                  <div class="proj-tags"><span>Featured</span><span>#{i+1}</span></div>
                </div>"""

    if not project_cards_html:
        project_cards_html = '<div class="proj-card"><h3>Add Your Projects</h3><p>Projects you add in the form will appear here.</p></div>'

    # Roles for typewriter
    roles = [title, "Creative Coder", "Problem Solver", "Digital Craftsman"]
    roles_js = str(roles).replace("'", '"')

    html = f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width,initial-scale=1">
<meta name="description" content="{name} — portfolio website">
<title>{name} | Portfolio</title>
<link rel="preconnect" href="https://fonts.googleapis.com">
<link rel="preconnect" href="https://fonts.gstatic.com" crossorigin>
<link href="https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700;800;900&family=Space+Grotesk:wght@300;400;500;600;700&family=JetBrains+Mono:wght@400;700&display=swap" rel="stylesheet">
<style>
:root {{
  --bg:{t['bg']};--card:{t['card']};--accent:{t['accent']};--accent2:{t['accent2']};
  --accent3:{t['accent3']};--text:{t['text']};--muted:{t['muted']};--border:{t['border']};
  --glow:{t['glow']};--gradient:{t['gradient']};
  --font-body:'Inter',sans-serif;--font-title:'Space Grotesk',sans-serif;--font-mono:'JetBrains Mono',monospace;
}}
*,*::before,*::after{{box-sizing:border-box;margin:0;padding:0}}
html{{scroll-behavior:smooth}}
body{{background:var(--bg);color:var(--text);font-family:var(--font-body);overflow-x:hidden;line-height:1.6}}

/* ── LOADER ── */
#loader{{position:fixed;inset:0;z-index:9999;background:var(--bg);display:flex;flex-direction:column;align-items:center;justify-content:center;gap:20px;transition:opacity 0.5s ease}}
#loader.done{{opacity:0;pointer-events:none}}
.ld-title{{font-family:var(--font-title);font-size:2.5rem;font-weight:700;background:linear-gradient(135deg,var(--accent),var(--accent2));-webkit-background-clip:text;-webkit-text-fill-color:transparent}}
.ld-sub{{font-family:var(--font-mono);font-size:0.85rem;color:var(--muted);letter-spacing:3px;text-transform:uppercase}}
.ld-bar{{width:260px;height:3px;background:rgba(255,255,255,0.08);border-radius:99px;overflow:hidden}}
.ld-bar-fill{{height:100%;width:0%;background:linear-gradient(90deg,var(--accent),var(--accent2));border-radius:99px;transition:width 0.04s linear;box-shadow:0 0 8px var(--glow)}}
.ld-pct{{font-family:var(--font-mono);font-size:0.8rem;color:var(--accent)}}
@keyframes ldFill{{to{{width:100%}}}}
.ld-rings{{display:flex;gap:12px}}
.ld-ring{{width:14px;height:14px;border-radius:50%;border:2px solid var(--accent);animation:ldPulse 1.2s ease-in-out infinite}}
.ld-ring:nth-child(2){{animation-delay:.2s;border-color:var(--accent2)}}
.ld-ring:nth-child(3){{animation-delay:.4s;border-color:var(--accent3)}}
.ld-ring:nth-child(4){{animation-delay:.6s}}
@keyframes ldPulse{{0%,100%{{transform:scale(1);opacity:1}}50%{{transform:scale(1.4);opacity:.5}}}}

/* ── CANVAS BG ── */
#bgCanvas{{position:fixed;inset:0;z-index:0;pointer-events:none}}

/* ── AURORA ── */
.aurora{{position:fixed;border-radius:50%;filter:blur(80px);pointer-events:none;z-index:0;animation:auroraFloat 8s ease-in-out infinite}}
.aurora-1{{width:600px;height:600px;background:{t['aurora1']};top:-200px;left:-150px;animation-delay:0s}}
.aurora-2{{width:500px;height:500px;background:{t['aurora2']};top:30%;right:-200px;animation-delay:3s}}
.aurora-3{{width:400px;height:400px;background:{t['aurora3']};bottom:-100px;left:40%;animation-delay:5s}}
@keyframes auroraFloat{{0%,100%{{transform:translate(0,0) scale(1)}}50%{{transform:translate(30px,20px) scale(1.08)}}}}

/* ── SCROLL PROGRESS ── */
#progress{{position:fixed;top:0;left:0;height:3px;background:linear-gradient(90deg,var(--accent),var(--accent2),var(--accent3));z-index:9998;width:0%;transition:width .1s}}

/* ── NAVBAR ── */
.nav{{position:fixed;top:20px;left:50%;transform:translateX(-50%);z-index:999;padding:12px 24px;display:flex;align-items:center;justify-content:space-between;background:rgba(5,10,20,0.65);backdrop-filter:blur(24px);border:1px solid rgba(255,255,255,0.1);border-radius:50px;width:90%;max-width:800px;box-shadow:0 10px 30px rgba(0,0,0,0.3);transition:all .3s}}
.nav:hover{{background:rgba(5,10,20,0.8);border-color:var(--accent);box-shadow:0 10px 40px var(--glow)}}
.nav-brand{{font-family:var(--font-title);font-weight:800;font-size:1.3rem;color:var(--text);text-decoration:none;letter-spacing:1px;display:flex;align-items:center;gap:6px}}
.nav-brand span{{color:var(--accent)}}
.nav-links{{display:flex;gap:30px;list-style:none;margin-left:auto}}
.nav-links a{{color:var(--text);text-decoration:none;font-size:0.95rem;font-weight:600;transition:all .3s;position:relative;padding:8px 12px;border-radius:20px}}
.nav-links a:hover{{color:var(--bg);background:var(--accent);box-shadow:0 0 15px var(--glow)}}

/* ── FULLSCREEN BTN ── */
.fs-btn{{position:fixed;bottom:30px;right:30px;width:55px;height:55px;border-radius:50%;background:linear-gradient(135deg,var(--accent),var(--accent2));color:#000;border:none;box-shadow:0 4px 25px var(--glow);cursor:pointer;z-index:9999;display:flex;align-items:center;justify-content:center;font-size:1.5rem;transition:all .3s}}
.fs-btn:hover{{transform:scale(1.1);box-shadow:0 6px 35px var(--glow);color:#fff}}

/* ── HERO ── */
.hero{{min-height:100vh;display:flex;align-items:center;padding:120px 5% 60px;position:relative;z-index:1;gap:60px}}
.hero-copy{{flex:1;max-width:600px}}
.hero-eyebrow{{font-family:var(--font-mono);font-size:0.78rem;color:var(--accent);letter-spacing:3px;text-transform:uppercase;margin-bottom:16px}}
.hero-name{{font-family:var(--font-title);font-size:clamp(2.5rem,6vw,4.5rem);font-weight:800;line-height:1.1;margin-bottom:12px;background:linear-gradient(135deg,var(--text) 0%,var(--accent) 100%);-webkit-background-clip:text;-webkit-text-fill-color:transparent}}
.hero-role{{font-size:1.4rem;color:var(--muted);margin-bottom:20px;font-weight:400}}
.hero-role span{{color:var(--accent);font-weight:600}}
.type-caret{{display:inline-block;width:2px;height:1.2em;background:var(--accent);margin-left:3px;animation:blink .7s step-end infinite;vertical-align:bottom}}
@keyframes blink{{0%,100%{{opacity:1}}50%{{opacity:0}}}}
.hero-desc{{color:var(--muted);font-size:1.05rem;line-height:1.7;margin-bottom:32px;max-width:480px}}
.hero-actions{{display:flex;gap:14px;flex-wrap:wrap;margin-bottom:36px}}
.btn{{padding:12px 28px;border-radius:50px;font-weight:600;font-size:0.95rem;text-decoration:none;transition:all .3s;display:inline-flex;align-items:center;gap:8px;cursor:pointer;border:none}}
.btn-primary{{background:linear-gradient(135deg,var(--accent),var(--accent2));color:#fff;box-shadow:0 0 25px var(--glow)}}
.btn-primary:hover{{transform:translateY(-3px);box-shadow:0 0 40px var(--glow)}}
.btn-secondary{{background:transparent;color:var(--accent);border:1px solid var(--border)}}
.btn-secondary:hover{{background:var(--accent);color:#000;transform:translateY(-3px)}}
.hero-stats{{display:flex;gap:36px}}
.stat{{display:flex;flex-direction:column}}
.stat-val{{font-family:var(--font-title);font-size:1.8rem;font-weight:800;color:var(--accent)}}
.stat-lbl{{font-size:0.75rem;color:var(--muted);text-transform:uppercase;letter-spacing:1px}}
.hero-visual{{flex:0 0 400px;position:relative}}
.portrait-frame{{width:340px;height:380px;position:relative;margin:0 auto}}
.portrait-orb{{position:absolute;border-radius:50%;animation:orbFloat 6s ease-in-out infinite}}
.portrait-orb-1{{width:320px;height:320px;background:radial-gradient(circle,{t['aurora2']} 0%,transparent 70%);top:-20px;left:-20px}}
.portrait-orb-2{{width:250px;height:250px;background:radial-gradient(circle,{t['aurora1']} 0%,transparent 70%);bottom:-30px;right:-30px;animation-delay:3s}}
@keyframes orbFloat{{0%,100%{{transform:translate(0,0) scale(1)}}50%{{transform:translate(10px,-15px) scale(1.05)}}}}
.portrait-glass{{width:100%;height:100%;border-radius:30px;background:var(--card);border:1px solid var(--border);backdrop-filter:blur(20px);display:flex;align-items:center;justify-content:center;font-size:7rem;box-shadow:0 0 60px var(--glow),inset 0 1px 0 rgba(255,255,255,0.1);position:relative;z-index:1;overflow:hidden}}
.portrait-glass::before{{content:'';position:absolute;inset:0;background:linear-gradient(135deg,rgba(255,255,255,0.05),transparent);pointer-events:none}}
.float-tag{{position:absolute;padding:6px 14px;border-radius:30px;font-size:0.75rem;font-weight:600;background:var(--card);border:1px solid var(--border);backdrop-filter:blur(10px);color:var(--accent);z-index:2;animation:tagFloat 4s ease-in-out infinite}}
.tag-tl{{top:-10px;left:-30px;animation-delay:0s}}
.tag-tr{{top:-10px;right:-30px;animation-delay:1s}}
.tag-bl{{bottom:20px;left:-40px;animation-delay:2s}}
.tag-br{{bottom:20px;right:-40px;animation-delay:3s}}
@keyframes tagFloat{{0%,100%{{transform:translateY(0)}}50%{{transform:translateY(-8px)}}}}

/* ── SECTIONS ── */
.section{{padding:90px 5%;position:relative;z-index:1}}
.section-head{{text-align:center;margin-bottom:60px}}
.eyebrow{{font-family:var(--font-mono);font-size:0.78rem;color:var(--accent);letter-spacing:3px;text-transform:uppercase;margin-bottom:12px;display:block}}
.section-head h2{{font-family:var(--font-title);font-size:clamp(1.8rem,4vw,2.8rem);font-weight:700;margin-bottom:14px}}
.section-head p{{color:var(--muted);max-width:500px;margin:0 auto}}

/* ── GLASS CARD ── */
.glass{{background:var(--card);border:1px solid var(--border);border-radius:20px;padding:32px;backdrop-filter:blur(20px);transition:all .3s;position:relative;overflow:hidden}}
.glass::before{{content:'';position:absolute;inset:0;background:linear-gradient(135deg,rgba(255,255,255,0.03),transparent);pointer-events:none}}
.glass:hover{{border-color:var(--accent);box-shadow:0 0 30px var(--glow);transform:translateY(-4px)}}
.grid-2{{display:grid;grid-template-columns:repeat(auto-fit,minmax(300px,1fr));gap:24px}}

/* ── SKILLS ── */
.skill-pill{{display:inline-block;padding:6px 16px;border-radius:30px;background:rgba(255,255,255,0.04);border:1px solid var(--border);color:var(--accent);font-size:0.85rem;margin:5px;transition:all .3s;animation:pillPop .4s ease both}}
.skill-pill:hover{{background:var(--accent);color:#000;border-color:var(--accent);transform:scale(1.05)}}
@keyframes pillPop{{from{{opacity:0;transform:scale(0.8)}}to{{opacity:1;transform:scale(1)}}}}
.skill-bar-item{{margin-bottom:18px}}
.skill-bar-label{{display:flex;justify-content:space-between;font-size:0.88rem;margin-bottom:7px;color:var(--text)}}
.skill-bar-track{{height:6px;background:rgba(255,255,255,0.07);border-radius:99px;overflow:hidden}}
.skill-bar-fill{{height:100%;background:linear-gradient(90deg,var(--accent),var(--accent2));border-radius:99px;width:0%;transition:width 1.2s cubic-bezier(.25,.46,.45,.94);box-shadow:0 0 10px var(--glow)}}

/* ── PROJECTS ── */
.proj-grid{{display:grid;grid-template-columns:repeat(auto-fill,minmax(280px,1fr));gap:24px}}
.proj-card{{background:var(--card);border:1px solid var(--border);border-radius:20px;padding:28px;position:relative;overflow:hidden;transition:all .3s;animation:fadeUp .6s ease both}}
.proj-card:hover{{border-color:var(--accent);transform:translateY(-6px);box-shadow:0 0 40px var(--glow)}}
.proj-glow{{position:absolute;width:200px;height:200px;background:radial-gradient(circle,var(--glow),transparent);top:-80px;right:-80px;pointer-events:none;opacity:0;transition:opacity .3s}}
.proj-card:hover .proj-glow{{opacity:1}}
.proj-card h3{{font-family:var(--font-title);font-size:1.1rem;font-weight:700;margin-bottom:10px;color:var(--accent)}}
.proj-card p{{color:var(--muted);font-size:0.88rem;line-height:1.6}}
.proj-tags{{display:flex;gap:8px;margin-top:16px}}
.proj-tags span{{padding:3px 12px;border-radius:20px;font-size:0.72rem;font-weight:600;background:rgba(255,255,255,0.05);border:1px solid var(--border);color:var(--accent2)}}
@keyframes fadeUp{{from{{opacity:0;transform:translateY(20px)}}to{{opacity:1;transform:translateY(0)}}}}

/* ── CONTACT ── */
.contact-info{{display:flex;flex-direction:column;gap:16px}}
.contact-row{{display:flex;align-items:center;gap:12px;padding:14px 18px;background:rgba(255,255,255,0.03);border:1px solid var(--border);border-radius:12px;transition:all .3s}}
.contact-row:hover{{border-color:var(--accent);background:rgba(255,255,255,0.05)}}
.contact-row .ci-icon{{font-size:1.2rem}}
.contact-row a{{color:var(--text);text-decoration:none;font-size:0.9rem;transition:color .2s}}
.contact-row:hover a{{color:var(--accent)}}

/* ── FOOTER ── */
footer{{position:relative;z-index:1;border-top:1px solid var(--border);padding:40px 5%;text-align:center}}
.ribbon{{overflow:hidden;white-space:nowrap;background:linear-gradient(90deg,var(--accent),var(--accent2),var(--accent3),var(--accent));background-size:200%;-webkit-background-clip:text;-webkit-text-fill-color:transparent;font-family:var(--font-mono);font-size:0.8rem;letter-spacing:3px;padding:12px 0;margin-bottom:24px;animation:ribbonMove 10s linear infinite}}
@keyframes ribbonMove{{to{{background-position:200%}}}}
.ribbon-track{{display:inline-block;animation:scrollRibbon 20s linear infinite}}
@keyframes scrollRibbon{{from{{transform:translateX(0)}}to{{transform:translateX(-50%)}}}}
.footer-copy{{color:var(--muted);font-size:0.82rem}}

/* ── REVEAL ANIMATION ── */
.reveal{{opacity:0;transform:translateY(30px);transition:all .7s cubic-bezier(.25,.46,.45,.94)}}
.reveal.visible{{opacity:1;transform:translateY(0)}}

@media(max-width:768px){{
  .hero{{flex-direction:column;text-align:center;padding-top:100px}}
  .hero-visual{{display:none}}
  .hero-stats{{justify-content:center}}
  .nav-links{{display:none}}
}}
</style>
</head>
<body>

<!-- LOADER -->
<div id="loader">
  <div class="ld-rings">
    <div class="ld-ring"></div><div class="ld-ring"></div>
    <div class="ld-ring"></div><div class="ld-ring"></div>
  </div>
  <div class="ld-title">{name.split()[0] if name.split() else name}</div>
  <div class="ld-sub">Loading Portfolio</div>
  <div class="ld-bar"><div class="ld-bar-fill"></div></div>
  <div class="ld-pct" id="ldPct">0%</div>
</div>

<!-- SCROLL PROGRESS -->
<div id="progress"></div>

<!-- CANVAS BG -->
<canvas id="bgCanvas"></canvas>

<!-- AURORA -->
<div class="aurora aurora-1"></div>
<div class="aurora aurora-2"></div>
<div class="aurora aurora-3"></div>

<!-- NAVBAR -->
<nav class="nav">
  <a class="nav-brand" href="#home">{name.split()[0] if name.split() else name}</a>
  <ul class="nav-links">
    <li><a href="#about">About</a></li>
    <li><a href="#skills">Skills</a></li>
    <li><a href="#projects">Projects</a></li>
    <li><a href="#contact">Contact</a></li>
  </ul>
</nav>

<main>
<!-- HERO -->
<section class="hero" id="home">
  <div class="hero-copy reveal">
    <p class="hero-eyebrow">✨ Welcome to my portfolio</p>
    <h1 class="hero-name">Hi, I'm {name}.</h1>
    <p class="hero-role">I am <span id="typeText">{title}</span><span class="type-caret"></span></p>
    <p class="hero-desc">{about}</p>
    <div class="hero-actions">
      <a class="btn btn-primary" href="#projects">View Work</a>
      <a class="btn btn-secondary" href="#contact">Contact Me</a>
    </div>
    <div class="hero-stats">
      <div class="stat"><span class="stat-val">{len(skills_list)}+</span><span class="stat-lbl">Skills</span></div>
      <div class="stat"><span class="stat-val">∞</span><span class="stat-lbl">Passion</span></div>
      <div class="stat"><span class="stat-val">24/7</span><span class="stat-lbl">Online</span></div>
    </div>
  </div>
  <div class="hero-visual reveal">
    <div class="portrait-frame">
      <div class="portrait-orb portrait-orb-1"></div>
      <div class="portrait-orb portrait-orb-2"></div>
      <div class="portrait-glass">👤</div>
      <div class="float-tag tag-tl">{skills_list[0] if skills_list else 'Developer'}</div>
      <div class="float-tag tag-tr">{skills_list[1] if len(skills_list)>1 else 'Designer'}</div>
      <div class="float-tag tag-bl">{skills_list[2] if len(skills_list)>2 else 'Creator'}</div>
      <div class="float-tag tag-br">{'AI' if 'ai' not in skills_list[0].lower() else skills_list[-1]}</div>
    </div>
  </div>
</section>

<!-- ABOUT -->
<section class="section" id="about">
  <div class="section-head reveal">
    <span class="eyebrow">About Me</span>
    <h2>The person behind the work.</h2>
    <p>Passionate about crafting digital experiences that are both beautiful and functional.</p>
  </div>
  <div class="grid-2">
    <div class="glass reveal">
      <h3 style="font-family:var(--font-title);margin-bottom:12px;color:var(--accent)">What I Do</h3>
      <p style="color:var(--muted);line-height:1.8">{about}</p>
      {'<p style="margin-top:14px;color:var(--muted)"><strong style="color:var(--text)">Education:</strong> ' + education.replace(chr(10),'<br>') + '</p>' if education else ''}
    </div>
    <div class="glass reveal">
      <h3 style="font-family:var(--font-title);margin-bottom:16px;color:var(--accent)">Tech Stack</h3>
      <div>{skills_pills}</div>
      {'<div style="margin-top:20px;padding-top:16px;border-top:1px solid var(--border);color:var(--muted);font-size:0.88rem"><strong style="color:var(--text)">Certifications:</strong><br>' + certs.replace(chr(10),'<br>') + '</div>' if certs else ''}
    </div>
  </div>
</section>

<!-- SKILLS -->
<section class="section" id="skills" style="background:linear-gradient(180deg,transparent,rgba(255,255,255,0.01),transparent)">
  <div class="section-head reveal">
    <span class="eyebrow">Skills</span>
    <h2>Tools & Strengths.</h2>
  </div>
  <div class="glass reveal" style="max-width:700px;margin:0 auto">
    {skill_bars_html}
  </div>
</section>

<!-- PROJECTS -->
<section class="section" id="projects">
  <div class="section-head reveal">
    <span class="eyebrow">Projects</span>
    <h2>Selected work.</h2>
  </div>
  <div class="proj-grid">
    {project_cards_html}
  </div>
</section>

<!-- CONTACT -->
<section class="section" id="contact">
  <div class="section-head reveal">
    <span class="eyebrow">Contact</span>
    <h2>Let's connect.</h2>
    <p>Open for collaborations, freelance, and full-time roles.</p>
  </div>
  <div style="max-width:500px;margin:0 auto">
    <div class="glass reveal">
      <div class="contact-info">
        {f'<div class="contact-row"><span class="ci-icon">📧</span><a href="mailto:{email}">{email}</a></div>' if email else ''}
        {f'<div class="contact-row"><span class="ci-icon">📱</span><span>{phone}</span></div>' if phone else ''}
        {f'<div class="contact-row"><span class="ci-icon">📍</span><span>{location}</span></div>' if location else ''}
        {f'<div class="contact-row"><span class="ci-icon">💼</span><a href="{linkedin}" target="_blank">LinkedIn</a></div>' if linkedin else ''}
        {f'<div class="contact-row"><span class="ci-icon">🐙</span><a href="{github}" target="_blank">GitHub</a></div>' if github else ''}
      </div>
    </div>
  </div>
</section>
</main>

<!-- FOOTER -->
<footer>
  <div class="ribbon"><div class="ribbon-track">
    ✦ {name.upper()} &nbsp; ✦ {title.upper()} &nbsp; ✦ {location.upper() or 'WORLD'} &nbsp;
    ✦ {name.upper()} &nbsp; ✦ {title.upper()} &nbsp; ✦ {location.upper() or 'WORLD'} &nbsp;
    ✦ {name.upper()} &nbsp; ✦ {title.upper()} &nbsp; ✦ {location.upper() or 'WORLD'} &nbsp;
  </div></div>
  <p class="footer-copy">&copy; 2026 {name}. All Rights Reserved. Built with EduSphere AI.</p>
</footer>
<!-- FULLSCREEN BTN -->
<button class="fs-btn" id="fsBtn" title="Toggle Fullscreen">⛶</button>
<script>
document.addEventListener('DOMContentLoaded', function() {{

// ── Loader: robust step-based counter, no CSS animation conflict
var _loaderEl  = document.getElementById('loader');
var _barFill   = document.querySelector('.ld-bar-fill');
var _pctTxt    = document.getElementById('ldPct');
var _loaderPct = 0;
function _loaderStep() {{
  _loaderPct += 2;
  if (_loaderPct > 100) _loaderPct = 100;
  if (_barFill) _barFill.style.width = _loaderPct + '%';
  if (_pctTxt)  _pctTxt.textContent  = _loaderPct + '%';
  if (_loaderPct < 100) {{
    setTimeout(_loaderStep, 18);
  }} else {{
    setTimeout(function() {{
      if (_loaderEl) {{
        _loaderEl.classList.add('done');
        setTimeout(function() {{ if (_loaderEl) _loaderEl.style.display = 'none'; }}, 550);
      }}
    }}, 300);
  }}
}}
setTimeout(_loaderStep, 80);

// ── Typewriter
(function () {{
  var roles = {roles_js};
  var roleIndex = 0;
  var charIndex = 0;
  var deleting = false;
  var element = document.getElementById('typeText');
  if (!element) return;

  function type() {{
    var current = roles[roleIndex];

    if (!deleting) {{
      element.textContent = current.substring(0, charIndex + 1);
      charIndex++;
      if (charIndex === current.length) {{
        deleting = true;
        return setTimeout(type, 1200);
      }}
    }} else {{
      element.textContent = current.substring(0, charIndex - 1);
      charIndex--;
      if (charIndex === 0) {{
        deleting = false;
        roleIndex = (roleIndex + 1) % roles.length;
      }}
    }}

    setTimeout(type, deleting ? 50 : 90);
  }}

  type();
}})();

// ── Particle canvas
(function() {{
  var cv  = document.getElementById('bgCanvas');
  if (!cv) return;
  var ctx = cv.getContext('2d');
  var W, H, pts = [];
  function resize() {{ W = cv.width = window.innerWidth; H = cv.height = window.innerHeight; }}
  resize();
  window.addEventListener('resize', resize);
  for (var i = 0; i < 70; i++) {{
    pts.push({{ x: Math.random()*W, y: Math.random()*H,
               vx: (Math.random()-.5)*.5, vy: (Math.random()-.5)*.5 }});
  }}
  function draw() {{
    ctx.clearRect(0, 0, W, H);
    for (var p = 0; p < pts.length; p++) {{
      pts[p].x += pts[p].vx; pts[p].y += pts[p].vy;
      if (pts[p].x < 0 || pts[p].x > W) pts[p].vx *= -1;
      if (pts[p].y < 0 || pts[p].y > H) pts[p].vy *= -1;
      ctx.beginPath(); ctx.arc(pts[p].x, pts[p].y, 1.5, 0, Math.PI*2);
      ctx.fillStyle = 'rgba({t["particle_color"]},0.6)'; ctx.fill();
    }}
    for (var a = 0; a < pts.length; a++) {{
      for (var b = a+1; b < pts.length; b++) {{
        var dx = pts[a].x-pts[b].x, dy = pts[a].y-pts[b].y;
        var d  = Math.sqrt(dx*dx+dy*dy);
        if (d < 120) {{
          ctx.beginPath(); ctx.moveTo(pts[a].x, pts[a].y); ctx.lineTo(pts[b].x, pts[b].y);
          ctx.strokeStyle = 'rgba({t["particle_color"]},'+(((1-d/120)*0.18)).toFixed(2)+')';
          ctx.lineWidth = 0.7; ctx.stroke();
        }}
      }}
    }}
    requestAnimationFrame(draw);
  }}
  draw();
}})();

// ── Scroll progress
var _progBar = document.getElementById('progress');
window.addEventListener('scroll', function() {{
  var sc = document.documentElement;
  var pct = sc.scrollTop / (sc.scrollHeight - sc.clientHeight) * 100;
  if (_progBar) _progBar.style.width = pct + '%';
}});

// ── Reveal on scroll
var revealEls = document.querySelectorAll('.reveal');
if (window.IntersectionObserver) {{
  var io = new IntersectionObserver(function(ents) {{
    ents.forEach(function(e) {{ if (e.isIntersecting) e.target.classList.add('visible'); }});
  }}, {{ threshold: 0.12 }});
  revealEls.forEach(function(el) {{ io.observe(el); }});
}} else {{
  revealEls.forEach(function(el) {{ el.classList.add('visible'); }});
}}

// ── Skill bars animate on scroll
var bars = document.querySelectorAll('.skill-bar-fill');
bars.forEach(function(b) {{ b.dataset.w = b.style.width; b.style.width = '0%'; }});
if (window.IntersectionObserver) {{
  var bIO = new IntersectionObserver(function(ents) {{
    ents.forEach(function(e) {{
      if (e.isIntersecting) {{ e.target.style.width = e.target.dataset.w || '70%'; e.target.style.transition='width 1.2s cubic-bezier(.25,.46,.45,.94)'; }}
    }});
  }}, {{ threshold: 0.3 }});
  bars.forEach(function(b) {{ bIO.observe(b); }});
}} else {{
  setTimeout(function() {{ bars.forEach(function(b) {{ b.style.width = b.dataset.w || '70%'; }}); }}, 800);
}}

// FULLSCREEN TOGGLE
var fsb = document.getElementById('fsBtn');
if (fsb) {{
  fsb.addEventListener('click', function() {{
    if (!document.fullscreenElement) {{
      document.documentElement.requestFullscreen().catch(err => {{
        console.log("Error attempting to enable fullscreen: " + err.message);
      }});
    }} else {{
      document.exitFullscreen();
    }}
  }});
  document.addEventListener('fullscreenchange', function() {{
    if (document.fullscreenElement) {{
      fsb.textContent = '✕';
      fsb.style.background = 'rgba(255,50,50,0.8)';
    }} else {{
      fsb.textContent = '⛶';
      fsb.style.background = 'linear-gradient(135deg,var(--accent),var(--accent2))';
    }}
  }});
}}

}}); // end DOMContentLoaded
</script>
</body>
</html>"""
    return html


def render_portfolio_builder() -> None:
    """🌐 Real-Time Portfolio Website Builder."""
    st.markdown(
        """
        <div style="text-align:center;margin-bottom:6px;">
          <span style="font-family:'JetBrains Mono',monospace;font-size:0.72rem;letter-spacing:3px;
                       color:var(--accent);text-transform:uppercase;">Portfolio Builder</span>
          <h3 style="margin:6px 0 4px;font-size:1.4rem;">🌐 Build Your Personal Website — Instantly</h3>
          <p style="color:var(--sub);font-size:0.88rem;">Fill your info, pick a theme, and get a full animated HTML site in seconds.</p>
        </div>
        """,
        unsafe_allow_html=True,
    )

    _card_open()
    col1, col2 = st.columns(2)

    with col1:
        p_name     = st.text_input("👤 Full Name",          placeholder="Sanjaya Kandel",          key="pf_name",     value=st.session_state.get("resume_name", ""))
        p_title    = st.text_input("🏷️ Professional Title", placeholder="Full-Stack Developer & AI Enthusiast", key="pf_title")
        p_email    = st.text_input("📧 Email",              placeholder="you@example.com",          key="pf_email",    value=st.session_state.get("resume_email", ""))
        p_phone    = st.text_input("📱 Phone",              placeholder="+977-9800000000",          key="pf_phone",    value=st.session_state.get("resume_phone", ""))
        p_location = st.text_input("📍 Location",           placeholder="Kathmandu, Nepal",         key="pf_location", value=st.session_state.get("resume_location", ""))

    with col2:
        p_about    = st.text_area(
            "💬 About / Bio",
            placeholder="I build thoughtful digital products and experiences that feel clear from the first interaction.",
            height=100,
            key="pf_about",
        )
        p_github   = st.text_input("🐙 GitHub URL",   placeholder="https://github.com/yourname",    key="pf_github")
        p_linkedin = st.text_input("💼 LinkedIn URL",  placeholder="https://linkedin.com/in/yourname", key="pf_linkedin", value=st.session_state.get("resume_linkedin", ""))
        p_skills   = st.text_input("💡 Skills (comma-separated)", placeholder="Python, React, AI, UI/UX, Node.js", key="pf_skills", value=st.session_state.get("resume_skills", ""))

    p_projects = st.text_area(
        "🚀 Projects (one per line, format: Name — Description)",
        placeholder="EduSphere AI — AI-powered learning platform built with Streamlit & GROQ\nPortfolio Site — Personal animated website with particle effects",
        height=100,
        key="pf_projects",
        value=st.session_state.get("resume_projects", ""),
    )
    p_education = st.text_input("🎓 Education", placeholder="BSc. CSIT — Lumbini ICT Campus (2020-2024)", key="pf_education", value=st.session_state.get("resume_education", ""))
    p_certs     = st.text_input("🏆 Certifications", placeholder="AWS Cloud Practitioner, Google IT Support", key="pf_certs", value=st.session_state.get("resume_certs", ""))

    _card_close()

    # Theme selector — visually styled
    st.markdown("#### 🎨 Choose Your Theme")
    theme_cols = st.columns(4)
    theme_options = [
        ("🌌", "Nebula Dark",    "#050a14", "#00f0ff", "#7c3aed"),
        ("🔥", "Crimson Forge",  "#0d0500", "#ff4500", "#ff8c00"),
        ("🌊", "Ocean Pulse",    "#020d1a", "#00d4ff", "#0077ff"),
        ("🌿", "Emerald Matrix", "#010d05", "#00ff41", "#39d353"),
    ]
    for i, (emoji, label, bg, acc, acc2) in enumerate(theme_options):
        with theme_cols[i]:
            st.markdown(
                f"""<div style="background:{bg};border:2px solid {acc};border-radius:16px;padding:14px;text-align:center;
                           box-shadow:0 0 20px {acc}40;cursor:pointer;margin-bottom:8px;">
                      <div style="font-size:1.8rem;margin-bottom:6px;">{emoji}</div>
                      <div style="color:{acc};font-weight:700;font-size:0.85rem;">{label}</div>
                      <div style="display:flex;gap:6px;justify-content:center;margin-top:8px;">
                        <span style="width:14px;height:14px;border-radius:50%;background:{acc};display:inline-block"></span>
                        <span style="width:14px;height:14px;border-radius:50%;background:{acc2};display:inline-block"></span>
                      </div>
                    </div>""",
                unsafe_allow_html=True,
            )

    p_theme = st.selectbox(
        "🖌️ Select Theme",
        ["🌌 Nebula Dark", "🔥 Crimson Forge", "🌊 Ocean Pulse", "🌿 Emerald Matrix"],
        key="pf_theme",
    )

    build_btn = st.button("⚡ Build My Website — Instantly!", key="btn_build_portfolio", use_container_width=True)

    if build_btn:
        if not p_name.strip():
            st.warning("Please enter your name to build the portfolio.")
            return

        data = {
            "name":      p_name.strip(),
            "title":     p_title.strip() or "Full-Stack Developer & AI Enthusiast",
            "email":     p_email.strip(),
            "phone":     p_phone.strip(),
            "location":  p_location.strip(),
            "github":    p_github.strip(),
            "linkedin":  p_linkedin.strip(),
            "about":     p_about.strip() or "I build thoughtful digital products and experiences.",
            "skills":    p_skills.strip() or "Python, JavaScript, React, AI, UI/UX",
            "projects":  p_projects.strip(),
            "education": p_education.strip(),
            "certs":     p_certs.strip(),
        }

        with logo_spinner("⚡ Compiling your portfolio website…"):
            import time as _time
            _time.sleep(0.3)
            html_output = _generate_portfolio_html(data, p_theme)

        st.session_state.portfolio_html   = html_output
        st.session_state.portfolio_name   = p_name
        st.success("✅ Your portfolio is ready! Scroll down to preview and download.", icon="🎉")

    # ── Display if built ──────────────────────────────────────────────────────
    if st.session_state.get("portfolio_html"):
        html_output = st.session_state.portfolio_html
        p_name_saved = st.session_state.get("portfolio_name", "portfolio")

        st.markdown("---")
        st.markdown("### 🖥️ Live Preview")
        st.caption("This is your actual website rendered live. Scroll inside the preview to explore all sections.")

        # Live preview iframe (custom iframe with allowfullscreen)
        import base64
        b64_html = base64.b64encode(html_output.encode("utf-8")).decode()
        iframe_code = f'<iframe src="data:text/html;base64,{b64_html}" width="100%" height="720" style="border:1px solid rgba(255,255,255,0.1); border-radius:12px; background:#000;" allowfullscreen="true" webkitallowfullscreen="true" mozallowfullscreen="true" allow="fullscreen"></iframe>'
        st.markdown(iframe_code, unsafe_allow_html=True)

        st.markdown("---")

        # ── Export row ────────────────────────────────────────────────────────
        st.markdown("### 📦 Export Your Website")
        col_dl, col_copy = st.columns(2)

        with col_dl:
            fname = f"{p_name_saved.replace(' ', '_').lower()}_portfolio.html"
            st.download_button(
                label="⬇️ Download HTML File",
                data=html_output.encode("utf-8"),
                file_name=fname,
                mime="text/html",
                key="dl_portfolio_html",
                use_container_width=True,
            )
            st.caption(f"Downloads as `{fname}` — open in any browser, no server needed!")

        with col_copy:
            # Claude-style copy: show the full HTML in a code area
            if st.button("📋 Show Full HTML to Copy", key="btn_show_html", use_container_width=True):
                st.session_state.show_portfolio_code = not st.session_state.get("show_portfolio_code", False)

        if st.session_state.get("show_portfolio_code", False):
            st.markdown("**📄 Your complete HTML file** — Select All (`Ctrl+A`) and Copy (`Ctrl+C`):")
            st.code(html_output, language="html")
            st.info("💡 Copy the code above and paste it into a `.html` file. Open in any browser to see your portfolio!", icon="💡")


def render_resume_builder() -> None:
    """📋 AI-Powered Resume & Portfolio Builder."""
    st.markdown("### 📋 Resume & Portfolio Builder")
    st.caption("Build your professional resume and generate an animated personal portfolio website.")

    tab_resume, tab_portfolio = st.tabs(["📋 Resume Builder", "🌐 Portfolio Website Builder"])

    with tab_resume:
        _render_resume_tab()

    with tab_portfolio:
        render_portfolio_builder()


def _render_resume_tab() -> None:
    """Inner resume builder content (Tab 1)."""
    st.markdown("##### 📋 AI-Powered Professional Resume")
    st.caption("Fill in your details and let AI craft a professional resume for you.")

    _card_open()
    col1, col2 = st.columns(2)

    with col1:
        full_name = st.text_input("👤 Full Name", placeholder="Sanjaya Kandel", key="resume_name")
        email = st.text_input("📧 Email", placeholder="your@email.com", key="resume_email")
        phone = st.text_input("📱 Phone", placeholder="+977-9800000000", key="resume_phone")
        location = st.text_input("📍 Location", placeholder="Kathmandu, Nepal", key="resume_location")
        linkedin = st.text_input("🔗 LinkedIn / Portfolio URL", placeholder="https://linkedin.com/in/yourname", key="resume_linkedin")

    with col2:
        objective = st.text_area(
            "🎯 Career Objective / Summary",
            placeholder="Brief summary of your career goals and key strengths…",
            height=100,
            key="resume_objective",
        )
        education = st.text_area(
            "🎓 Education",
            placeholder="BSc. CSIT — Lumbini ICT Campus (2020-2024)\nGPA: 3.5/4.0",
            height=100,
            key="resume_education",
        )
        skills = st.text_input(
            "💡 Key Skills",
            placeholder="Python, JavaScript, React, Machine Learning, SQL",
            key="resume_skills",
        )

    experience = st.text_area(
        "💼 Work Experience",
        placeholder="Software Developer — TechCorp Nepal (2024-Present)\n"
        "• Built REST APIs using Django and PostgreSQL\n"
        "• Led a team of 3 junior developers",
        height=120,
        key="resume_experience",
    )

    projects = st.text_area(
        "🚀 Projects",
        placeholder="EduSphere AI — AI-powered educational platform\n"
        "• Built with Streamlit, GROQ API, FAISS\n"
        "• Features: RAG chatbot, quiz generator, resume builder",
        height=100,
        key="resume_projects",
    )

    certifications = st.text_input(
        "🏆 Certifications / Awards",
        placeholder="AWS Cloud Practitioner, Google IT Support Certificate",
        key="resume_certs",
    )

    resume_style = st.selectbox(
        "📐 Resume Style",
        ["Professional Modern", "Academic / Research", "Creative Tech", "Minimal Clean"],
        key="resume_style",
    )

    _card_close()

    col_gen, col_preview = st.columns([1, 2])

    with col_gen:
        generate_btn = st.button("✨ Generate Resume", key="btn_generate_resume", use_container_width=True)

    if generate_btn:
        if not full_name.strip():
            st.warning("Please enter at least your full name.")
            return

        prompt = f"""Generate a professional resume in clean markdown format using the following details:

**Full Name:** {full_name}
**Email:** {email}
**Phone:** {phone}
**Location:** {location}
**LinkedIn / Portfolio:** {linkedin}

**Career Objective:**
{objective}

**Education:**
{education}

**Skills:** {skills}

**Work Experience:**
{experience}

**Projects:**
{projects}

**Certifications / Awards:** {certifications}

**Style Preference:** {resume_style}

Please create a well-structured, professional resume with:
1. Clear section headings
2. Bullet points for experiences and projects
3. Professional language and action verbs
4. Proper formatting for ATS compatibility
5. A brief professional summary if objective is provided

Output in clean, well-formatted markdown."""

        with logo_spinner("✨ AI is crafting your resume…"):
            resume_text = groq_chat(
                prompt,
                system="You are a professional resume writer and career advisor. "
                       "Create polished, ATS-compatible resumes that highlight the candidate's strengths. "
                       "Use clean markdown formatting with clear sections.",
                max_tokens=3000,
            )

        st.session_state.generated_resume = resume_text

    # Display resume if generated
    if st.session_state.get("generated_resume"):
        resume_text = st.session_state.generated_resume

        with col_preview:
            _card_open()
            st.markdown("#### 📄 Your Generated Resume")
            st.markdown("---")
            st.markdown(resume_text)
            _card_close()

        _render_export_buttons(resume_text, f"resume_{full_name.replace(' ', '_')}")


# ==============================================================================
# MODULE 13 — AI Image Generator
# ==============================================================================

def render_image_generator() -> None:
    """🎨 AI Image Generator — Rewrites prompt via Groq and pulls image from Pollinations."""
    st.markdown("### 🎨 AI Image Generator")
    
    col1, col2 = st.columns([1, 1])
    
    with col1:
        _card_open()
        st.markdown("#### 📝 Describe Your Image Concept")
        
        user_prompt = st.text_area(
            "Enter image concept / description",
            placeholder="e.g. A futuristic learning university on Mars with students studying using neon hologram displays...",
            height=120,
            key="img_gen_prompt_input"
        )
        
        enhance_prompt = st.checkbox("✨ Enhance prompt with Groq AI (Recommended)", value=True)
        
        style_preset = st.selectbox(
            "🎨 Render Style Preset",
            [
                "📸 Realistic Photorealistic / Cinematic",
                "🎨 Digital Concept Art",
                "✏️ Anime & Illustration",
                "🖌️ Classical Oil Painting"
            ],
            index=0
        )
        
        aspect_ratio = st.selectbox(
            "📐 Aspect Ratio",
            ["1:1 (Square)", "16:9 (Widescreen)", "9:16 (Portrait)", "4:3 (Classic)"],
            index=1
        )
        
        gen_clicked = st.button("🎨 Generate Image", use_container_width=True)
        _card_close()
        
    with col2:
        _card_open()
        st.markdown("#### 🖼️ Image Output")
        
        if gen_clicked:
            if not user_prompt.strip():
                st.warning("⚠️ Please provide a description first.")
            else:
                enhanced = user_prompt
                if enhance_prompt:
                    with logo_spinner("Expanding and refining prompt via Groq..."):
                        style_instructions = {
                            "📸 Realistic Photorealistic / Cinematic": (
                                "Make the output a highly detailed photorealistic masterpiece. "
                                "Specify camera parameters like 'shot on 35mm lens, f/1.8, cinematic volumetric lighting, 8k resolution, sharp focus, realistic textures'."
                            ),
                            "🎨 Digital Concept Art": "Make the output a gorgeous digital concept art piece, with vibrant colors, dramatic lighting, and creative detailing.",
                            "✏️ Anime & Illustration": "Make the output a beautiful anime style illustration, clean line art, and soft cel shading.",
                            "🖌️ Classical Oil Painting": "Make the output look like a classical oil painting masterpiece, with detailed brush strokes and canvas texture."
                        }[style_preset]
                        
                        enhance_instructions = (
                            "You are a professional prompt engineer for Stable Diffusion/Midjourney. "
                            "Expand the user's short concept into a detailed visual prompt. "
                            f"Focus style directions: {style_instructions} "
                            "Keep your response strictly to the enhanced prompt text under 100 words. Do NOT include any intro or chat."
                        )
                        enhanced = groq_chat(user_prompt, system=enhance_instructions)
                else:
                    # Append default style suffix
                    style_suffix = {
                        "📸 Realistic Photorealistic / Cinematic": ", photorealistic, cinematic lighting, 8k resolution, shot on 35mm lens, highly detailed, sharp focus, masterpiece",
                        "🎨 Digital Concept Art": ", digital concept art, trending on artstation, detailed, vibrant colors",
                        "✏️ Anime & Illustration": ", anime illustration, clean lines, detailed graphic style",
                        "🖌️ Classical Oil Painting": ", classical oil painting, fine art, visible brush strokes, canvas texture"
                    }[style_preset]
                    enhanced = user_prompt + style_suffix
                
                # Aspect Ratio to Dimensions
                dim_map = {
                    "1:1 (Square)": (1024, 1024),
                    "16:9 (Widescreen)": (1344, 768),
                    "9:16 (Portrait)": (768, 1344),
                    "4:3 (Classic)": (1024, 768)
                }
                width, height = dim_map[aspect_ratio]
                
                # Build URL using Pollinations
                import urllib.parse
                safe_prompt = urllib.parse.quote(enhanced)
                import time
                seed = int(time.time())
                img_url = f"https://image.pollinations.ai/prompt/{safe_prompt}?width={width}&height={height}&nologo=true&seed={seed}"
                
                with logo_spinner("Generating image via high-speed CDN..."):
                    # Render image
                    st.image(img_url, caption=f"Generated Image: {user_prompt}", use_container_width=True)
                    
                    try:
                        import requests
                        from io import BytesIO
                        from PIL import Image

                        resp = requests.get(img_url, timeout=20)
                        if resp.status_code == 200:
                            img_data = resp.content
                            img = Image.open(BytesIO(img_data))
                            
                            # Convert to PNG, JPG, WebP bytes
                            png_io = BytesIO()
                            img.save(png_io, format="PNG")
                            png_bytes = png_io.getvalue()
                            
                            jpg_io = BytesIO()
                            img.convert("RGB").save(jpg_io, format="JPEG")
                            jpg_bytes = jpg_io.getvalue()
                            
                            webp_io = BytesIO()
                            img.save(webp_io, format="WEBP")
                            webp_bytes = webp_io.getvalue()
                            
                            st.markdown("##### 📥 Download in Different Formats:")
                            col_d1, col_d2, col_d3 = st.columns(3)
                            with col_d1:
                                st.download_button("💾 Download .png", png_bytes, file_name="generated_image.png", mime="image/png", use_container_width=True)
                            with col_d2:
                                st.download_button("💾 Download .jpg", jpg_bytes, file_name="generated_image.jpg", mime="image/jpeg", use_container_width=True)
                            with col_d3:
                                st.download_button("💾 Download .webp", webp_bytes, file_name="generated_image.webp", mime="image/webp", use_container_width=True)
                        else:
                            st.error("Failed to download image bytes for format conversion.")
                    except Exception as e:
                        st.error(f"Error converting image formats: {e}")
        else:
            st.info("ℹ️ Enter a description and click Generate to see the image output here.")
        _card_close()


# ==============================================================================
# MODULE 14 — Interactive 3D Globe
# ==============================================================================

def render_globe_map() -> None:
    """🌍 Interactive 3D Globe Map for Academic Learning."""
    st.markdown("### 🌍 Interactive 3D Globe & Weather Intelligence")

    col1, col2 = st.columns([2, 1])

    # Predefined places of academic/scientific interest
    PREDEFINED_PLACES = {
        "Mount Everest": {"lat": 27.9881, "lng": 86.9250, "desc": "Highest point on Earth, located in the Himalayas on the border of Nepal and China."},
        "Great Pyramids of Giza": {"lat": 29.9792, "lng": 31.1342, "desc": "Ancient structures located near Cairo, Egypt, built as tombs for Pharaohs."},
        "Mariana Trench": {"lat": 11.3493, "lng": 142.1996, "desc": "Deepest known point in Earth's oceans, located in the Western Pacific."},
        "Amazon Rainforest": {"lat": -3.4653, "lng": -62.2159, "desc": "World's largest tropical rainforest, famous for its biodiverse ecosystem."},
        "CERN (Hadron Collider)": {"lat": 46.2333, "lng": 6.0491, "desc": "World's largest particle physics laboratory, located on the France-Switzerland border."}
    }

    with col1:
        _card_open()
        st.markdown("#### 🗺️ 3D Globe Viewer (Drag to rotate, scroll to zoom)")

        import json
        points_data = [
            {"lat": val["lat"], "lng": val["lng"], "name": name, "color": "red", "size": 0.5}
            for name, val in PREDEFINED_PLACES.items()
        ]

        html_code = f"""
        <link href="https://cesium.com/downloads/cesiumjs/releases/1.105/Build/Cesium/Widgets/widgets.css" rel="stylesheet">
        <style>
            html, body {{
                margin: 0;
                padding: 0;
                height: 100%;
                overflow: hidden;
                background: #000;
                font-family: 'Space Grotesk', sans-serif;
            }}
            #mapWrapper {{
                position: relative;
                width: 100%;
                height: 500px;
                border-radius: 8px;
                overflow: hidden;
            }}
            #mapWrapper:fullscreen {{
                width: 100% !important;
                height: 100% !important;
                border-radius: 0 !important;
            }}
            #cesiumContainer {{
                position: absolute;
                top: 0;
                left: 0;
                width: 100%;
                height: 100%;
            }}
            #searchPanel {{
                position: absolute;
                top: 12px;
                left: 50px;
                z-index: 1000;
                display: flex;
                gap: 6px;
                background: rgba(15, 23, 42, 0.85);
                backdrop-filter: blur(8px);
                border: 1px solid rgba(255, 255, 255, 0.12);
                border-radius: 6px;
                padding: 6px 10px;
                width: 280px;
                box-shadow: 0 4px 15px rgba(0,0,0,0.5);
            }}
            #mapSearchInput {{
                flex: 1;
                background: transparent;
                border: none;
                outline: none;
                color: #ececf1;
                font-size: 0.8rem;
                font-family: inherit;
            }}
            #mapSearchBtn {{
                background: #00f0ff;
                color: #0f172a;
                border: none;
                border-radius: 4px;
                padding: 4px 10px;
                font-size: 0.75rem;
                font-weight: bold;
                cursor: pointer;
                transition: all 0.2s;
            }}
            #mapSearchBtn:hover {{
                box-shadow: 0 0 8px rgba(0, 240, 255, 0.6);
            }}
            #detailsPanel {{
                position: absolute;
                bottom: 12px;
                left: 12px;
                right: 12px;
                z-index: 1000;
                background: rgba(15, 23, 42, 0.85);
                backdrop-filter: blur(8px);
                border: 1px solid rgba(255, 255, 255, 0.12);
                border-radius: 8px;
                padding: 10px 14px;
                font-size: 0.82rem;
                color: #ececf1;
                line-height: 1.4;
                box-shadow: 0 4px 15px rgba(0, 0, 0, 0.6);
            }}
            .cesium-viewer-bottom {{ display: none !important; }}
        </style>
        <div id="mapWrapper">
            <div id="cesiumContainer"></div>
            <div id="searchPanel">
                 <input type="text" id="mapSearchInput" placeholder="Search any place or question...">
                 <button id="mapSearchBtn">Search</button>
            </div>
            <div id="detailsPanel">
                🌍 <b>Interactive 3D Globe</b><br/>
                Click any marker or any place on the sphere to fly to it and view historical/geographical details here.
            </div>
        </div>
        <script src="https://cesium.com/downloads/cesiumjs/releases/1.105/Build/Cesium/Cesium.js"></script>
        <script>
            Cesium.Ion.defaultAccessToken = '';

            try {{
                const viewer = new Cesium.Viewer('cesiumContainer', {{
                    imageryProvider: new Cesium.UrlTemplateImageryProvider({{
                        url: 'https://mt1.google.com/vt/lyrs=y&x={{x}}&y={{y}}&z={{z}}',
                        maximumLevel: 20
                    }}),
                    fullscreenElement: document.getElementById('mapWrapper'),
                    fullscreenButton: true,
                    baseLayerPicker: false,
                    geocoder: false,
                    navigationHelpButton: false,
                    homeButton: false,
                    sceneModePicker: false,
                    timeline: false,
                    animation: false
                }});

                viewer.resize();
                setTimeout(() => {{ viewer.resize(); }}, 200);

                viewer.camera.setView({{
                    destination: Cesium.Cartesian3.fromDegrees(86.9250, 27.9881, 10000000.0)
                }});

                const pointsData = {json.dumps(points_data)};
                pointsData.forEach(p => {{
                    viewer.entities.add({{
                        position: Cesium.Cartesian3.fromDegrees(p.lng, p.lat),
                        billboard: {{
                            image: 'https://img.icons8.com/color/48/marker.png',
                            width: 32,
                            height: 32
                        }},
                        label: {{
                            text: p.name,
                            font: '14px Space Grotesk, sans-serif',
                            fillColor: Cesium.Color.AQUA,
                            outlineColor: Cesium.Color.BLACK,
                            outlineWidth: 2,
                            style: Cesium.LabelStyle.FILL_AND_OUTLINE,
                            verticalOrigin: Cesium.VerticalOrigin.BOTTOM,
                            pixelOffset: new Cesium.Cartesian2(0, -16)
                        }}
                    }});
                }});

                // selectedEntity handler for markers
                viewer.selectedEntityChanged.addEventListener(function(entity) {{
                    if (Cesium.defined(entity) && Cesium.defined(entity.label)) {{
                        const name = entity.label.text.getValue();
                        const descMap = {{
                            "Mount Everest": "Highest peak on Earth, located in the Himalayas. Famous for mountaineering and unique sub-zero alpine ecosystems.",
                            "Great Pyramids of Giza": "Ancient Egyptian pyramids near Cairo. Built during the Old Kingdom, they are marvels of ancient engineering.",
                            "Mariana Trench": "The deepest oceanic trench on Earth, situated in the western Pacific Ocean. Famous for extreme pressure and unique deep-sea life.",
                            "Amazon Rainforest": "World's largest tropical rainforest, stretching across South America. Home to unparalleled biodiversity and crucial global carbon sinks.",
                            "CERN (Hadron Collider)": "World's largest particle physics laboratory near Geneva. Famous for the Large Hadron Collider (LHC) and discovering the Higgs Boson."
                        }};
                        const desc = descMap[name] || "Famous academic and geographic site.";
                        document.getElementById('detailsPanel').innerHTML = "📍 <b>" + name + "</b><br/>" + desc;

                        const cartographic = Cesium.Ellipsoid.WGS84.cartesianToCartographic(entity.position.getValue(viewer.clock.currentTime));
                        const lat = Cesium.Math.toDegrees(cartographic.latitude);
                        const lng = Cesium.Math.toDegrees(cartographic.longitude);
                        window.parent.postMessage({{
                            type: 'globe_click',
                            name: name,
                            lat: lat,
                            lng: lng
                        }}, '*');
                    }}
                }});

                const handler = new Cesium.ScreenSpaceEventHandler(viewer.scene.canvas);
                handler.setInputAction(function (click) {{
                    viewer.selectedEntity = undefined;
                    const pickedPosition = viewer.camera.pickEllipsoid(click.position, viewer.scene.globe.ellipsoid);
                    if (pickedPosition) {{
                        const cartographic = Cesium.Cartographic.fromCartesian(pickedPosition);
                        const lat = Cesium.Math.toDegrees(cartographic.latitude);
                        const lng = Cesium.Math.toDegrees(cartographic.longitude);
                        const coordStr = lat.toFixed(4) + ', ' + lng.toFixed(4);

                        document.getElementById('detailsPanel').innerHTML = "📍 <b>Coordinates: " + coordStr + "</b><br/>Click 'Analyze Location' in the sidebar to search detailed historical and geographical knowledge via Groq & Tavily AI.";

                        window.parent.postMessage({{
                            type: 'globe_click',
                            name: coordStr,
                            lat: lat,
                            lng: lng
                        }}, '*');
                    }}
                }}, Cesium.ScreenSpaceEventType.LEFT_CLICK);

                // Custom search input listeners
                const searchInput = document.getElementById('mapSearchInput');
                const searchBtn = document.getElementById('mapSearchBtn');

                function triggerSearch() {{
                    const query = searchInput.value.trim();
                    if (!query) return;

                    document.getElementById('detailsPanel').innerHTML = "🔍 <b>Searching: " + query + "...</b><br/>Fetching summary details and flying to target location...";

                    const wikiUrl = 'https://en.wikipedia.org/api/rest_v1/page/summary/' + encodeURIComponent(query);
                    fetch(wikiUrl)
                        .then(res => res.json())
                        .then(wikiData => {{
                            let desc = "No direct summary found. Click 'Analyze Location' in the sidebar to search via Groq AI.";
                            if (wikiData.extract) {{
                                desc = wikiData.extract;
                            }}

                            const geoUrl = 'https://nominatim.openstreetmap.org/search?format=json&q=' + encodeURIComponent(query) + '&limit=1';
                            fetch(geoUrl)
                                .then(r => r.json())
                                .then(geoData => {{
                                    if (geoData && geoData.length > 0) {{
                                        const lat = parseFloat(geoData[0].lat);
                                        const lon = parseFloat(geoData[0].lon);

                                        document.getElementById('detailsPanel').innerHTML = "📍 <b>" + (wikiData.title || query) + "</b><br/>" + desc;

                                        viewer.camera.flyTo({{
                                            destination: Cesium.Cartesian3.fromDegrees(lon, lat, 25000.0),
                                            duration: 2.0
                                        }});

                                        window.parent.postMessage({{
                                            type: 'globe_click',
                                            name: query,
                                            lat: lat,
                                            lng: lon
                                        }}, '*');
                                    }} else {{
                                        document.getElementById('detailsPanel').innerHTML = "📍 <b>" + (wikiData.title || query) + "</b><br/>" + desc + "<br/><span style='color:#ef4444;'>Failed to geocode location coordinates on map.</span>";
                                    }}
                                }})
                                .catch(err => {{
                                    document.getElementById('detailsPanel').innerHTML = "📍 <b>" + (wikiData.title || query) + "</b><br/>" + desc;
                                }});
                        }})
                        .catch(err => {{
                            document.getElementById('detailsPanel').innerHTML = "❌ <b>Error</b><br/>Failed to fetch search results.";
                        }});
                }}

                searchBtn.addEventListener('click', triggerSearch);
                searchInput.addEventListener('keypress', function(e) {{
                    if (e.key === 'Enter') {{
                        triggerSearch();
                    }}
                }});
            }} catch (e) {{
                console.error("Cesium failed to load", e);
                document.getElementById('cesiumContainer').innerHTML = "<div style='color:#ef4444; padding:20px; font-family:sans-serif;'>Failed to load 3D Globe: " + e.message + "</div>";
            }}
        </script>
        """
        st.components.v1.html(html_code, height=520)
        _card_close()

        st.markdown(
            """
            <script>
                if (!window.globeListenerAdded) {
                    window.addEventListener('message', function(event) {
                        if (event.data && event.data.type === 'globe_click') {
                            const name = event.data.name;
                            const nameInput = document.querySelector('input[aria-label="Location Name or Custom Coordinates"]');
                            if (nameInput) {
                                const nativeSetter = Object.getOwnPropertyDescriptor(HTMLInputElement.prototype, 'value').set;
                                nativeSetter.call(nameInput, name);
                                nameInput.dispatchEvent(new Event('input', { bubbles: true }));
                                setTimeout(() => {
                                    const analyzeBtn = Array.from(document.querySelectorAll('button')).find(btn => btn.textContent.includes('Analyze Location'));
                                    if (analyzeBtn) {
                                        analyzeBtn.click();
                                    }
                                }, 100);
                            }
                        }
                    });
                    window.globeListenerAdded = true;
                }
            </script>
            """,
            unsafe_allow_html=True
        )

    with col2:
        # ── Location Analyzer ──
        _card_open()
        st.markdown("#### 🔬 Location Analyzer")

        selected_place = st.text_input(
            "Location Name or Custom Coordinates",
            value="Mount Everest",
            help="Click a red label on the globe or type a custom place name/coordinates."
        )

        analyze_clicked = st.button("🔍 Analyze Location", use_container_width=True)
        _card_close()

        if analyze_clicked or selected_place:
            _card_open()
            st.markdown(f"##### 📚 Academic Report: **{selected_place}**")

            with logo_spinner(f"Retrieving research details for {selected_place}..."):
                search_query = f"{selected_place} geographical scientific historical facts summary"
                web_results = duckduckgo_search(search_query)
                web_context = ""
                if web_results:
                    web_context = "\n".join([res['snippet'] for res in web_results])

                system_role = (
                    "You are a helpful science, history, and geography AI tutor. "
                    "Write a neat, structured academic report on the requested location containing: "
                    "1. Coordinates & Geographic Overview "
                    "2. Historical & Cultural significance "
                    "3. Scientific / Geological / Environmental importance "
                    "4. Intriguing facts for students."
                )
                report = groq_chat(
                    f"Context: {web_context}\n\nProvide an academic report for: {selected_place}",
                    system=system_role
                )
                st.markdown(report)
            _card_close()



# ==============================================================================
# MODULE 14b — Standalone Weather Forecast (branch of 3D Globe)
# ==============================================================================

def render_weather_forecast() -> None:
    """⛅ Weather Forecast — live weather and 3-day forecast for any city, with AI insight."""
    st.markdown("### ⛅ Weather Forecast")
    st.caption("🌍 Branch of the Interactive 3D Globe · Powered by Open-Meteo (free, no API key required)")

    import urllib.request
    import urllib.parse
    import json as _json
    import datetime as _dt

    # ── WMO code → emoji + label ──
    def wx_icon(code: int) -> str:
        if code == 0:          return "☀️"
        if code in (1, 2, 3):  return "⛅"
        if code in (45, 48):   return "🌫️"
        if code in (51, 53, 55, 61, 63, 65): return "🌧️"
        if code in (71, 73, 75): return "❄️"
        if code in (80, 81, 82): return "🌦️"
        if code in (95, 96, 99): return "⛈️"
        return "🌡️"

    def wx_label(code: int) -> str:
        labels = {0:"Clear sky", 1:"Mainly clear", 2:"Partly cloudy", 3:"Overcast",
                  45:"Fog", 48:"Icy fog", 51:"Light drizzle", 53:"Drizzle",
                  55:"Dense drizzle", 61:"Slight rain", 63:"Moderate rain", 65:"Heavy rain",
                  71:"Slight snow", 73:"Moderate snow", 75:"Heavy snow",
                  80:"Slight showers", 81:"Moderate showers", 82:"Violent showers",
                  95:"Thunderstorm", 96:"Thunderstorm w/ hail", 99:"Heavy thunderstorm"}
        return labels.get(code, "Unknown")

    # ── Search bar ──
    col_search, col_btn = st.columns([4, 1], gap="small")
    with col_search:
        weather_city = st.text_input(
            "🏙️ Enter City Name",
            value=st.session_state.get("wx_city", "Kathmandu"),
            key="wf_city_input",
            placeholder="e.g. London, Tokyo, New York, Kathmandu...",
            label_visibility="collapsed"
        )
    with col_btn:
        get_btn = st.button("🌡️ Get Weather", use_container_width=True, key="wf_get_btn")

    if not (get_btn and weather_city.strip()):
        st.markdown(
            """
            <div style="text-align:center; padding:40px; color:var(--sub); font-size:0.9rem;">
                🌍 Enter a city name above and click <b>Get Weather</b> to see live conditions
            </div>
            """,
            unsafe_allow_html=True
        )
        return

    st.session_state["wx_city"] = weather_city.strip()

    with st.spinner(f"Fetching weather for {weather_city}..."):
        try:
            # Geocoding
            geo_url = (
                "https://geocoding-api.open-meteo.com/v1/search"
                f"?name={urllib.parse.quote(weather_city)}&count=5&language=en&format=json"
            )
            with urllib.request.urlopen(geo_url, timeout=8) as resp:
                geo_data = _json.loads(resp.read())

            if not geo_data.get("results"):
                st.error(f"❌ City '{weather_city}' not found. Try a different spelling.")
                return

            r = geo_data["results"][0]
            lat_w, lon_w = r["latitude"], r["longitude"]
            country = r.get("country", "")
            admin  = r.get("admin1", "")
            city_name = r["name"]

            # Open-Meteo API
            wx_url = (
                f"https://api.open-meteo.com/v1/forecast"
                f"?latitude={lat_w}&longitude={lon_w}"
                f"&current=temperature_2m,relative_humidity_2m,wind_speed_10m,"
                f"weather_code,apparent_temperature,precipitation,wind_direction_10m,"
                f"surface_pressure,visibility"
                f"&hourly=temperature_2m,weather_code"
                f"&daily=temperature_2m_max,temperature_2m_min,weather_code,"
                f"precipitation_sum,wind_speed_10m_max,sunrise,sunset,uv_index_max"
                f"&timezone=auto&forecast_days=7"
            )
            with urllib.request.urlopen(wx_url, timeout=10) as resp2:
                wx = _json.loads(resp2.read())

            cur   = wx["current"]
            daily = wx["daily"]

            temp   = cur["temperature_2m"]
            feels  = cur["apparent_temperature"]
            hum    = cur["relative_humidity_2m"]
            wind   = cur["wind_speed_10m"]
            wdir   = cur.get("wind_direction_10m", 0)
            precip = cur["precipitation"]
            press  = cur.get("surface_pressure", "—")
            vis    = cur.get("visibility", "—")
            code   = cur["weather_code"]
            icon   = wx_icon(code)
            label  = wx_label(code)

            # ── Main hero card ──
            st.markdown(
                f"""
                <div style="background:linear-gradient(135deg,rgba(0,180,255,0.12),rgba(100,60,255,0.08));
                            border:1px solid rgba(0,200,255,0.2); border-radius:20px; padding:28px 32px;
                            margin-bottom:20px;">
                    <div style="display:flex; align-items:center; justify-content:space-between; flex-wrap:wrap; gap:16px;">
                        <div>
                            <div style="font-size:0.78rem; color:var(--sub); letter-spacing:1px; text-transform:uppercase; margin-bottom:4px;">
                                Live Conditions · {_dt.datetime.now().strftime('%H:%M, %d %b %Y')}
                            </div>
                            <div style="font-size:1.3rem; font-weight:700; color:var(--accent); margin-bottom:2px;">
                                📍 {city_name}, {admin}, {country}
                            </div>
                            <div style="font-size:0.8rem; color:var(--sub);">
                                {lat_w:.4f}°N · {lon_w:.4f}°E · {label}
                            </div>
                        </div>
                        <div style="text-align:right;">
                            <div style="font-size:4rem; line-height:1;">{icon}</div>
                            <div style="font-size:2.6rem; font-weight:900; color:var(--text); font-family:'Orbitron',monospace;">{temp}°C</div>
                            <div style="font-size:0.85rem; color:var(--sub);">Feels like {feels}°C</div>
                        </div>
                    </div>
                    <div style="display:grid; grid-template-columns:repeat(4,1fr); gap:14px; margin-top:22px;">
                        <div style="background:rgba(255,255,255,0.05); border-radius:12px; padding:12px; text-align:center;">
                            <div style="font-size:1.3rem;">💧</div>
                            <div style="font-size:1rem; font-weight:700; color:var(--text);">{hum}%</div>
                            <div style="font-size:0.68rem; color:var(--sub);">Humidity</div>
                        </div>
                        <div style="background:rgba(255,255,255,0.05); border-radius:12px; padding:12px; text-align:center;">
                            <div style="font-size:1.3rem;">💨</div>
                            <div style="font-size:1rem; font-weight:700; color:var(--text);">{wind} km/h</div>
                            <div style="font-size:0.68rem; color:var(--sub);">Wind ({wdir}°)</div>
                        </div>
                        <div style="background:rgba(255,255,255,0.05); border-radius:12px; padding:12px; text-align:center;">
                            <div style="font-size:1.3rem;">🌧️</div>
                            <div style="font-size:1rem; font-weight:700; color:var(--text);">{precip} mm</div>
                            <div style="font-size:0.68rem; color:var(--sub);">Precipitation</div>
                        </div>
                        <div style="background:rgba(255,255,255,0.05); border-radius:12px; padding:12px; text-align:center;">
                            <div style="font-size:1.3rem;">🔵</div>
                            <div style="font-size:1rem; font-weight:700; color:var(--text);">{press} hPa</div>
                            <div style="font-size:0.68rem; color:var(--sub);">Pressure</div>
                        </div>
                    </div>
                </div>
                """,
                unsafe_allow_html=True
            )

            # ── 7-Day Forecast ──
            st.markdown("#### 📅 7-Day Forecast")
            days_count = min(7, len(daily["time"]))
            day_cols = st.columns(days_count)
            for i in range(days_count):
                day_label = _dt.date.fromisoformat(daily["time"][i]).strftime("%a\n%d %b")
                dc = daily["weather_code"][i]
                sunrise = daily.get("sunrise", ["—"])[i].split("T")[-1][:5] if "sunrise" in daily else "—"
                sunset  = daily.get("sunset", ["—"])[i].split("T")[-1][:5] if "sunset" in daily else "—"
                uv = daily.get("uv_index_max", [0])[i] if "uv_index_max" in daily else 0
                with day_cols[i]:
                    is_today = i == 0
                    border = "border:1px solid rgba(0,240,255,0.35);" if is_today else "border:1px solid rgba(255,255,255,0.07);"
                    st.markdown(
                        f"""
                        <div style="background:rgba(255,255,255,{'0.08' if is_today else '0.03'});
                                    {border} border-radius:14px; padding:12px 8px; text-align:center;">
                            <div style="font-size:0.68rem; color:{'var(--accent)' if is_today else 'var(--sub)'}; font-weight:{'700' if is_today else '400'}; white-space:pre-line; margin-bottom:6px;">
                                {day_label}{'\\n(Today)' if is_today else ''}
                            </div>
                            <div style="font-size:1.8rem; margin:4px 0;">{wx_icon(dc)}</div>
                            <div style="color:var(--accent); font-weight:700; font-size:0.88rem;">{daily['temperature_2m_max'][i]}°</div>
                            <div style="color:var(--sub); font-size:0.78rem;">{daily['temperature_2m_min'][i]}°</div>
                            <div style="font-size:0.65rem; color:var(--sub); margin-top:6px; line-height:1.6;">
                                💧{daily['precipitation_sum'][i]}mm<br>
                                ☀{sunrise}<br>
                                🌙{sunset}
                            </div>
                        </div>
                        """,
                        unsafe_allow_html=True
                    )

            # ── Hourly temperature chart (simple HTML sparkline) ──
            st.markdown("<div style='margin-top:20px;'></div>", unsafe_allow_html=True)
            st.markdown("#### 📈 Hourly Temperature (Next 24h)")
            hourly_temps = wx["hourly"]["temperature_2m"][:24]
            hourly_times = [t.split("T")[1][:5] for t in wx["hourly"]["time"][:24]]
            chart_data = {"Time": hourly_times, "Temp (°C)": hourly_temps}
            try:
                import pandas as pd
                df = pd.DataFrame(chart_data).set_index("Time")
                st.line_chart(df, use_container_width=True, height=180)
            except Exception:
                st.caption("Chart unavailable (pandas not installed).")

            # ── AI Weather Commentary ──
            st.markdown("<div style='margin-top:20px;'></div>", unsafe_allow_html=True)
            _card_open()
            st.markdown("#### 🤖 AI Weather Insight")
            with logo_spinner("Generating AI weather report..."):
                wx_prompt = (
                    f"Current weather in {city_name}, {country}: {temp}°C (feels {feels}°C), "
                    f"humidity {hum}%, wind {wind} km/h ({wdir}°), precipitation {precip}mm, "
                    f"pressure {press} hPa. Condition: {label}. "
                    f"7-day max temps: {daily['temperature_2m_max'][:7]}. "
                    f"Write a friendly 3-sentence weather summary, mention any notable conditions, "
                    f"and give 2 practical tips for students going outside today."
                )
                insight = groq_chat(wx_prompt, system="You are a helpful and friendly meteorology assistant.")
                st.markdown(insight)
            _card_close()

            # ── Other cities (if multiple results) ──
            if len(geo_data["results"]) > 1:
                st.markdown("<div style='margin-top:12px;'></div>", unsafe_allow_html=True)
                with st.expander("📍 Other matching cities"):
                    for alt in geo_data["results"][1:]:
                        st.markdown(
                            f"**{alt['name']}** — {alt.get('admin1','')}, {alt.get('country','')} "
                            f"({alt['latitude']:.3f}°, {alt['longitude']:.3f}°)"
                        )

        except Exception as err:
            st.error(f"❌ Weather fetch failed: {err}")





        get_weather = st.button("🌡️ Get Weather", use_container_width=True, key="get_weather_btn")

        if get_weather and weather_city.strip():
            with st.spinner(f"Fetching weather for {weather_city}..."):
                try:
                    import urllib.request
                    import json as _json

                    # Open-Meteo geocoding (free, no key)
                    geo_url = (
                        f"https://geocoding-api.open-meteo.com/v1/search"
                        f"?name={urllib.parse.quote(weather_city)}&count=1&language=en&format=json"
                    )
                    import urllib.parse
                    with urllib.request.urlopen(geo_url, timeout=8) as resp:
                        geo_data = _json.loads(resp.read())

                    if not geo_data.get("results"):
                        st.error(f"❌ City '{weather_city}' not found.")
                    else:
                        r = geo_data["results"][0]
                        lat_w, lon_w = r["latitude"], r["longitude"]
                        country = r.get("country", "")
                        admin = r.get("admin1", "")

                        # Open-Meteo weather API (free, no key)
                        wx_url = (
                            f"https://api.open-meteo.com/v1/forecast"
                            f"?latitude={lat_w}&longitude={lon_w}"
                            f"&current=temperature_2m,relative_humidity_2m,wind_speed_10m,"
                            f"weather_code,apparent_temperature,precipitation"
                            f"&daily=temperature_2m_max,temperature_2m_min,weather_code,"
                            f"precipitation_sum,wind_speed_10m_max"
                            f"&timezone=auto&forecast_days=4"
                        )
                        with urllib.request.urlopen(wx_url, timeout=8) as resp2:
                            wx = _json.loads(resp2.read())

                        cur = wx["current"]
                        daily = wx["daily"]

                        # WMO weather code → emoji
                        def wx_icon(code):
                            if code == 0: return "☀️"
                            if code in (1, 2, 3): return "⛅"
                            if code in (45, 48): return "🌫️"
                            if code in (51, 53, 55, 61, 63, 65): return "🌧️"
                            if code in (71, 73, 75): return "❄️"
                            if code in (80, 81, 82): return "🌦️"
                            if code in (95, 96, 99): return "⛈️"
                            return "🌡️"

                        temp = cur["temperature_2m"]
                        feels = cur["apparent_temperature"]
                        hum = cur["relative_humidity_2m"]
                        wind = cur["wind_speed_10m"]
                        precip = cur["precipitation"]
                        icon = wx_icon(cur["weather_code"])

                        st.markdown(
                            f"""
                            <div style="background:rgba(0,240,255,0.06); border:1px solid rgba(0,240,255,0.18);
                                        border-radius:14px; padding:16px 18px; margin-bottom:12px;">
                                <div style="font-size:2.6rem; text-align:center; margin-bottom:4px;">{icon}</div>
                                <div style="text-align:center; font-size:1rem; font-weight:700;
                                            color:var(--accent); margin-bottom:2px;">
                                    {r['name']}, {admin}, {country}
                                </div>
                                <div style="text-align:center; font-size:2rem; font-weight:900;
                                            color:var(--text); margin-bottom:8px;">{temp}°C</div>
                                <div style="display:grid; grid-template-columns:1fr 1fr; gap:8px; font-size:0.78rem; color:var(--sub);">
                                    <div>🌡️ Feels like: <b style="color:var(--text);">{feels}°C</b></div>
                                    <div>💧 Humidity: <b style="color:var(--text);">{hum}%</b></div>
                                    <div>💨 Wind: <b style="color:var(--text);">{wind} km/h</b></div>
                                    <div>🌧️ Precip: <b style="color:var(--text);">{precip} mm</b></div>
                                </div>
                            </div>
                            """,
                            unsafe_allow_html=True
                        )

                        # 3-day forecast
                        st.markdown("**📅 3-Day Forecast**")
                        import datetime as _dt
                        days = daily["time"][1:4]
                        max_t = daily["temperature_2m_max"][1:4]
                        min_t = daily["temperature_2m_min"][1:4]
                        codes = daily["weather_code"][1:4]
                        precips = daily["precipitation_sum"][1:4]
                        winds = daily["wind_speed_10m_max"][1:4]

                        fc_html = '<div style="display:grid; grid-template-columns:1fr 1fr 1fr; gap:8px; margin-top:8px;">'
                        for i in range(3):
                            day_label = _dt.date.fromisoformat(days[i]).strftime("%a %d %b")
                            fc_html += f"""
                            <div style="background:rgba(255,255,255,0.04); border:1px solid rgba(255,255,255,0.08);
                                        border-radius:10px; padding:10px 8px; text-align:center; font-size:0.78rem;">
                                <div style="font-size:0.7rem; color:var(--sub); margin-bottom:4px;">{day_label}</div>
                                <div style="font-size:1.6rem;">{wx_icon(codes[i])}</div>
                                <div style="color:var(--accent); font-weight:700;">{max_t[i]}°C</div>
                                <div style="color:var(--sub);">{min_t[i]}°C</div>
                                <div style="color:var(--sub); font-size:0.68rem; margin-top:4px;">
                                    💧{precips[i]}mm 💨{winds[i]}km/h
                                </div>
                            </div>
                            """
                        fc_html += "</div>"
                        st.markdown(fc_html, unsafe_allow_html=True)

                        # AI weather commentary
                        with logo_spinner("Generating AI weather insight..."):
                            wx_prompt = (
                                f"The current weather in {r['name']}, {country} is {temp}°C (feels like {feels}°C), "
                                f"humidity {hum}%, wind {wind} km/h, precipitation {precip}mm. "
                                f"Give a 2-sentence friendly weather summary and one practical tip for students."
                            )
                            wx_insight = groq_chat(wx_prompt, system="You are a helpful meteorology assistant.")
                            st.info(f"🤖 **AI Insight:** {wx_insight}")

                except Exception as wx_err:
                    st.error(f"❌ Weather fetch failed: {wx_err}")

        _card_close()


        _card_open()
        st.markdown("#### 🗺️ 3D Globe Viewer (Drag to rotate, scroll to zoom)")
        
        # Predefined places of academic/scientific interest
        PREDEFINED_PLACES = {
            "Mount Everest": {"lat": 27.9881, "lng": 86.9250, "desc": "Highest point on Earth, located in the Himalayas on the border of Nepal and China."},
            "Great Pyramids of Giza": {"lat": 29.9792, "lng": 31.1342, "desc": "Ancient structures located near Cairo, Egypt, built as tombs for Pharaohs."},
            "Mariana Trench": {"lat": 11.3493, "lng": 142.1996, "desc": "Deepest known point in Earth's oceans, located in the Western Pacific."},
            "Amazon Rainforest": {"lat": -3.4653, "lng": -62.2159, "desc": "World's largest tropical rainforest, famous for its biodiverse ecosystem."},
            "CERN (Hadron Collider)": {"lat": 46.2333, "lng": 6.0491, "desc": "World's largest particle physics laboratory, located on the France-Switzerland border."}
        }
        import json
        points_data = [
            {"lat": val["lat"], "lng": val["lng"], "name": name, "color": "red", "size": 0.5}
            for name, val in PREDEFINED_PLACES.items()
        ]
        
        html_code = f"""
        <link href="https://cesium.com/downloads/cesiumjs/releases/1.105/Build/Cesium/Widgets/widgets.css" rel="stylesheet">
        <style>
            html, body {{
                margin: 0;
                padding: 0;
                height: 100%;
                overflow: hidden;
                background: #000;
                font-family: 'Space Grotesk', sans-serif;
            }}
            #mapWrapper {{
                position: relative;
                width: 100%;
                height: 500px;
                border-radius: 8px;
                overflow: hidden;
            }}
            #mapWrapper:fullscreen {{
                width: 100% !important;
                height: 100% !important;
                border-radius: 0 !important;
            }}
            #cesiumContainer {{
                position: absolute;
                top: 0;
                left: 0;
                width: 100%;
                height: 100%;
            }}
            #searchPanel {{
                position: absolute;
                top: 12px;
                left: 50px;
                z-index: 1000;
                display: flex;
                gap: 6px;
                background: rgba(15, 23, 42, 0.85);
                backdrop-filter: blur(8px);
                border: 1px solid rgba(255, 255, 255, 0.12);
                border-radius: 6px;
                padding: 6px 10px;
                width: 280px;
                box-shadow: 0 4px 15px rgba(0,0,0,0.5);
            }}
            #mapSearchInput {{
                flex: 1;
                background: transparent;
                border: none;
                outline: none;
                color: #ececf1;
                font-size: 0.8rem;
                font-family: inherit;
            }}
            #mapSearchBtn {{
                background: #00f0ff;
                color: #0f172a;
                border: none;
                border-radius: 4px;
                padding: 4px 10px;
                font-size: 0.75rem;
                font-weight: bold;
                cursor: pointer;
                transition: all 0.2s;
            }}
            #mapSearchBtn:hover {{
                box-shadow: 0 0 8px rgba(0, 240, 255, 0.6);
            }}
            #detailsPanel {{
                position: absolute;
                bottom: 12px;
                left: 12px;
                right: 12px;
                z-index: 1000;
                background: rgba(15, 23, 42, 0.85);
                backdrop-filter: blur(8px);
                border: 1px solid rgba(255, 255, 255, 0.12);
                border-radius: 8px;
                padding: 10px 14px;
                font-size: 0.82rem;
                color: #ececf1;
                line-height: 1.4;
                box-shadow: 0 4px 15px rgba(0, 0, 0, 0.6);
            }}
            .cesium-viewer-bottom {{ display: none !important; }}
        </style>
        <div id="mapWrapper">
            <div id="cesiumContainer"></div>
            <div id="searchPanel">
                 <input type="text" id="mapSearchInput" placeholder="Search any place or question...">
                 <button id="mapSearchBtn">Search</button>
            </div>
            <div id="detailsPanel">
                🌍 <b>Interactive 3D Globe</b><br/>
                Click any marker or any place on the sphere to fly to it and view historical/geographical details here.
            </div>
        </div>
        <script src="https://cesium.com/downloads/cesiumjs/releases/1.105/Build/Cesium/Cesium.js"></script>
        <script>
            Cesium.Ion.defaultAccessToken = '';
            
            try {{
                const viewer = new Cesium.Viewer('cesiumContainer', {{
                    imageryProvider: new Cesium.UrlTemplateImageryProvider({{
                        url: 'https://mt1.google.com/vt/lyrs=y&x={{x}}&y={{y}}&z={{z}}',
                        maximumLevel: 20
                    }}),
                    fullscreenElement: document.getElementById('mapWrapper'),
                    fullscreenButton: true,
                    baseLayerPicker: false,
                    geocoder: false,
                    navigationHelpButton: false,
                    homeButton: false,
                    sceneModePicker: false,
                    timeline: false,
                    animation: false
                }});
                
                viewer.resize();
                setTimeout(() => {{ viewer.resize(); }}, 200);

                viewer.camera.setView({{
                    destination: Cesium.Cartesian3.fromDegrees(86.9250, 27.9881, 10000000.0)
                }});

                const pointsData = {json.dumps(points_data)};
                pointsData.forEach(p => {{
                    viewer.entities.add({{
                        position: Cesium.Cartesian3.fromDegrees(p.lng, p.lat),
                        billboard: {{
                            image: 'https://img.icons8.com/color/48/marker.png',
                            width: 32,
                            height: 32
                        }},
                        label: {{
                            text: p.name,
                            font: '14px Space Grotesk, sans-serif',
                            fillColor: Cesium.Color.AQUA,
                            outlineColor: Cesium.Color.BLACK,
                            outlineWidth: 2,
                            style: Cesium.LabelStyle.FILL_AND_OUTLINE,
                            verticalOrigin: Cesium.VerticalOrigin.BOTTOM,
                            pixelOffset: new Cesium.Cartesian2(0, -16)
                        }}
                    }});
                }});

                // selectedEntity handler for markers
                viewer.selectedEntityChanged.addEventListener(function(entity) {{
                    if (Cesium.defined(entity) && Cesium.defined(entity.label)) {{
                        const name = entity.label.text.getValue();
                        const descMap = {{
                            "Mount Everest": "Highest peak on Earth, located in the Himalayas. Famous for mountaineering and unique sub-zero alpine ecosystems.",
                            "Great Pyramids of Giza": "Ancient Egyptian pyramids near Cairo. Built during the Old Kingdom, they are marvels of ancient engineering.",
                            "Mariana Trench": "The deepest oceanic trench on Earth, situated in the western Pacific Ocean. Famous for extreme pressure and unique deep-sea life.",
                            "Amazon Rainforest": "World's largest tropical rainforest, stretching across South America. Home to unparalleled biodiversity and crucial global carbon sinks.",
                            "CERN (Hadron Collider)": "World's largest particle physics laboratory near Geneva. Famous for the Large Hadron Collider (LHC) and discovering the Higgs Boson."
                        }};
                        const desc = descMap[name] || "Famous academic and geographic site.";
                        document.getElementById('detailsPanel').innerHTML = "📍 <b>" + name + "</b><br/>" + desc;
                        
                        const cartographic = Cesium.Ellipsoid.WGS84.cartesianToCartographic(entity.position.getValue(viewer.clock.currentTime));
                        const lat = Cesium.Math.toDegrees(cartographic.latitude);
                        const lng = Cesium.Math.toDegrees(cartographic.longitude);
                        window.parent.postMessage({{
                            type: 'globe_click',
                            name: name,
                            lat: lat,
                            lng: lng
                        }}, '*');
                    }}
                }});

                const handler = new Cesium.ScreenSpaceEventHandler(viewer.scene.canvas);
                handler.setInputAction(function (click) {{
                    viewer.selectedEntity = undefined;
                    const pickedPosition = viewer.camera.pickEllipsoid(click.position, viewer.scene.globe.ellipsoid);
                    if (pickedPosition) {{
                        const cartographic = Cesium.Cartographic.fromCartesian(pickedPosition);
                        const lat = Cesium.Math.toDegrees(cartographic.latitude);
                        const lng = Cesium.Math.toDegrees(cartographic.longitude);
                        const coordStr = lat.toFixed(4) + ', ' + lng.toFixed(4);
                        
                        document.getElementById('detailsPanel').innerHTML = "📍 <b>Coordinates: " + coordStr + "</b><br/>Click 'Analyze Location' in the sidebar to search detailed historical and geographical knowledge via Groq & Tavily AI.";
                        
                        window.parent.postMessage({{
                            type: 'globe_click',
                            name: coordStr,
                            lat: lat,
                            lng: lng
                        }}, '*');
                    }}
                }}, Cesium.ScreenSpaceEventType.LEFT_CLICK);

                // Custom search input listeners
                const searchInput = document.getElementById('mapSearchInput');
                const searchBtn = document.getElementById('mapSearchBtn');

                function triggerSearch() {{
                    const query = searchInput.value.trim();
                    if (!query) return;
                    
                    document.getElementById('detailsPanel').innerHTML = "🔍 <b>Searching: " + query + "...</b><br/>Fetching summary details and flying to target location...";
                    
                    const wikiUrl = 'https://en.wikipedia.org/api/rest_v1/page/summary/' + encodeURIComponent(query);
                    fetch(wikiUrl)
                        .then(res => res.json())
                        .then(wikiData => {{
                            let desc = "No direct summary found. Click 'Analyze Location' in the sidebar to search via Groq AI.";
                            if (wikiData.extract) {{
                                desc = wikiData.extract;
                            }}
                            
                            const geoUrl = 'https://nominatim.openstreetmap.org/search?format=json&q=' + encodeURIComponent(query) + '&limit=1';
                            fetch(geoUrl)
                                .then(r => r.json())
                                .then(geoData => {{
                                    if (geoData && geoData.length > 0) {{
                                        const lat = parseFloat(geoData[0].lat);
                                        const lon = parseFloat(geoData[0].lon);
                                        
                                        document.getElementById('detailsPanel').innerHTML = "📍 <b>" + (wikiData.title || query) + "</b><br/>" + desc;
                                        
                                        viewer.camera.flyTo({{
                                            destination: Cesium.Cartesian3.fromDegrees(lon, lat, 25000.0),
                                            duration: 2.0
                                        }});
                                        
                                        window.parent.postMessage({{
                                            type: 'globe_click',
                                            name: query,
                                            lat: lat,
                                            lng: lon
                                        }}, '*');
                                    }} else {{
                                        document.getElementById('detailsPanel').innerHTML = "📍 <b>" + (wikiData.title || query) + "</b><br/>" + desc + "<br/><span style='color:#ef4444;'>Failed to geocode location coordinates on map.</span>";
                                    }}
                                }})
                                .catch(err => {{
                                    document.getElementById('detailsPanel').innerHTML = "📍 <b>" + (wikiData.title || query) + "</b><br/>" + desc;
                                }});
                        }})
                        .catch(err => {{
                            document.getElementById('detailsPanel').innerHTML = "❌ <b>Error</b><br/>Failed to fetch search results.";
                        }});
                }}

                searchBtn.addEventListener('click', triggerSearch);
                searchInput.addEventListener('keypress', function(e) {{
                    if (e.key === 'Enter') {{
                        triggerSearch();
                    }}
                }});
            }} catch (e) {{
                console.error("Cesium failed to load", e);
                document.getElementById('cesiumContainer').innerHTML = "<div style='color:#ef4444; padding:20px; font-family:sans-serif;'>Failed to load 3D Globe: " + e.message + "</div>";
            }}
        </script>
        """
        st.components.v1.html(html_code, height=520)
        _card_close()

# ==============================================================================
# MODULE 15 — Cyber Security Panel
# ==============================================================================

def render_cyber_panel() -> None:
    """Cyber Security Panel - Analyze Spam, Phishing, Malware headers, and Threat Feeds."""
    st.markdown("### 🛡️ Cyber Security Threat & File Analysis Panel")
    
    tab_email, tab_malware, tab_feeds = st.tabs([
        "📧 Email & Link Safety",
        "🪱 Malware & File Diagnostics",
        "📡 Live Threat Intelligence"
    ])
    
    with tab_email:
        _card_open()
        st.markdown("#### 📧 Email Header / Text / Link Reputation Analyzer")
        
        email_content = st.text_area(
            "Paste Email Content (Subject, Body, Headers, or Links)",
            placeholder="e.g. Dear customer, your account is suspended. Click here http://bank-verify-security.com to login...",
            height=150,
            key="cyber_email_text"
        )
        
        analyze_email = st.button("🛡️ Analyze Threat Vector", key="btn_analyze_email")
        _card_close()
        
        if analyze_email:
            if not email_content.strip():
                st.warning("Please paste some email content or links to analyze.")
            else:
                _card_open()
                st.markdown("##### 🔍 Cyber Analyst Report")
                
                with logo_spinner("Running threat model analysis..."):
                    import re
                    links = re.findall(r'https?://[^\s<>"]+|www\.[^\s<>"]+', email_content)
                    link_context = ""
                    if links:
                        link_context = "Found Links:\n" + "\n".join(links) + "\n\n"
                    
                    system_role = (
                        "You are an expert Cyber Security Analyst and Incident Responder. "
                        "Inspect the provided email/link threat vector. Detect signs of: "
                        "- Phishing indicators (urgent tone, generic greetings, spelling errors, mismatched domains) "
                        "- Spam markers "
                        "- Link reputation (check if domains look spoofed or suspicious) "
                        "Assign a Threat Level (Low / Medium / High / Critical) and outline mitigations."
                    )
                    
                    report = groq_chat(
                        f"{link_context}Email/Link Content to analyze:\n{email_content}",
                        system=system_role
                    )
                    st.markdown(report)
                _card_close()
                
    with tab_malware:
        _card_open()
        st.markdown("#### 🪱 File Structure & Malware Header Integrity Diagnostics")
        st.info("Upload any file to analyze its headers (Magic Bytes) for integrity and potential masquerading.")
        
        uploaded_file = st.file_uploader(
            "Upload File to Scan",
            key="cyber_file_uploader",
            help="Select any file (PNG, PDF, EXE, TXT, etc.). We inspect its header signature safely."
        )
        _card_close()
        
        if uploaded_file:
            _card_open()
            file_bytes = uploaded_file.getvalue()
            file_size = len(file_bytes)
            
            header_hex = file_bytes[:16].hex().upper()
            header_space = " ".join([header_hex[i:i+2] for i in range(0, len(header_hex), 2)])
            
            signatures = {
                "89 50 4E 47 0D 0A 1A 0A": "PNG Image File",
                "25 50 44 46": "PDF Document",
                "FF D8 FF": "JPEG Image File",
                "4D 5A": "Windows Executable (EXE / DLL)",
                "50 4B 03 04": "ZIP / Office OpenXML Archive (DOCX/PPTX/ZIP)",
                "7F 45 4C 46": "ELF Executable (Linux/Unix)",
                "49 44 33": "MP3 Audio File",
                "52 61 72 21 1A 07": "RAR Compressed Archive"
            }
            
            detected_format = "Unknown / Custom Binary Data"
            for sig, label in signatures.items():
                sig_clean = sig.replace(" ", "")
                if header_hex.startswith(sig_clean):
                    detected_format = label
                    break
                    
            st.markdown(f"**File Name:** `{uploaded_file.name}`")
            st.markdown(f"**File Size:** `{file_size} bytes`")
            st.markdown(f"**Magic Bytes (Hex Header):** `{header_space}`")
            st.markdown(f"**Detected Signature Format:** `{detected_format}`")
            
            with logo_spinner("Evaluating file header integrity and safety profile..."):
                prompt = (
                    f"File Name: {uploaded_file.name}\n"
                    f"File Size: {file_size} bytes\n"
                    f"Header Hex: {header_hex}\n"
                    f"Signature Match: {detected_format}\n\n"
                    "Analyze if the file extension matches the signature (masquerading check), "
                    "explain what this file signature means, and assess security concerns if any."
                )
                analysis = groq_chat(
                    prompt,
                    system="You are a Malware Analysis specialist. Provide a brief, professional integrity report."
                )
                st.markdown("##### 🧬 File Integrity Report")
                st.markdown(analysis)
            _card_close()
            
    with tab_feeds:
        _card_open()
        st.markdown("#### 📡 Global Cyber Security Threats & Vulnerability Feed")
        st.info("Fetches real-time security alerts, CVE releases, and cyber attack reports using Tavily Search.")
        
        search_keyword = st.text_input("Vulnerability / Threat Keyword Search", value="Recent ransomware campaigns zero-day exploits")
        fetch_clicked = st.button("📡 Fetch Threat Intel Feed", use_container_width=True)
        _card_close()
        
        if fetch_clicked or search_keyword:
            _card_open()
            st.markdown(f"##### 📢 Intel Report for: *{search_keyword}*")
            
            with logo_spinner(f"Scanning web for threat intel on '{search_keyword}'..."):
                results = duckduckgo_search(f"{search_keyword} security threat alert news cve")
                if not results:
                    st.warning("No recent threat reports found.")
                else:
                    context_pieces = []
                    for idx, res in enumerate(results[:5], 1):
                        context_pieces.append(
                            f"**[{idx}] {res['title']}**\n"
                            f"Source: {res['link']}\n"
                            f"Intel Summary: {res['snippet']}\n"
                        )
                        st.markdown(f"🔗 [{res['title']}]({res['link']})")
                        st.write(res['snippet'])
                        st.markdown("---")
                        
                    intel_context = "\n".join(context_pieces)
                    report = groq_chat(
                        f"Intel Context:\n{intel_context}\n\nProvide a concise executive threat brief summarizing the main risks, active CVEs, and recommended protection steps.",
                        system="You are a Cyber Threat Intelligence Specialist. Provide a clear, actionable brief."
                    )
                    st.markdown("##### 📝 Threat Intelligence Executive Brief")
                    st.markdown(report)
            _card_close()



# ==============================================================================
# ==============================================================================
# ==============================================================================
# MODULE: Presentation Generator
# ==============================================================================

def render_presentation_generator() -> None:
    """AI Presentation Generator - Create beautiful, downloadable PowerPoint slides from any prompt."""
    st.markdown("### 🎞️ AI Presentation Generator")
    st.markdown(
        '<div style="color:var(--sub); font-size:0.88rem; margin-bottom:16px;">'
        "Type a topic or detailed description — EduSphere AI will build a stunning, fully-coloured "
        "PowerPoint presentation with high-contrast text containers that solve readability issues on premium wave backgrounds."
        "</div>",
        unsafe_allow_html=True,
    )

    col_left, col_right = st.columns([1.6, 1], gap="medium")

    with col_right:
        _card_open()
        st.markdown("#### ⚙️ Slide Options")
        slide_count = st.slider("Number of Content Slides", min_value=4, max_value=14, value=7)
        include_chart = st.checkbox("📊 Include Chart Slide", value=True)
        include_table = st.checkbox("📋 Include Data Table Slide", value=True)
        theme_choice = st.selectbox(
            "🎨 Presentation Design & Theme",
            [
                "🌊 Rhythm Blue Wave (Premium)",
                "🌌 Deep Space (Dark)",
                "🌅 Sunrise (Orange)",
                "🌿 Nature (Green)",
                "💜 Neon Violet",
                "🔷 Corporate Blue",
            ],
        )
        st.markdown("---")
        st.markdown(
            '<div style="font-size:0.76rem; color:var(--sub);">'
            "✨ Dynamic high-contrast glassmorphism stencils<br>"
            "💧 All slides include <b>EduSphere AI</b> watermark<br>"
            "📥 Download as <b>.pptx</b> — opens in MS PowerPoint, Google Slides, LibreOffice"
            "</div>",
            unsafe_allow_html=True,
        )
        _card_close()

    with col_left:
        _card_open()
        st.markdown("#### 📝 Describe Your Presentation")
        user_topic = st.text_area(
            "Prompt",
            placeholder=(
                "Examples:\n"
                "• Make a presentation about climate change and its global impact\n"
                "• Slides on Machine Learning: types, applications, and future\n"
                "• Business pitch for a food delivery startup in Nepal\n"
                "• History of Nepal — geography, culture, achievements"
            ),
            height=160,
            label_visibility="collapsed",
        )
        gen_btn = st.button("✨ Generate Presentation", use_container_width=True, type="primary")
        _card_close()

    if gen_btn and user_topic.strip():
        with logo_spinner("🤖 EduSphere AI is designing your presentation…"):
            _generate_and_show_ppt(
                user_topic.strip(), slide_count, include_chart, include_table, theme_choice
            )
    elif gen_btn:
        st.warning("⚠️ Please enter a topic or description first.")


# ─── PPT helpers ──────────────────────────────────────────────────────────────

def _ppt_theme_colors(theme_choice: str) -> dict:
    themes = {
        "🌊 Rhythm Blue Wave (Premium)": dict(
            bg_title=(8, 10, 36), bg_content=(12, 14, 45), bg_alt=(16, 18, 55),
            accent=(0, 240, 255), accent2=(236, 72, 153),
            text=(248, 250, 252), sub_text=(148, 163, 184),
            chart_colors=["#00f0ff", "#ec4899", "#a855f7", "#10b981", "#fb923c", "#38bdf8"],
        ),
        "🌌 Deep Space (Dark)": dict(
            bg_title=(5, 5, 20), bg_content=(10, 12, 35), bg_alt=(15, 10, 40),
            accent=(0, 180, 255), accent2=(139, 92, 246),
            text=(240, 248, 255), sub_text=(180, 200, 230),
            chart_colors=["#00b4ff", "#8b5cf6", "#ec4899", "#10b981", "#f97316", "#f59e0b"],
        ),
        "🌅 Sunrise (Orange)": dict(
            bg_title=(20, 10, 5), bg_content=(35, 18, 8), bg_alt=(45, 20, 5),
            accent=(249, 115, 22), accent2=(251, 191, 36),
            text=(255, 250, 235), sub_text=(220, 180, 130),
            chart_colors=["#f97316", "#fbbf24", "#ef4444", "#ec4899", "#8b5cf6", "#06b6d4"],
        ),
        "🌿 Nature (Green)": dict(
            bg_title=(5, 18, 10), bg_content=(8, 28, 15), bg_alt=(10, 35, 20),
            accent=(52, 211, 153), accent2=(16, 185, 129),
            text=(236, 255, 244), sub_text=(160, 220, 180),
            chart_colors=["#34d399", "#10b981", "#06b6d4", "#84cc16", "#fbbf24", "#f97316"],
        ),
        "💜 Neon Violet": dict(
            bg_title=(8, 5, 20), bg_content=(12, 8, 30), bg_alt=(18, 10, 40),
            accent=(168, 85, 247), accent2=(236, 72, 153),
            text=(250, 245, 255), sub_text=(200, 170, 240),
            chart_colors=["#a855f7", "#ec4899", "#8b5cf6", "#f43f5e", "#06b6d4", "#10b981"],
        ),
        "🔷 Corporate Blue": dict(
            bg_title=(5, 15, 40), bg_content=(8, 22, 60), bg_alt=(12, 30, 75),
            accent=(59, 130, 246), accent2=(14, 165, 233),
            text=(240, 246, 255), sub_text=(170, 200, 240),
            chart_colors=["#3b82f6", "#0ea5e9", "#6366f1", "#8b5cf6", "#10b981", "#f97316"],
        ),
    }
    return themes.get(theme_choice, themes["🌌 Deep Space (Dark)"])


def _wm(slide, _tc):
    """Add EduSphere AI watermark to bottom-right of every slide."""
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    tb = slide.shapes.add_textbox(Inches(9.8), Inches(6.9), Inches(3.0), Inches(0.35))
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = "edusphere ai"
    r.font.name = "Montserrat"
    r.font.size = Pt(8.5)
    r.font.color.rgb = RGBColor(120, 140, 180)
    r.font.bold = True
    r.font.italic = True


def _set_bg(slide, rgb):
    from pptx.dml.color import RGBColor
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*rgb)


def _rect(slide, l, t, w, h, rgb):
    from pptx.util import Inches
    from pptx.dml.color import RGBColor
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = RGBColor(*rgb)
    s.line.fill.background()
    return s


def _tb(slide, text, l, t, w, h, size, bold=False, color=(255, 255, 255), align=None, italic=False):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    if align is None:
        align = PP_ALIGN.LEFT
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Space Grotesk"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = RGBColor(*color)
    return box


def _apply_animation_tag(slide):
    try:
        slide.slide_show_transition.type = 3  # MSO_ANIMATION.WHEEL / morph simulation
        slide.slide_show_transition.speed = 1 # fast
    except Exception:
        pass


def _generate_and_show_ppt(
    topic: str, slide_count: int, include_chart: bool, include_table: bool, theme_choice: str
) -> None:
    import io
    import json
    import re
    import os
    import datetime as dtm

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
    except ImportError as exc:
        st.error(f"❌ Missing library: {exc}. Run `pip install python-pptx matplotlib`")
        return

    tc = _ppt_theme_colors(theme_choice)
    A = tc["accent"]
    A2 = tc["accent2"]
    TX = tc["text"]
    SB = tc["sub_text"]
    BG0 = tc["bg_title"]
    BG1 = tc["bg_content"]
    BG2 = tc["bg_alt"]
    CC = tc["chart_colors"]

    CENTER = PP_ALIGN.CENTER
    LEFT = PP_ALIGN.LEFT

    # ── AI: slide outline with extra bullet point details ────────────────────────
    ai_prompt = (
        f'Create a high-quality PowerPoint presentation outline for: "{topic}"\n'
        f"Generate exactly {slide_count} slides.\n"
        "Each content slide MUST have: \n"
        "1. A short, punchy title\n"
        "2. Exactly 3-4 structured, descriptive bullet points. Every bullet point MUST contain a bold title followed by a detail explanation, like: '**Topic Header:** Detailed explanation text.'\n"
        "Return ONLY a valid JSON array format, no other text:\n"
        '[{"slide_num":1,"title":"Title","points":["**Header 1:** Detailed explanation","**Header 2:** Detailed explanation"]}]'
    )
    try:
        raw = groq_chat(ai_prompt, system="Return ONLY valid JSON arrays. No markdown formatting, no commentary.")
        raw = re.sub(r"```(?:json)?", "", raw).strip().rstrip("`").strip()
        slides_data = json.loads(raw)
        if not isinstance(slides_data, list):
            raise ValueError("Not a list")
    except Exception:
        slides_data = [
            {
                "slide_num": i + 1,
                "title": f"Core Pillar {i + 1}",
                "points": [
                    f"**Strategic Objective {i + 1}:** Detailed explanation of the primary goal and milestone targets.",
                    "**Data Grounding:** Evidence-based statistics supporting this core pillar and framework.",
                    "**Market Alignment:** Exploring consumer fitment, demographics, and execution strategies.",
                ],
            }
            for i in range(slide_count)
        ]

    # ── AI: chart data ─────────────────────────────────────────────────────────
    chart_data = None
    if include_chart:
        try:
            cr = groq_chat(
                f'For "{topic}", create 5 data points for a bar chart. '
                'Return ONLY JSON: {"title":"Chart Title","labels":["A","B","C","D","E"],"values":[45,78,62,91,55]}',
                system="Return ONLY valid JSON. No markdown.",
            )
            cr = re.sub(r"```(?:json)?", "", cr).strip().rstrip("`").strip()
            chart_data = json.loads(cr)
        except Exception:
            chart_data = {
                "title": f"Key Metrics: {topic[:35]}",
                "labels": ["Pillar 1", "Pillar 2", "Pillar 3", "Pillar 4", "Pillar 5"],
                "values": [65, 82, 55, 90, 73],
            }

    # ── AI: table data ─────────────────────────────────────────────────────────
    table_data = None
    if include_table:
        try:
            tr = groq_chat(
                f'For "{topic}", create a comparison table. '
                'Return ONLY JSON: {"title":"Table Title","headers":["Col1","Col2","Col3"],'
                '"rows":[["v1","v2","v3"],["v4","v5","v6"],["v7","v8","v9"],["v10","v11","v12"]]}',
                system="Return ONLY valid JSON. No markdown.",
            )
            tr = re.sub(r"```(?:json)?", "", tr).strip().rstrip("`").strip()
            table_data = json.loads(tr)
        except Exception:
            table_data = {
                "title": f"Overview: {topic[:35]}",
                "headers": ["Aspect", "Details", "Impact Assessment"],
                "rows": [
                    ["Pillar 1", "Implementation Stage A", "High Return"],
                    ["Pillar 2", "Implementation Stage B", "Medium Stability"],
                    ["Pillar 3", "Implementation Stage C", "High Return"],
                    ["Pillar 4", "Implementation Stage D", "Low Risk"],
                ],
            }

    # ── Build PowerPoint ───────────────────────────────────────────────────────
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    bg_img_path = "assets/ppt_wave_bg.png"
    has_custom_bg = os.path.exists(bg_img_path)

    # Helper function to place the premium blue wave background image
    def apply_slide_bg(slide, bg_color):
        if theme_choice == "🌊 Rhythm Blue Wave (Premium)" and has_custom_bg:
            slide.shapes.add_picture(bg_img_path, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
        else:
            _set_bg(slide, bg_color)

    # === Slide 1: Title ===
    sl = prs.slides.add_slide(blank)
    apply_slide_bg(sl, BG0)
    _apply_animation_tag(sl)

    # Semi-transparent glassmorphic title card to separate text from background wave flares
    title_glass = sl.shapes.add_shape(1, Inches(0.4), Inches(1.8), Inches(12.53), Inches(3.2))
    title_glass.fill.solid()
    title_glass.fill.fore_color.rgb = RGBColor(15, 23, 42) # dark slate base
    title_glass.line.color.rgb = RGBColor(*A)
    title_glass.line.width = Pt(1.5)

    _rect(sl, 0, 0, 13.33, 0.09, A)           # top bar
    _rect(sl, 0, 0.09, 0.09, 7.41, A2)         # left strip
    _tb(sl, topic.upper(), 0.5, 2.3, 12.33, 1.6, 38, bold=True, color=TX, align=CENTER)
    _rect(sl, 4.5, 3.95, 4.3, 0.05, A)
    _tb(sl, "Powered by EduSphere AI  ·  Premium Presentation Studio", 0.5, 4.15, 12.33, 0.55,
        15, italic=True, color=SB, align=CENTER)
    _tb(sl, str(dtm.datetime.now().year), 0.5, 6.6, 12.33, 0.4, 11, color=SB, align=CENTER, italic=True)
    _wm(sl, tc)

    # === Content Slides ===
    for idx, sinfo in enumerate(slides_data):
        sl = prs.slides.add_slide(blank)
        _apply_animation_tag(sl)
        bg_c = BG2 if idx % 2 == 0 else BG1
        apply_slide_bg(sl, bg_c)
        _rect(sl, 0, 0, 13.33, 0.07, A)

        # Slide number badge
        badge = sl.shapes.add_shape(1, Inches(11.9), Inches(0.12), Inches(1.2), Inches(0.35))
        badge.fill.solid()
        badge.fill.fore_color.rgb = RGBColor(*A2)
        badge.line.fill.background()
        _tb(sl, f"{idx+1}/{len(slides_data)}", 11.95, 0.14, 1.1, 0.3, 8,
            bold=True, color=BG0, align=CENTER)

        slide_title = sinfo.get("title", f"Slide {idx+1}")
        
        # High contrast title box backing
        title_backing = sl.shapes.add_shape(1, Inches(0.28), Inches(0.15), Inches(11.4), Inches(0.85))
        title_backing.fill.solid()
        title_backing.fill.fore_color.rgb = RGBColor(12, 15, 35)
        title_backing.line.fill.background()
        
        _tb(sl, slide_title, 0.38, 0.22, 11.2, 0.85, 27, bold=True, color=A)
        _rect(sl, 0.38, 0.95, 5.0, 0.04, A2)

        # Content Rendering (Left bar visual + bullet details parser)
        _rect(sl, 0.35, 1.5, 0.05, 5.0, A2)
        points = sinfo.get("points", [])
        
        for b_idx, pt in enumerate(points[:5]):
            top_pos = 1.48 + b_idx * 1.35
            
            # Semi-transparent backing panel for each bullet point to ensure max readability on light waves/stars
            panel = sl.shapes.add_shape(1, Inches(0.45), Inches(top_pos - 0.08), Inches(12.38), Inches(1.15))
            panel.fill.solid()
            panel.fill.fore_color.rgb = RGBColor(16, 20, 48) # high contrast dark navy card
            panel.line.color.rgb = RGBColor(40, 50, 90)
            panel.line.width = Pt(1.0)
            
            # Sub-elements for bullet points (Bold Title + Details layout)
            if "**" in pt:
                parts = pt.split("**", 2)
                bold_header = parts[1].replace(":", "").strip()
                detail_text = parts[2].strip()
            else:
                bold_header = "Objective"
                detail_text = pt

            # Bullet Indicator icon
            _tb(sl, "⚡", 0.58, top_pos, 0.4, 0.4, 13, bold=True, color=A)
            
            # Content textbox with header in accent, body in white
            tb_box = sl.shapes.add_textbox(Inches(0.95), Inches(top_pos - 0.05), Inches(11.8), Inches(1.1))
            tf = tb_box.text_frame
            tf.word_wrap = True
            
            # Paragraph 1: Bold subheader
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = bold_header
            r.font.name = "Montserrat"
            r.font.size = Pt(15.5)
            r.font.bold = True
            r.font.color.rgb = RGBColor(*A)

            # Paragraph 2: Detailed text body
            p2 = tf.add_paragraph()
            r2 = p2.add_run()
            r2.text = detail_text
            r2.font.name = "Space Grotesk"
            r2.font.size = Pt(13)
            r2.font.color.rgb = RGBColor(*TX)

        _wm(sl, tc)

    # === Chart Slide ===
    if include_chart and chart_data:
        sl = prs.slides.add_slide(blank)
        _apply_animation_tag(sl)
        apply_slide_bg(sl, BG0)
        _rect(sl, 0, 0, 13.33, 0.07, A)
        
        # High contrast title box backing
        title_backing = sl.shapes.add_shape(1, Inches(0.28), Inches(0.1), Inches(12.77), Inches(0.85))
        title_backing.fill.solid()
        title_backing.fill.fore_color.rgb = RGBColor(12, 15, 35)
        title_backing.line.fill.background()

        _tb(sl, "📊 " + chart_data.get("title", "Data Visualization"),
            0.35, 0.14, 12.8, 0.78, 24, bold=True, color=A)
        _rect(sl, 0.35, 0.85, 4.5, 0.04, A2)

        fig, ax = plt.subplots(figsize=(11.5, 5.5))
        fig.patch.set_facecolor(f"#{BG0[0]:02x}{BG0[1]:02x}{BG0[2]:02x}")
        ax.set_facecolor(f"#{BG1[0]:02x}{BG1[1]:02x}{BG1[2]:02x}")
        labels = chart_data.get("labels", [])
        values = chart_data.get("values", [])
        bar_colors_list = [CC[i % len(CC)] for i in range(len(labels))]
        bars = ax.bar(labels, values, color=bar_colors_list, edgecolor="none", width=0.55)
        for br, val in zip(bars, values):
            ax.text(
                br.get_x() + br.get_width() / 2, br.get_height() + 1,
                str(val), ha="center", va="bottom",
                color="white", fontsize=11, fontweight="bold",
            )
        ax.set_xticks(range(len(labels)))
        ax.set_xticklabels(labels, color=f"#{TX[0]:02x}{TX[1]:02x}{TX[2]:02x}", fontsize=10)
        ax.tick_params(axis="y", colors=f"#{SB[0]:02x}{SB[1]:02x}{SB[2]:02x}")
        for sp in ["top", "right"]:
            ax.spines[sp].set_visible(False)
        ax.spines["left"].set_color(f"#{SB[0]:02x}{SB[1]:02x}{SB[2]:02x}")
        ax.spines["bottom"].set_color(f"#{SB[0]:02x}{SB[1]:02x}{SB[2]:02x}")
        ax.yaxis.grid(True, alpha=0.2, color="white")
        ax.set_axisbelow(True)
        plt.tight_layout(pad=0.5)
        chart_buf = io.BytesIO()
        plt.savefig(chart_buf, format="png", dpi=150, bbox_inches="tight",
                    facecolor=fig.get_facecolor())
        plt.close(fig)
        chart_buf.seek(0)
        
        # High contrast picture framing card
        pic_frame = sl.shapes.add_shape(1, Inches(0.35), Inches(0.95), Inches(12.63), Inches(6.25))
        pic_frame.fill.solid()
        pic_frame.fill.fore_color.rgb = RGBColor(10, 12, 32)
        pic_frame.line.color.rgb = RGBColor(40, 50, 90)
        pic_frame.line.width = Pt(1.5)

        sl.shapes.add_picture(chart_buf, Inches(0.4), Inches(1.0), Inches(12.5), Inches(6.1))
        _wm(sl, tc)

    # === Table Slide ===
    if include_table and table_data:
        from pptx.enum.text import PP_ALIGN as _PA
        sl = prs.slides.add_slide(blank)
        _apply_animation_tag(sl)
        apply_slide_bg(sl, BG1)
        _rect(sl, 0, 0, 13.33, 0.07, A)
        
        # High contrast title box backing
        title_backing = sl.shapes.add_shape(1, Inches(0.28), Inches(0.1), Inches(12.77), Inches(0.85))
        title_backing.fill.solid()
        title_backing.fill.fore_color.rgb = RGBColor(12, 15, 35)
        title_backing.line.fill.background()

        _tb(sl, "📋 " + table_data.get("title", "Data Overview"),
            0.35, 0.14, 12.8, 0.78, 24, bold=True, color=A)
        _rect(sl, 0.35, 0.85, 4.5, 0.04, A2)
        headers = table_data.get("headers", [])
        rows = table_data.get("rows", [])
        if headers and rows:
            nc = len(headers)
            nr = len(rows) + 1
            
            # High contrast table framing panel
            tbl_frame = sl.shapes.add_shape(1, Inches(0.35), Inches(0.95), Inches(12.6), Inches(min(6.2, nr * 0.94 + 0.2)))
            tbl_frame.fill.solid()
            tbl_frame.fill.fore_color.rgb = RGBColor(10, 12, 32)
            tbl_frame.line.color.rgb = RGBColor(40, 50, 90)
            tbl_frame.line.width = Pt(1.5)

            tbl = sl.shapes.add_table(
                nr, nc, Inches(0.45), Inches(1.05), Inches(12.4), Inches(min(5.8, nr * 0.9))
            ).table
            # Header row
            for ci, hdr in enumerate(headers):
                cell = tbl.cell(0, ci)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(*A)
                p = cell.text_frame.paragraphs[0]
                p.alignment = _PA.CENTER
                r = p.add_run()
                r.text = str(hdr)
                r.font.name = "Montserrat"
                r.font.size = Pt(13)
                r.font.bold = True
                r.font.color.rgb = RGBColor(*BG0)
            # Data rows
            for ri, row_vals in enumerate(rows):
                row_bg = (
                    (min(A[0] // 8, 255), min(A[1] // 8, 255), min(A[2] // 8, 255))
                    if ri % 2 == 0
                    else (min(BG2[0] + 15, 255), min(BG2[1] + 15, 255), min(BG2[2] + 15, 255))
                )
                for ci in range(min(nc, len(row_vals))):
                    cell = tbl.cell(ri + 1, ci)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(*row_bg)
                    p = cell.text_frame.paragraphs[0]
                    p.alignment = _PA.CENTER
                    r = p.add_run()
                    r.text = str(row_vals[ci])
                    r.font.name = "Space Grotesk"
                    r.font.size = Pt(12)
                    r.font.color.rgb = RGBColor(*TX)
        _wm(sl, tc)

    # === Thank You Slide ===
    sl = prs.slides.add_slide(blank)
    _apply_animation_tag(sl)
    apply_slide_bg(sl, BG0)
    _rect(sl, 0, 0, 13.33, 0.08, A)
    _rect(sl, 0, 0.08, 0.08, 7.42, A2)

    # High contrast card backing for closing slide text
    close_card = sl.shapes.add_shape(1, Inches(2.5), Inches(1.8), Inches(8.33), Inches(3.8))
    close_card.fill.solid()
    close_card.fill.fore_color.rgb = RGBColor(12, 15, 35)
    close_card.line.color.rgb = RGBColor(*A2)
    close_card.line.width = Pt(1.5)

    _tb(sl, "🙏 Thank You", 0.5, 2.3, 12.3, 1.0, 42, bold=True, color=A, align=CENTER)
    _rect(sl, 4.5, 3.55, 4.3, 0.05, A2)
    _tb(sl, topic, 0.5, 3.65, 12.3, 0.6, 18, italic=True, color=SB, align=CENTER)
    _tb(sl, "Powered by EduSphere AI", 0.5, 4.5, 12.3, 0.45, 13, italic=True, color=SB, align=CENTER)
    _wm(sl, tc)

    # ── Save & download ────────────────────────────────────────────────────────
    ppt_buf = io.BytesIO()
    prs.save(ppt_buf)
    ppt_buf.seek(0)

    total = (
        1
        + len(slides_data)
        + (1 if include_chart and chart_data else 0)
        + (1 if include_table and table_data else 0)
        + 1
    )
    st.success(f"Presentation generated! **{total} slides** ready.")
    st.markdown(
        '<div style="background:var(--card); border:1px solid rgba(255,255,255,0.12); border-radius:14px; '
        'padding:22px; margin:14px 0; text-align:center;">'
        '<div style="font-size:2.8rem; margin-bottom:8px;">*</div>'
        '<div style="font-size:1.1rem; font-weight:700; color:var(--text); margin-bottom:6px;">'
        'Your Presentation is Ready!'
        '</div>'
        '<div style="font-size:0.85rem; color:var(--sub);">'
        '{total} slides &nbsp;·&nbsp; {theme_choice} theme &nbsp;·&nbsp; EduSphere AI watermark on every slide'
        '</div>'
        '</div>'.format(total=total, theme_choice=theme_choice),
        unsafe_allow_html=True,
    )
    safe = "".join(c if c.isalnum() or c in " _-" else "_" for c in topic[:40]).strip().replace(" ", "_")
    st.download_button(
        label="📥 Download Presentation (.pptx)",
        data=ppt_buf.getvalue(),
        file_name=f"EduSphere_{safe}.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        use_container_width=True,
        type="primary",
    )

    # Slide structure preview
    st.markdown("#### 📋 Slide Structure Preview")
    all_titles = ["🎯 Title Slide"] + [s.get("title", f"Slide {i+1}") for i, s in enumerate(slides_data)]
    if include_chart and chart_data:
        all_titles.append("📊 Chart Slide")
    if include_table and table_data:
        all_titles.append("📋 Table Slide")
    all_titles.append("🙏 Thank You")

    pcols = st.columns(min(4, len(all_titles)))
    for i, ttl in enumerate(all_titles[:16]):
        with pcols[i % len(pcols)]:
            st.markdown(
                f'<div style="background:var(--card); border:1px solid rgba(255,255,255,0.08);'
                f'border-radius:8px; padding:8px 10px; margin-bottom:8px; font-size:0.74rem; color:var(--sub);">'
                f'<span style="color:var(--accent); font-weight:700;">#{i+1}</span> {ttl}'
                f"</div>",
                unsafe_allow_html=True,
            )
